"""Read layer (ops/read.py) against the real corpus, with python-pptx as the
independent oracle for text extraction and slide enumeration.

Oracle discipline (Phase 0 lesson): python-pptx's Presentation.slides
RENUMBERS slide partnames on mere access, so every oracle helper builds a
FRESH Presentation instance and resolves order through the sldIdLst-rId
pattern; partname comparisons never reuse an instance that touched .slides.
python-pptx renders a:br as a vertical tab ("\\x0b"); ours renders a
newline, so oracle text is normalized before comparison.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest

from kitchensink4ppt.core.errors import PptMcpError, TargetNotFound
from kitchensink4ppt.core.package import PptxPackage, qn
from kitchensink4ppt.ops import read

CORPUS = Path(__file__).parents[1] / "corpus"

PPTX_NAMES = [
    "proposal_defense.pptx",
    "nsu_pcsj.pptx",
    "unitar_final.pptx",
    "military_brief.pptx",
    "pmr_tables.pptx",
]


def _work_copy(name: str, tmp_path: Path) -> Path:
    src = CORPUS / name
    if not src.exists():
        pytest.skip(f"corpus file missing: {name}")
    work = tmp_path / name
    shutil.copy2(src, work)
    return work


# ------------------------------------------------------------------- oracles


def _oracle_slide_ids_and_parts(path: Path) -> list[tuple[int, str]]:
    """(slide_id, part name) in presentation order, via python-pptx's own
    rels machinery, NEVER via Presentation.slides (which renumbers)."""
    from pptx import Presentation

    part = Presentation(str(path)).part
    lst = part._element.get_or_add_sldIdLst()
    return [
        (int(s.get("id")), str(part.related_part(s.rId).partname)[1:])
        for s in lst
    ]


def _oracle_slide_texts(path: Path) -> list[str]:
    """Per-slide reading-order text via python-pptx, mirroring our reading
    rules: spTree order, groups recursed, tables as tab/newline grids."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))  # fresh instance for this assertion only
    out = []

    def walk(shapes, parts):
        for sh in shapes:
            try:
                st = sh.shape_type
            except Exception:
                st = None
            if st == MSO_SHAPE_TYPE.GROUP:
                walk(sh.shapes, parts)
            elif getattr(sh, "has_text_frame", False):
                t = sh.text_frame.text.replace("\x0b", "\n")
                if t:
                    parts.append(t)
            elif getattr(sh, "has_table", False):
                t = "\n".join(
                    "\t".join(
                        c.text_frame.text.replace("\x0b", "\n") for c in r.cells
                    )
                    for r in sh.table.rows
                )
                if t:
                    parts.append(t)

    for slide in prs.slides:
        parts: list[str] = []
        walk(slide.shapes, parts)
        out.append("\n".join(parts))
    return out


def _oracle_table_count(path: Path) -> int:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    n = 0

    def walk(shapes):
        nonlocal n
        for sh in shapes:
            try:
                st = sh.shape_type
            except Exception:
                st = None
            if st == MSO_SHAPE_TYPE.GROUP:
                walk(sh.shapes)
            elif getattr(sh, "has_table", False):
                n += 1

    for slide in prs.slides:
        walk(slide.shapes)
    return n


# ------------------------------------------------------- get_text vs oracle


@pytest.mark.parametrize(
    "name", ["proposal_defense.pptx", "pmr_tables.pptx", "military_brief.pptx"]
)
def test_get_text_matches_python_pptx_oracle(name, tmp_path):
    work = _work_copy(name, tmp_path)
    ours = read.get_text(PptxPackage(work))
    oracle = _oracle_slide_texts(work)
    assert ours["slide_count"] == len(oracle)
    for i, (mine, ref) in enumerate(
        zip((s["text"] for s in ours["slides"]), oracle)
    ):
        assert mine == ref, f"{name} slide {i}: text diverges from oracle"


@pytest.mark.parametrize("name", PPTX_NAMES)
def test_slide_enumeration_matches_oracle(name, tmp_path):
    work = _work_copy(name, tmp_path)
    listed = read.list_elements(PptxPackage(work), "slides")
    oracle = _oracle_slide_ids_and_parts(work)
    assert listed["count"] == len(oracle) > 0
    for item, (sid, part) in zip(listed["items"], oracle):
        assert item["slide_id"] == sid
        assert item["part"] == part


# --------------------------------------------------------- presentation info


def test_get_presentation_info_proposal_deck(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    info = read.get_presentation_info(PptxPackage(work))
    assert info["slide_count"] == 26
    assert info["slide_size"]["cx"] > 0 and info["slide_size"]["cx_in"] > 1
    assert info["masters"], "no masters reported"
    assert all(m["layouts"] for m in info["masters"]), "master without layouts"
    assert isinstance(info["sections"], list)


def test_get_slide_info_shapes_and_placeholders(make_deck):
    pkg = PptxPackage(make_deck())
    info = read.get_slide_info(pkg, 0)  # title slide
    assert info["index"] == 0 and info["slide_id"] >= 256
    assert info["layout_part"].startswith("ppt/slideLayouts/")
    assert info["shape_count"] == len(info["shapes"]) > 0
    types = {s["type"] for s in info["shapes"]}
    assert "placeholder" in types
    assert info["placeholders"], "title slide reported no placeholders"
    assert {"title", "ctrTitle"} & {
        p["type"] for p in info["placeholders"]
    }, "no title placeholder found"
    # z-positions are the 0-based spTree order of top-level shapes.
    top = [s for s in info["shapes"] if "group_id" not in s]
    assert [s["z"] for s in top] == list(range(len(top)))
    # Shape ids unique per slide; geometry (when present) has EMU + inches.
    ids = [s["id"] for s in info["shapes"]]
    assert len(ids) == len(set(ids))
    for s in info["shapes"]:
        if s["geometry"]:
            assert s["geometry"]["cx_in"] == pytest.approx(
                s["geometry"]["cx"] / 914400, abs=0.01
            )


def test_slide_addressing_by_id_equals_by_index(make_deck):
    pkg = PptxPackage(make_deck())
    by_index = read.get_slide_info(pkg, 1)
    by_id = read.get_slide_info(pkg, {"slide_id": by_index["slide_id"]})
    assert by_id == by_index


def test_out_of_range_and_bad_selectors_refused(make_deck):
    pkg = PptxPackage(make_deck())
    n = len(read.slide_table(pkg))
    with pytest.raises(TargetNotFound, match=f"out of range, presentation has {n}"):
        read.get_slide_info(pkg, 99)
    with pytest.raises(TargetNotFound, match="no slide with slide_id"):
        read.get_slide_info(pkg, {"slide_id": 9})
    with pytest.raises(PptMcpError):
        read.get_slide_info(pkg, "3")
    with pytest.raises(PptMcpError, match="unknown element kind"):
        read.list_elements(pkg, "widgets")


def test_hidden_slide_flag(make_deck):
    pkg = PptxPackage(make_deck())
    part = read.slide_table(pkg)[2]["part"]
    assert read.get_slide_info(pkg, 2)["hidden"] is False
    pkg.root(part).set("show", "0")  # in-memory only; read layer never saves
    assert read.get_slide_info(pkg, 2)["hidden"] is True
    assert read.list_elements(pkg, "slides")["items"][2]["hidden"] is True


# ------------------------------------------------------------- list_elements


@pytest.mark.parametrize("name", ["pmr_tables.pptx", "military_brief.pptx"])
def test_list_tables_matches_oracle(name, tmp_path):
    """Table count agrees with python-pptx. NOTE (corpus finding): the REAL
    pmr_tables.pptx contains no native a:tbl tables at all (its 'tables'
    are drawn shapes); military_brief.pptx carries 9 real ones, so it is
    the native-table coverage deck. The synthetic stand-ins have one."""
    work = _work_copy(name, tmp_path)
    pkg = PptxPackage(work)
    tables = read.list_elements(pkg, "tables")
    assert tables["count"] == _oracle_table_count(work)
    for item in tables["items"]:
        assert item["type"] == "table"
        assert item["rows"] >= 1 and item["cols"] >= 1


def test_list_tables_nonempty_on_synthetic_deck(make_deck):
    pkg = PptxPackage(make_deck())
    tables = read.list_elements(pkg, "tables")
    assert tables["count"] == 1
    assert (tables["items"][0]["rows"], tables["items"][0]["cols"]) == (3, 3)


def test_list_images_and_layouts_and_masters(make_deck):
    pkg = PptxPackage(make_deck())
    images = read.list_elements(pkg, "images")
    assert images["count"] >= 1
    assert images["items"][0]["media_part"].startswith("ppt/media/")
    layouts = read.list_elements(pkg, "layouts")
    masters = read.list_elements(pkg, "masters")
    assert masters["count"] >= 1
    assert layouts["count"] == sum(m["layout_count"] for m in masters["items"])
    assert all(
        it["master_part"] == masters["items"][0]["part"]
        for it in layouts["items"]
    )


def test_list_placeholders_scope_single_slide(make_deck):
    pkg = PptxPackage(make_deck())
    all_ph = read.list_elements(pkg, "placeholders")
    one = read.list_elements(pkg, "placeholders", scope=0)
    assert one["count"] < all_ph["count"]
    assert all(it["slide_index"] == 0 for it in one["items"])
    assert all("placeholder_type" in it for it in one["items"])


def test_list_notes(make_deck):
    pkg = PptxPackage(make_deck())
    notes = read.list_elements(pkg, "notes")
    # make_corpus puts notes on the title slide (0) and picture slide (3).
    assert {it["slide_index"] for it in notes["items"]} == {0, 3}
    assert notes["items"][0]["text"] == "Synthetic speaker notes, slide one."


def test_get_text_include_notes(make_deck):
    pkg = PptxPackage(make_deck())
    out = read.get_text(pkg, include_notes=True)
    assert out["slides"][0]["notes"] == "Synthetic speaker notes, slide one."
    assert "[Notes] Synthetic speaker notes, slide one." in out["text"]
    # Default excludes notes entirely.
    assert "notes" not in read.get_text(pkg)["slides"][0]


# ----------------------------------------------------------------- find_text


def test_find_text_offsets_are_exact(make_deck):
    pkg = PptxPackage(make_deck())
    out = read.find_text(pkg, "Synthetic")
    assert out["count"] >= 1
    title_hits = [
        m for m in out["matches"] if m["slide_index"] == 0 and m["where"] == "slide"
    ]
    assert title_hits, "no hit on the title slide"
    hit = title_hits[0]
    assert (hit["start"], hit["end"]) == (0, len("Synthetic"))
    assert hit["match"] == "Synthetic"
    assert "Synthetic" in hit["context"]
    assert hit["shape_id"] and hit["slide_id"] >= 256
    # Notes hits are labeled and carry paragraph offsets too.
    notes_hits = [m for m in out["matches"] if m["where"] == "notes"]
    assert notes_hits and notes_hits[0]["start"] == 0


def test_find_text_in_tables_carries_row_col(make_deck):
    pkg = PptxPackage(make_deck())
    out = read.find_text(pkg, r"r\dc\d", regex=True)
    table_hits = [m for m in out["matches"] if m["where"] == "table"]
    assert len(table_hits) == 9  # 3x3 synthetic table, one label per cell
    for m in table_hits:
        assert m["match"] == f"r{m['row']}c{m['col']}"


def test_find_text_scope_and_refusals(make_deck):
    pkg = PptxPackage(make_deck())
    scoped = read.find_text(pkg, "Synthetic", scope=1)
    assert all(m["slide_index"] == 1 for m in scoped["matches"])
    with pytest.raises(PptMcpError, match="non-empty"):
        read.find_text(pkg, "")
    with pytest.raises(PptMcpError, match="invalid regex"):
        read.find_text(pkg, "(unclosed", regex=True)


def test_find_text_on_real_deck(tmp_path):
    """Every match's offsets must slice back out of our own extraction."""
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    out = read.find_text(pkg, "the")
    for m in out["matches"][:50]:
        assert m["match"] == "the"
        assert m["end"] - m["start"] == 3


# -------------------------------------------------------- military stress run


def test_military_brief_full_read_layer_stress(tmp_path):
    """The read layer must survive the heaviest corpus deck end to end."""
    work = _work_copy("military_brief.pptx", tmp_path)
    pkg = PptxPackage(work)
    info = read.get_presentation_info(pkg)
    assert info["slide_count"] >= 30
    for kind in (
        "slides", "shapes", "placeholders", "tables", "charts",
        "images", "notes", "sections", "layouts", "masters",
    ):
        out = read.list_elements(pkg, kind)
        assert out["count"] == len(out["items"])
    for i in range(info["slide_count"]):
        si = read.get_slide_info(pkg, i)
        assert si["index"] == i
    text = read.get_text(pkg, include_notes=True)
    assert text["slide_count"] == info["slide_count"]
    found = read.find_text(pkg, "the")
    assert found["count"] == len(found["matches"])
