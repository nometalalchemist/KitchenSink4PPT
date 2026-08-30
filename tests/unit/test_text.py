"""Phase 3 ops: the text engine (placeholders, text boxes, formatting,
bullets, search/replace, autofit reporting), exercised on the real corpus via
tmp copies (originals are never opened for saving) plus deliberately built
structures the corpus lacks (fragmented runs, fields).

Every mutated output is saved, which runs pkg._validate_payload; round-trip
tests reopen the saved file and re-read through the read layer (and the
python-pptx oracle where it can see the property).
"""

from __future__ import annotations

import collections
import re
import shutil
from pathlib import Path

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import (
    AmbiguousTarget,
    PptMcpError,
    TargetNotFound,
    UnsupportedStructure,
)
from kitchensink4ppt.core.package import PptxPackage, qn
from kitchensink4ppt.ops import _runmap
from kitchensink4ppt.ops.read import (
    find_text,
    get_slide_info,
    get_text,
    list_elements,
    paragraph_text,
    shape_text,
    txbody_paragraphs,
)
from kitchensink4ppt.ops.slides import delete_slide, insert_slide
from kitchensink4ppt.ops.text import (
    format_text,
    get_autofit_state,
    insert_textbox,
    search_and_replace,
    set_bullets,
    set_placeholder_text,
)

CORPUS = Path(__file__).parents[1] / "corpus"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _work_copy(name: str, tmp_path: Path) -> Path:
    src = CORPUS / name
    if not src.exists():
        pytest.skip(f"corpus file missing: {name}")
    work = tmp_path / name
    shutil.copy2(src, work)
    return work


def _shape_elem(pkg: PptxPackage, slide_index: int, shape_id: int):
    part = pkg.slide_parts()[slide_index]
    for cnvpr in pkg.root(part).iter(qn("p:cNvPr")):
        if int(cnvpr.get("id", "0")) == shape_id:
            return cnvpr.getparent().getparent()
    raise AssertionError(f"shape {shape_id} not found on slide {slide_index}")


def _slide_with_title_and_body(pkg: PptxPackage) -> dict:
    """Insert a fresh slide from the first layout whose skeleton yields both
    a title-family and a text-family placeholder; failed attempts are
    deleted again."""
    n_layouts = list_elements(pkg, "layouts")["count"]
    for i in range(n_layouts):
        result = insert_slide(pkg, i)
        types = {p["type"] for p in result["placeholders"]}
        if types & {"title", "ctrTitle"} and types & {"body", "obj"}:
            return result
        delete_slide(pkg, {"slide_id": result["slide_id"]})
    pytest.skip("no layout offers both a title and a body placeholder")


def _add_table(pkg: PptxPackage, slide_index: int, rows: list[list[str]]) -> int:
    """Minimal a:tbl graphicFrame appended to a slide (the corpus decks
    carry no native tables, verified 2026-08-30); returns the shape id."""
    part = pkg.slide_parts()[slide_index]
    sp_tree = pkg.root(part).find(f"{qn('p:cSld')}/{qn('p:spTree')}")
    sid = pkg.next_shape_id(part)
    frame = etree.SubElement(sp_tree, qn("p:graphicFrame"))
    nv = etree.SubElement(frame, qn("p:nvGraphicFramePr"))
    cnv = etree.SubElement(nv, qn("p:cNvPr"))
    cnv.set("id", str(sid))
    cnv.set("name", f"Table {sid}")
    gfp = etree.SubElement(nv, qn("p:cNvGraphicFramePr"))
    locks = etree.SubElement(gfp, qn("a:graphicFrameLocks"))
    locks.set("noGrp", "1")
    etree.SubElement(nv, qn("p:nvPr"))
    xfrm = etree.SubElement(frame, qn("p:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", "914400")
    off.set("y", "914400")
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(1828800 * len(rows[0])))
    ext.set("cy", str(370840 * len(rows)))
    graphic = etree.SubElement(frame, qn("a:graphic"))
    data = etree.SubElement(graphic, qn("a:graphicData"))
    data.set("uri", "http://schemas.openxmlformats.org/drawingml/2006/table")
    tbl = etree.SubElement(data, qn("a:tbl"))
    tblpr = etree.SubElement(tbl, qn("a:tblPr"))
    tblpr.set("firstRow", "1")
    tblpr.set("bandRow", "1")
    grid = etree.SubElement(tbl, qn("a:tblGrid"))
    for _ in rows[0]:
        col = etree.SubElement(grid, qn("a:gridCol"))
        col.set("w", "1828800")
    for row in rows:
        tr = etree.SubElement(tbl, qn("a:tr"))
        tr.set("h", "370840")
        for cell in row:
            tc = etree.SubElement(tr, qn("a:tc"))
            body = etree.SubElement(tc, qn("a:txBody"))
            etree.SubElement(body, qn("a:bodyPr"))
            etree.SubElement(body, qn("a:lstStyle"))
            p = etree.SubElement(body, qn("a:p"))
            r = etree.SubElement(p, qn("a:r"))
            t = etree.SubElement(r, qn("a:t"))
            t.text = cell
            etree.SubElement(tc, qn("a:tcPr"))
    pkg.mark_dirty(part)
    return sid


def _fragment_first_run(pkg: PptxPackage, slide_index: int, shape_id: int, cuts):
    """Split the first run of the shape's first paragraph at the given
    character offsets, building a deliberately fragmented paragraph."""
    elem = _shape_elem(pkg, slide_index, shape_id)
    p = elem.find(f"{qn('p:txBody')}/{qn('a:p')}")
    r = p.find(qn("a:r"))
    t = r.find(qn("a:t"))
    text = t.text
    pieces = []
    prev = 0
    for cut in cuts:
        pieces.append(text[prev:cut])
        prev = cut
    pieces.append(text[prev:])
    t.text = pieces[0]
    anchor = r
    for piece in pieces[1:]:
        new_r = etree.Element(qn("a:r"))
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = piece
        anchor.addnext(new_r)
        anchor = new_r
    part = pkg.slide_parts()[slide_index]
    pkg.mark_dirty(part)
    return p


# ------------------------------------------------------ set_placeholder_text


def test_set_placeholder_text_fresh_slide_roundtrip(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    new = _slide_with_title_and_body(pkg)
    idx = new["index"]

    title = set_placeholder_text(pkg, idx, "title", "Delta Model Overview")
    assert title["paragraphs"] == 1
    body_sel = next(
        p["idx"] for p in new["placeholders"] if p["type"] in ("body", "obj")
    )
    body = set_placeholder_text(
        pkg, idx, body_sel, "Point one\n\tSub point\nPoint two"
    )
    assert body["paragraphs"] == 3
    pkg.save()

    reopened = PptxPackage(work)
    info = get_slide_info(reopened, idx)
    texts = {}
    for ph in info["placeholders"]:
        elem = _shape_elem(reopened, idx, ph["id"])
        texts[ph["id"]] = shape_text(elem)
    assert "Delta Model Overview" in texts.values()
    assert "Point one\nSub point\nPoint two" in texts.values()
    # The tab became an outline level, not literal text.
    body_elem = _shape_elem(reopened, idx, body["shape_id"])
    paras = txbody_paragraphs(body_elem)
    ppr = paras[1].find(qn("a:pPr"))
    assert ppr is not None and ppr.get("lvl") == "1"
    assert paras[0].find(qn("a:pPr")) is None  # level 0 writes no pPr


def test_set_placeholder_text_paragraph_dicts(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    new = _slide_with_title_and_body(pkg)
    body_sel = next(
        p["idx"] for p in new["placeholders"] if p["type"] in ("body", "obj")
    )
    result = set_placeholder_text(
        pkg,
        new["index"],
        body_sel,
        paragraphs=[
            {"text": "Alpha", "level": 0},
            {"text": "Beta", "level": 2},
            {"text": "", "level": 0},
        ],
    )
    assert result["paragraphs"] == 3
    pkg.save()
    reopened = PptxPackage(work)
    elem = _shape_elem(reopened, new["index"], result["shape_id"])
    paras = txbody_paragraphs(elem)
    assert [paragraph_text(p) for p in paras] == ["Alpha", "Beta", ""]
    assert paras[1].find(qn("a:pPr")).get("lvl") == "2"


def test_set_placeholder_text_selector_errors(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    new = _slide_with_title_and_body(pkg)
    idx = new["index"]
    with pytest.raises(TargetNotFound, match="placeholders present"):
        set_placeholder_text(pkg, idx, "sldImg", "x")
    with pytest.raises(PptMcpError, match="exactly one of"):
        set_placeholder_text(pkg, idx, "title")
    with pytest.raises(PptMcpError, match="exactly one of"):
        set_placeholder_text(
            pkg, idx, "title", "x", paragraphs=[{"text": "y"}]
        )
    # Force an ambiguity: clone the body placeholder under a fresh shape id.
    body_id = next(
        p["shape_id"] for p in new["placeholders"] if p["type"] in ("body", "obj")
    )
    part = pkg.slide_parts()[idx]
    elem = _shape_elem(pkg, idx, body_id)
    import copy as _copy

    clone = _copy.deepcopy(elem)
    clone.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}").set(
        "id", str(pkg.next_shape_id(part))
    )
    clone.find(f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}").set("idx", "98")
    elem.getparent().append(clone)
    ph_type = next(
        p["type"] for p in new["placeholders"] if p["shape_id"] == body_id
    )
    with pytest.raises(AmbiguousTarget, match="Address by idx"):
        set_placeholder_text(pkg, idx, ph_type, "x")


# ------------------------------------------------------------ insert_textbox


def test_insert_textbox_geometry_and_roundtrip(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    result = insert_textbox(
        pkg, 0, "Hello deck\nSecond line", 1.0, 1.5, 4.0, 1.25
    )
    geo = result["geometry"]
    assert geo["x"] == 914400
    assert geo["y"] == 1371600
    assert geo["cx"] == 3657600
    assert geo["cy"] == 1143000
    pkg.save()

    reopened = PptxPackage(work)
    info = get_slide_info(reopened, 0)
    rec = next(s for s in info["shapes"] if s["id"] == result["shape_id"])
    assert rec["type"] == "textbox"
    assert rec["geometry"]["x"] == 914400
    assert rec["geometry"]["cx_in"] == 4.0
    elem = _shape_elem(reopened, 0, result["shape_id"])
    assert shape_text(elem) == "Hello deck\nSecond line"
    # python-pptx oracle sees the same box.
    from pptx import Presentation
    from pptx.util import Emu

    slide = Presentation(str(work)).slides[0]
    box = next(s for s in slide.shapes if s.shape_id == result["shape_id"])
    assert box.left == Emu(914400)
    assert box.width == Emu(3657600)


def test_insert_textbox_unit_rules(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    # ints >= 10000 are EMU; small ints are inches; unit= overrides.
    r1 = insert_textbox(pkg, 0, "a", 914400, 914400, 914400, 914400)
    assert r1["geometry"]["x"] == 914400 and r1["geometry"]["cx"] == 914400
    r2 = insert_textbox(pkg, 0, "b", 1, 1, 2, 1)
    assert r2["geometry"]["x"] == 914400 and r2["geometry"]["cx"] == 1828800
    r3 = insert_textbox(pkg, 0, "c", 100, 100, 200, 100, unit="emu")
    assert r3["geometry"]["x"] == 100 and r3["geometry"]["cx"] == 200
    with pytest.raises(PptMcpError, match="unit must be"):
        insert_textbox(pkg, 0, "d", 1, 1, 1, 1, unit="cm")
    with pytest.raises(PptMcpError, match="positive"):
        insert_textbox(pkg, 0, "e", 1, 1, 0, 1)
    pkg.save()


def test_insert_textbox_with_formatting(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    result = insert_textbox(
        pkg,
        0,
        "Styled",
        0.5,
        0.5,
        3.0,
        0.5,
        font="Georgia",
        size_pt=24,
        bold=True,
        color="#1F4E79",
        align="center",
    )
    pkg.save()
    reopened = PptxPackage(work)
    elem = _shape_elem(reopened, 0, result["shape_id"])
    rpr = elem.find(f"{qn('p:txBody')}/{qn('a:p')}/{qn('a:r')}/{qn('a:rPr')}")
    assert rpr.get("sz") == "2400" and rpr.get("b") == "1"
    assert rpr.find(qn("a:latin")).get("typeface") == "Georgia"
    assert rpr.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}").get("val") == "1F4E79"
    ppr = elem.find(f"{qn('p:txBody')}/{qn('a:p')}/{qn('a:pPr')}")
    assert ppr.get("algn") == "ctr"


# --------------------------------------------------------------- format_text


def test_format_range_spanning_fragmented_runs(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "alliance dynamics matter", 1, 1, 4, 1)
    # Fragment into "allian" + "ce dyn" + "amics matter".
    _fragment_first_run(pkg, 0, box["shape_id"], (6, 12))
    elem = _shape_elem(pkg, 0, box["shape_id"])
    assert len(txbody_paragraphs(elem)[0].findall(qn("a:r"))) == 3

    # Bold chars [9, 17) = "namics ma"[no: 'dynamics ' region] across runs
    # 2 and 3; text must be unchanged and bolding must not bleed.
    result = format_text(
        pkg, 0, box["shape_id"], paragraph=0, start=9, end=17, bold=True
    )
    assert result["runs_formatted"] >= 2
    assert shape_text(elem) == "alliance dynamics matter"
    bold_chars = []
    pos = 0
    for r in txbody_paragraphs(elem)[0].findall(qn("a:r")):
        text = r.find(qn("a:t")).text or ""
        rpr = r.find(qn("a:rPr"))
        if rpr is not None and rpr.get("b") == "1":
            bold_chars.append((pos, pos + len(text)))
        pos += len(text)
    covered = set()
    for lo, hi in bold_chars:
        covered.update(range(lo, hi))
    assert covered == set(range(9, 17))
    pkg.save()


def test_format_whole_shape_and_theme_color_survives_save(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "Theme colored\nSecond", 1, 1, 4, 1)
    result = format_text(
        pkg, 0, box["shape_id"], color="accent1", size_pt=20, italic=True
    )
    assert result["paragraphs"] == 2
    pkg.save()

    reopened = PptxPackage(work)
    elem = _shape_elem(reopened, 0, box["shape_id"])
    schemes = elem.findall(
        f".//{qn('a:rPr')}/{qn('a:solidFill')}/{qn('a:schemeClr')}"
    )
    assert schemes and all(s.get("val") == "accent1" for s in schemes)
    for rpr in elem.iter(qn("a:rPr")):
        assert rpr.get("sz") == "2000" and rpr.get("i") == "1"


def test_format_rpr_children_written_in_schema_order(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "ordering", 1, 1, 3, 1, font="Calibri")
    # The run already has a:latin; adding a fill must insert BEFORE it.
    format_text(pkg, 0, box["shape_id"], color="FF0000")
    elem = _shape_elem(pkg, 0, box["shape_id"])
    rpr = elem.find(f"{qn('p:txBody')}/{qn('a:p')}/{qn('a:r')}/{qn('a:rPr')}")
    tags = [etree.QName(c).localname for c in rpr]
    assert tags.index("solidFill") < tags.index("latin")
    pkg.save()


def test_format_text_line_spacing_and_align(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "spaced out", 1, 1, 3, 1)
    format_text(pkg, 0, box["shape_id"], line_spacing=1.5, align="right")
    elem = _shape_elem(pkg, 0, box["shape_id"])
    ppr = elem.find(f"{qn('p:txBody')}/{qn('a:p')}/{qn('a:pPr')}")
    assert ppr.get("algn") == "r"
    assert ppr.find(f"{qn('a:lnSpc')}/{qn('a:spcPct')}").get("val") == "150000"
    assert list(ppr)[0].tag == qn("a:lnSpc")  # lnSpc is rank-first in pPr
    format_text(pkg, 0, box["shape_id"], line_spacing=18)
    assert ppr.find(f"{qn('a:lnSpc')}/{qn('a:spcPts')}").get("val") == "1800"
    pkg.save()


def test_format_text_errors(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "short", 1, 1, 3, 1)
    sid = box["shape_id"]
    with pytest.raises(PptMcpError, match="nothing to do"):
        format_text(pkg, 0, sid)
    with pytest.raises(PptMcpError, match="per paragraph"):
        format_text(pkg, 0, sid, start=0, end=2, bold=True)
    with pytest.raises(PptMcpError, match="together"):
        format_text(pkg, 0, sid, paragraph=0, start=0, bold=True)
    with pytest.raises(TargetNotFound, match="out of bounds"):
        format_text(pkg, 0, sid, paragraph=0, start=0, end=99, bold=True)
    with pytest.raises(TargetNotFound, match="paragraph 5 out of range"):
        format_text(pkg, 0, sid, paragraph=5, bold=True)
    with pytest.raises(PptMcpError, match="theme token"):
        format_text(pkg, 0, sid, color="chartreuse")
    with pytest.raises(TargetNotFound, match="no shape"):
        format_text(pkg, 0, 99999, bold=True)


def test_format_text_refuses_tables(tmp_path):
    work = _work_copy("pmr_tables.pptx", tmp_path)
    pkg = PptxPackage(work)
    sid = _add_table(pkg, 0, [["Head A", "Head B"], ["one", "two"]])
    assert list_elements(pkg, "tables", scope=0)["count"] == 1
    with pytest.raises(UnsupportedStructure, match="table"):
        format_text(pkg, 0, sid, bold=True)


# --------------------------------------------------------------- set_bullets


def test_set_bullets_char_roundtrip(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "One\nTwo\nThree", 1, 1, 4, 2)
    result = set_bullets(
        pkg, 0, box["shape_id"], "char", char="§", char_font="Wingdings"
    )
    assert result["paragraphs_updated"] == 3
    pkg.save()
    reopened = PptxPackage(work)
    elem = _shape_elem(reopened, 0, box["shape_id"])
    for p in txbody_paragraphs(elem):
        ppr = p.find(qn("a:pPr"))
        assert ppr.find(qn("a:buChar")).get("char") == "§"
        assert ppr.find(qn("a:buFont")).get("typeface") == "Wingdings"
        tags = [etree.QName(c).localname for c in ppr]
        assert tags.index("buFont") < tags.index("buChar")  # schema order


def test_set_bullets_autonum_none_override_and_options(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "First\nSecond\nThird", 1, 1, 4, 2)
    sid = box["shape_id"]
    set_bullets(pkg, 0, sid, "autonum", num_type="romanUcPeriod", start_at=3)
    # Per-paragraph override: paragraph 1 alone flips to no bullet.
    set_bullets(pkg, 0, sid, "none", paragraphs=1)
    elem = _shape_elem(pkg, 0, sid)
    paras = txbody_paragraphs(elem)
    an = paras[0].find(f"{qn('a:pPr')}/{qn('a:buAutoNum')}")
    assert an.get("type") == "romanUcPeriod" and an.get("startAt") == "3"
    assert paras[1].find(f"{qn('a:pPr')}/{qn('a:buNone')}") is not None
    assert paras[1].find(f"{qn('a:pPr')}/{qn('a:buAutoNum')}") is None
    assert paras[2].find(f"{qn('a:pPr')}/{qn('a:buAutoNum')}") is not None

    set_bullets(
        pkg, 0, sid, "char", paragraphs=[2], size_pct=80, color="accent2", level=1
    )
    ppr = paras[2].find(qn("a:pPr"))
    assert ppr.find(qn("a:buSzPct")).get("val") == "80000"
    assert ppr.find(f"{qn('a:buClr')}/{qn('a:schemeClr')}").get("val") == "accent2"
    assert ppr.get("lvl") == "1"
    assert ppr.get("marL") == str(228600 + 457200) and ppr.get("indent") == "-228600"
    pkg.save()


def test_set_bullets_errors(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "One", 1, 1, 3, 1)
    sid = box["shape_id"]
    with pytest.raises(PptMcpError, match="style must be"):
        set_bullets(pkg, 0, sid, "dots")
    with pytest.raises(PptMcpError, match="autonumber scheme"):
        set_bullets(pkg, 0, sid, "autonum", num_type="fancyNumbers")
    with pytest.raises(PptMcpError, match="size_pct"):
        set_bullets(pkg, 0, sid, "char", size_pct=500)
    with pytest.raises(PptMcpError, match="start_at"):
        set_bullets(pkg, 0, sid, "autonum", start_at=0)
    with pytest.raises(TargetNotFound, match="out of range"):
        set_bullets(pkg, 0, sid, "char", paragraphs=[7])


# -------------------------------------------------------- search_and_replace


def _most_common_term(pkg: PptxPackage) -> str:
    """A recurring real word from the deck (length >= 5, letters only)."""
    text = get_text(pkg)["text"]
    words = collections.Counter(
        w for w in re.findall(r"[A-Za-z]{5,}", text)
    )
    for word, count in words.most_common():
        if count >= 2:
            return word
    pytest.skip("deck has no recurring term to exercise replace")


def test_replace_deck_parity_and_no_self_reference_loop(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    term = _most_common_term(pkg)
    expected = find_text(pkg, term, include_notes=False)["count"]
    assert expected >= 2

    # Replacement CONTAINS the search term: the KS4W self-referencing
    # lesson; snapshot-once right-to-left must not re-match its own output.
    result = search_and_replace(pkg, term, term + "-plus")
    assert result["total"] == expected
    assert sum(s["count"] for s in result["slides"]) == expected
    pkg.save()

    reopened = PptxPackage(work)
    assert find_text(reopened, term + "-plus", include_notes=False)["count"] == expected
    assert find_text(reopened, term + "-plus-plus", include_notes=False)["count"] == 0


def test_replace_fragmented_run_match(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "the alliance holds", 1, 1, 4, 1)
    # "alliance" spans all three fragments: "the alli" + "an" + "ce holds".
    _fragment_first_run(pkg, 0, box["shape_id"], (8, 10))
    result = search_and_replace(pkg, "alliance", "coalition", scope=0)
    assert result["total"] == 1
    elem = _shape_elem(pkg, 0, box["shape_id"])
    assert shape_text(elem) == "the coalition holds"
    pkg.save()


def test_replace_scope_regex_and_case(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    insert_textbox(pkg, 0, "Alpha BETA beta\nalpha beta", 1, 1, 4, 1)
    insert_textbox(pkg, 1, "beta elsewhere", 1, 1, 4, 1)

    # Scoped, case-insensitive literal.
    result = search_and_replace(pkg, "beta", "gamma", scope=0, match_case=False)
    assert result["total"] == 3
    assert result["slides"][0]["slide_index"] == 0
    # Slide 1 untouched by the scope.
    assert find_text(pkg, "beta elsewhere", scope=1)["count"] == 1

    # Regex with group references.
    r2 = search_and_replace(
        pkg, r"(Alpha) gamma", r"\1 delta", scope=0, regex=True
    )
    assert r2["total"] == 1
    assert find_text(pkg, "Alpha delta", scope=0)["count"] == 1

    with pytest.raises(PptMcpError, match="non-empty"):
        search_and_replace(pkg, "", "x")
    with pytest.raises(PptMcpError, match="invalid regex"):
        search_and_replace(pkg, "(unclosed", "x", regex=True)
    pkg.save()


def test_replace_reaches_table_cells(tmp_path):
    work = _work_copy("pmr_tables.pptx", tmp_path)
    pkg = PptxPackage(work)
    _add_table(
        pkg, 0, [["metric", "value"], ["metric growth", "metric spread"]]
    )
    hits = find_text(pkg, "metric", scope=0, include_notes=False)
    assert {m["where"] for m in hits["matches"]} == {"table"}
    assert hits["count"] == 3
    result = search_and_replace(pkg, "metric", "measure", scope=0)
    assert result["total"] == 3
    pkg.save()
    reopened = PptxPackage(work)
    assert find_text(reopened, "measure", scope=0)["count"] == 3
    assert find_text(reopened, "metric", scope=0)["count"] == 0


def test_replace_notes_flag(tmp_path):
    for name in ("proposal_defense.pptx", "nsu_pcsj.pptx", "unitar_final.pptx"):
        work = _work_copy(name, tmp_path)
        pkg = PptxPackage(work)
        notes_hits = [
            m
            for m in find_text(pkg, r"[A-Za-z]{5,}", regex=True)["matches"]
            if m["where"] == "notes"
        ]
        if notes_hits:
            break
    else:
        pytest.skip("no corpus deck has speaker notes text")
    term = notes_hits[0]["match"]
    all_hits = find_text(pkg, term, include_notes=True)["count"]
    slide_hits = find_text(pkg, term, include_notes=False)["count"]

    without = search_and_replace(pkg, term, term + "X")
    assert without["total"] == slide_hits

    pkg2 = PptxPackage(work)  # fresh, unmutated
    with_notes = search_and_replace(pkg2, term, term + "X", include_notes=True)
    assert with_notes["total"] == all_hits
    pkg2.save()


def test_replace_refuses_field_overlap(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "see page ", 1, 1, 4, 1)
    elem = _shape_elem(pkg, 0, box["shape_id"])
    p = elem.find(f"{qn('p:txBody')}/{qn('a:p')}")
    fld = etree.SubElement(p, qn("a:fld"))
    fld.set("id", "{1F783F0A-1111-2222-3333-444444444444}")
    fld.set("type", "slidenum")
    t = etree.SubElement(fld, qn("a:t"))
    t.text = "12"
    pkg.mark_dirty(pkg.slide_parts()[0])

    # find_text sees the cached text; replace must refuse to touch it.
    assert find_text(pkg, "page 12", scope=0)["count"] == 1
    with pytest.raises(UnsupportedStructure, match="field"):
        search_and_replace(pkg, "page 12", "page 13", scope=0)
    # A replace clear of the field works.
    ok = search_and_replace(pkg, "see page", "on slide", scope=0)
    assert ok["total"] == 1
    pkg.save()


# ------------------------------------------------------------- runmap direct


def test_runmap_visible_text_matches_paragraph_text(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    checked = 0
    for part in pkg.slide_parts():
        for p in pkg.root(part).iter(qn("a:p")):
            text, _segments = _runmap.build_map(p)
            assert text == paragraph_text(p)
            checked += 1
    assert checked > 0


def test_runmap_replace_covering_break_slot():
    p = etree.fromstring(
        f'<a:p xmlns:a="{A}"><a:r><a:t>one</a:t></a:r><a:br/>'
        f"<a:r><a:t>two</a:t></a:r></a:p>"
    )
    text, segments = _runmap.build_map(p)
    assert text == "one\ntwo"
    _runmap.replace_range(p, segments, 3, 4, " ")  # the br slot alone
    assert paragraph_text(p) == "one two"
    assert p.find(qn("a:br")) is None


def test_runmap_empty_runs_are_transparent():
    p = etree.fromstring(
        f'<a:p xmlns:a="{A}"><a:r><a:t></a:t></a:r>'
        f"<a:r><a:t>text</a:t></a:r><a:r><a:t/></a:r></a:p>"
    )
    text, segments = _runmap.build_map(p)
    assert text == "text"
    assert len(segments) == 1


# ------------------------------------------------------------------- autofit


def test_get_autofit_state_real_placeholder(tmp_path):
    work = _work_copy("nsu_pcsj.pptx", tmp_path)
    pkg = PptxPackage(work)
    state = get_autofit_state(pkg, 0)
    assert "rendering cache" in state["caveat"]
    assert state["shapes"], "slide 0 should carry text-bearing shapes"
    for rec in state["shapes"]:
        assert rec["autofit"] in ("normAutofit", "spAutoFit", "none", "inherited")
        if rec["autofit"] == "normAutofit":
            assert 0 < rec["font_scale_pct"] <= 100
        overflow = rec.get("overflow")
        if overflow is not None and overflow.get("likely_overflow") is not None:
            assert overflow["heuristic"] is True
            assert "fill_ratio" in overflow


def test_get_autofit_state_parses_both_percent_forms(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    box = insert_textbox(pkg, 0, "x" * 400, 1, 1, 2, 0.5)
    elem = _shape_elem(pkg, 0, box["shape_id"])
    bodypr = elem.find(f"{qn('p:txBody')}/{qn('a:bodyPr')}")
    norm = etree.SubElement(bodypr, qn("a:normAutofit"))
    norm.set("fontScale", "62500")
    norm.set("lnSpcReduction", "20000")
    state = get_autofit_state(pkg, 0, box["shape_id"])
    rec = state["shapes"][0]
    assert rec["autofit"] == "normAutofit"
    assert rec["font_scale_pct"] == 62.5
    assert rec["line_spacing_reduction_pct"] == 20.0
    # 400 chars in a 2x0.5 inch box overflows by any honest estimate.
    assert rec["overflow"]["likely_overflow"] is True

    norm.set("fontScale", "62.5%")  # later-edition percent string form
    state2 = get_autofit_state(pkg, 0, box["shape_id"])
    assert state2["shapes"][0]["font_scale_pct"] == 62.5


def test_get_autofit_state_shape_selector_errors(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    with pytest.raises(TargetNotFound, match="no shape"):
        get_autofit_state(pkg, 0, 99999)
