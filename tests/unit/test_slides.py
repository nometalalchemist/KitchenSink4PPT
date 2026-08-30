"""Phase 2 ops: slide CRUD + create-from-template, exercised on the real
corpus (via tmp copies; corpus files are never opened for saving) and on
synthetic decks for structures the corpus lacks (charts, sections, custom
shows, jump hyperlinks; the corpus decks carry none of these, verified
2026-08-30).

Every mutated output is saved, which runs pkg._validate_payload (dangling
rels, unresolvable sldIds); most tests also reopen the file and cross-check
against the python-pptx oracle.
"""

from __future__ import annotations

import io
import posixpath
import shutil
import uuid
import zipfile
from pathlib import Path

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import (
    PptMcpError,
    TargetNotFound,
    UnsupportedStructure,
)
from kitchensink4ppt.core.package import (
    NSMAP,
    PRESENTATION_PART,
    PptxPackage,
    qn,
    rels_name,
    resolve_target,
)
from kitchensink4ppt.ops.slides import (
    CT_PRESENTATION_MAIN,
    RT_CHART,
    RT_CHART_COLORS,
    RT_CHART_STYLE,
    RT_NOTES_SLIDE,
    RT_PACKAGE,
    RT_SLIDE,
    create_presentation,
    delete_slide,
    duplicate_slide,
    insert_slide,
    move_slide,
    reorder_slides,
    set_slide_hidden,
)

CORPUS = Path(__file__).parents[1] / "corpus"


def _work_copy(name: str, tmp_path: Path) -> Path:
    src = CORPUS / name
    if not src.exists():
        pytest.skip(f"corpus file missing: {name}")
    work = tmp_path / name
    shutil.copy2(src, work)
    return work


def _oracle_slide_count(path: Path) -> int:
    from pptx import Presentation

    return len(list(Presentation(str(path)).slides))


def _heaviest_slide_index(pkg: PptxPackage) -> int:
    parts = pkg.slide_parts()

    def weight(part: str) -> int:
        root = pkg.root(part)
        return sum(1 for _ in root.iter(qn("p:sp"))) + 10 * sum(
            1 for _ in root.iter(qn("p:grpSp"))
        )

    return max(range(len(parts)), key=lambda i: weight(parts[i]))


def _creation_ids(pkg: PptxPackage, part: str) -> set[str]:
    root = pkg.root(part)
    return {
        el.get("id")
        for el in root.iter(qn("a16:creationId"), qn("p14:creationId"))
    }


def _rel_targets(pkg: PptxPackage, part: str, rel_type: str) -> list[str]:
    name = rels_name(part)
    if not pkg.has_part(name):
        return []
    return [
        resolve_target(part, rel.get("Target", ""))
        for rel in pkg.root(name)
        if rel.get("Type") == rel_type and rel.get("TargetMode") != "External"
    ]


# --------------------------------------------------------- synthetic fixtures


def _add_sections(pkg: PptxPackage, split_at: int) -> None:
    """Two sections, A = slides[:split_at], B = the rest."""
    pres = pkg.presentation()
    ids = [e.get("id") for e in pres.find(qn("p:sldIdLst"))]
    ext_lst = pres.find(qn("p:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(pres, qn("p:extLst"))
    ext = etree.SubElement(ext_lst, qn("p:ext"))
    ext.set("uri", "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}")
    p14 = NSMAP["p14"]
    sec_lst = etree.SubElement(ext, f"{{{p14}}}sectionLst", nsmap={"p14": p14})
    for name, chunk in (("A", ids[:split_at]), ("B", ids[split_at:])):
        sec = etree.SubElement(sec_lst, f"{{{p14}}}section")
        sec.set("name", name)
        sec.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        lst = etree.SubElement(sec, f"{{{p14}}}sldIdLst")
        for sid in chunk:
            el = etree.SubElement(lst, f"{{{p14}}}sldId")
            el.set("id", sid)
    pkg.mark_dirty(PRESENTATION_PART)


def _section_state(pkg: PptxPackage) -> dict[str, list[str]]:
    """{section name: [slide ids]} in listed order."""
    pres = pkg.presentation()
    out: dict[str, list[str]] = {}
    for ext in pres.iter(qn("p:ext")):
        sec_lst = ext.find(qn("p14:sectionLst"))
        if sec_lst is None:
            continue
        for sec in sec_lst.findall(qn("p14:section")):
            ids = [
                e.get("id")
                for e in sec.iter(qn("p14:sldId"))
            ]
            out[sec.get("name")] = ids
    return out


def _assert_section_invariants(pkg: PptxPackage) -> None:
    """Every slide in exactly one section; concatenated section lists (in
    listed order) equal deck order."""
    state = _section_state(pkg)
    if not state:
        return
    deck = [e.get("id") for e in pkg.presentation().find(qn("p:sldIdLst"))]
    flat = [sid for ids in state.values() for sid in ids]
    assert sorted(flat) == sorted(set(flat)), "a slide id appears twice"
    assert set(flat) == set(deck), "section membership does not cover the deck"
    assert flat == deck, "section lists do not mirror deck order"


def _add_custom_show(pkg: PptxPackage, rids: list[str]) -> None:
    cust = etree.Element(qn("p:custShowLst"))
    show = etree.SubElement(cust, qn("p:custShow"))
    show.set("name", "Demo")
    show.set("id", "0")
    lst = etree.SubElement(show, qn("p:sldLst"))
    for rid in rids:
        el = etree.SubElement(lst, qn("p:sld"))
        el.set(qn("r:id"), rid)
    pkg._insert_presentation_child(cust)
    pkg.mark_dirty(PRESENTATION_PART)


def _slide_rids(pkg: PptxPackage) -> list[str]:
    return [
        e.get(qn("r:id"))
        for e in pkg.presentation().find(qn("p:sldIdLst"))
    ]


def _add_jump_hyperlink(pkg: PptxPackage, from_part: str, to_part: str) -> str:
    """Shape-level jump hyperlink (a:hlinkClick on the first sp's cNvPr)."""
    target = posixpath.relpath(to_part, posixpath.dirname(from_part))
    rid = pkg.add_relationship(from_part, RT_SLIDE, target)
    sp = pkg.root(from_part).find(
        f"{qn('p:cSld')}/{qn('p:spTree')}/{qn('p:sp')}"
    )
    cnv = sp.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}")
    hl = etree.SubElement(cnv, qn("a:hlinkClick"))
    hl.set(qn("r:id"), rid)
    hl.set("action", "ppaction://hlinksldjump")
    pkg.mark_dirty(from_part)
    return rid


_C = NSMAP["c"]
_MIN_CHART = (
    f'<c:chartSpace xmlns:c="{_C}" xmlns:a="{NSMAP["a"]}" xmlns:r="{NSMAP["r"]}">'
    "<c:chart><c:plotArea><c:layout/>"
    '<c:pieChart><c:varyColors val="1"/>'
    '<c:ser><c:idx val="0"/><c:order val="0"/>'
    "<c:val><c:numRef><c:f>Sheet1!$B$1:$B$3</c:f>"
    '<c:numCache><c:ptCount val="3"/>'
    '<c:pt idx="0"><c:v>1</c:v></c:pt>'
    '<c:pt idx="1"><c:v>2</c:v></c:pt>'
    '<c:pt idx="2"><c:v>3</c:v></c:pt>'
    "</c:numCache></c:numRef></c:val></c:ser>"
    "</c:pieChart></c:plotArea>"
    '<c:plotVisOnly val="1"/></c:chart>'
    '<c:externalData r:id="rId3"><c:autoUpdate val="0"/></c:externalData>'
    "</c:chartSpace>"
).encode()


def _attach_chart(pkg: PptxPackage, slide_part: str) -> str:
    """Wire a minimal chart part + colors/style companions + embedded xlsx
    onto a slide, mirroring PowerPoint's package topology (the copy-machinery
    test target; chart internals are not COM-validated here)."""
    chart = "ppt/charts/chart1.xml"
    pkg.add_part_with_content_type(
        chart,
        _MIN_CHART,
        "application/vnd.openxmlformats-officedocument.drawingml.chart+xml",
    )
    pkg.add_part_with_content_type(
        "ppt/charts/colors1.xml",
        b'<?xml version="1.0"?><cs:colorStyle xmlns:cs="http://schemas.microsoft.com/office/drawing/2012/chartStyle" meth="cycle" id="10"/>',
        "application/vnd.ms-office.chartcolorstyle+xml",
    )
    pkg.add_part_with_content_type(
        "ppt/charts/style1.xml",
        b'<?xml version="1.0"?><cs:chartStyle xmlns:cs="http://schemas.microsoft.com/office/drawing/2012/chartStyle" id="201"/>',
        "application/vnd.ms-office.chartstyle+xml",
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("stub.txt", "opaque workbook stand-in")
    xlsx = "ppt/embeddings/Microsoft_Excel_Worksheet1.xlsx"
    pkg.add_part_with_content_type(
        xlsx,
        buf.getvalue(),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    assert pkg.add_relationship(chart, RT_CHART_COLORS, "colors1.xml") == "rId1"
    assert pkg.add_relationship(chart, RT_CHART_STYLE, "style1.xml") == "rId2"
    assert (
        pkg.add_relationship(chart, RT_PACKAGE, "../embeddings/Microsoft_Excel_Worksheet1.xlsx")
        == "rId3"
    )
    rid = pkg.add_relationship(slide_part, RT_CHART, "../charts/chart1.xml")

    sp_tree = pkg.root(slide_part).find(f"{qn('p:cSld')}/{qn('p:spTree')}")
    gf = etree.SubElement(sp_tree, qn("p:graphicFrame"))
    nv = etree.SubElement(gf, qn("p:nvGraphicFramePr"))
    cnv = etree.SubElement(nv, qn("p:cNvPr"))
    cnv.set("id", str(pkg.next_shape_id(slide_part)))
    cnv.set("name", "Chart 1")
    etree.SubElement(nv, qn("p:cNvGraphicFramePr"))
    etree.SubElement(nv, qn("p:nvPr"))
    xfrm = etree.SubElement(gf, qn("p:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", "0")
    off.set("y", "0")
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", "3000000")
    ext.set("cy", "2000000")
    graphic = etree.SubElement(gf, qn("a:graphic"))
    gdata = etree.SubElement(graphic, qn("a:graphicData"))
    gdata.set("uri", _C)
    cchart = etree.SubElement(gdata, f"{{{_C}}}chart", nsmap={"c": _C})
    cchart.set(qn("r:id"), rid)
    pkg.mark_dirty(slide_part)
    return chart


# ---------------------------------------------------------------- insert_slide


def test_insert_slide_at_end_and_position(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    before = pkg.slide_parts()
    res = insert_slide(pkg, 0)
    assert res["index"] == len(before)
    assert res["placeholders"], "layout 0 should yield cloneable placeholders"
    res2 = insert_slide(pkg, 1, position=0)
    assert res2["index"] == 0
    pkg.save(do_backup=False)

    pkg2 = PptxPackage(doc)
    parts = pkg2.slide_parts()
    assert len(parts) == len(before) + 2
    assert parts[0] == res2["part"]
    assert parts[-1] == res["part"]
    assert _oracle_slide_count(doc) == len(before) + 2
    # placeholder skeleton: cloned sp elements carry p:ph, no latent types
    ph_path = f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}"
    phs = [
        sp.find(ph_path)
        for sp in pkg2.root(res["part"]).iter(qn("p:sp"))
    ]
    assert phs and all(ph is not None for ph in phs)
    assert all(ph.get("type", "obj") not in ("dt", "ftr", "sldNum") for ph in phs)


def test_insert_slide_by_layout_name(make_deck):
    from kitchensink4ppt.ops.slides import _layouts

    doc = make_deck()
    pkg = PptxPackage(doc)
    layouts = _layouts(pkg)
    part, name = layouts[1]
    assert name, "python-pptx default layouts are named"
    res = insert_slide(pkg, name)
    assert res["layout"] == part
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_insert_slide_bad_layout_refused(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    with pytest.raises(TargetNotFound, match="available layouts"):
        insert_slide(pkg, "No Such Layout Anywhere")
    with pytest.raises(TargetNotFound, match="out of range"):
        insert_slide(pkg, 999)


# ------------------------------------------------------------- duplicate_slide


def test_duplicate_heaviest_proposal_slide_and_independence(tmp_path):
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    n_before = len(pkg.slide_parts())
    idx = _heaviest_slide_index(pkg)
    src_part = pkg.slide_parts()[idx]
    src_ids = _creation_ids(pkg, src_part)
    src_notes = _rel_targets(pkg, src_part, RT_NOTES_SLIDE)

    res = duplicate_slide(pkg, idx)
    assert res["index"] == idx + 1
    new_part = res["part"]
    assert new_part != src_part
    pkg.save(do_backup=False)

    pkg2 = PptxPackage(work)
    parts = pkg2.slide_parts()
    assert len(parts) == n_before + 1
    assert parts[idx + 1] == new_part
    assert _oracle_slide_count(work) == n_before + 1

    # creationId GUIDs regenerated 1:1, never shared with the source
    new_ids = _creation_ids(pkg2, new_part)
    assert len(new_ids) == len(src_ids)
    assert not (new_ids & src_ids)

    # notesSlide, if the source had one, was deep-copied not shared
    if src_notes:
        clone_notes = _rel_targets(pkg2, new_part, RT_NOTES_SLIDE)
        assert clone_notes and clone_notes[0] != src_notes[0]

    # independence: editing the clone leaves the source part byte-identical
    src_bytes = pkg2.raw_part(src_part)
    cnv = pkg2.root(new_part).find(
        f"{qn('p:cSld')}/{qn('p:spTree')}/{qn('p:nvGrpSpPr')}/{qn('p:cNvPr')}"
    )
    cnv.set("name", "edited-clone")
    pkg2.mark_dirty(new_part)
    pkg2.save(do_backup=False)

    pkg3 = PptxPackage(work)
    assert pkg3.raw_part(src_part) == src_bytes
    assert b"edited-clone" in pkg3.raw_part(new_part)


def test_repeated_duplication_no_duplicate_creation_ids(tmp_path):
    """The repeated-duplication repair-prompt scenario: three clones of the
    same slide must carry pairwise-disjoint creationId GUID sets."""
    work = _work_copy("proposal_defense.pptx", tmp_path)
    pkg = PptxPackage(work)
    idx = _heaviest_slide_index(pkg)
    parts = [pkg.slide_parts()[idx]]
    for _ in range(3):
        parts.append(duplicate_slide(pkg, idx)["part"])
    pkg.save(do_backup=False)
    pkg2 = PptxPackage(work)
    id_sets = [_creation_ids(pkg2, p) for p in parts]
    if not any(id_sets):
        pytest.skip("deck carries no creationId extensions (synthetic corpus)")
    for i in range(len(id_sets)):
        for j in range(i + 1, len(id_sets)):
            assert not (id_sets[i] & id_sets[j]), (parts[i], parts[j])


def test_duplicate_slide_with_notes_deep_copies(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    src_part = pkg.slide_parts()[0]  # title slide carries notes
    src_notes = _rel_targets(pkg, src_part, RT_NOTES_SLIDE)
    assert src_notes, "make_deck slide 0 should have speaker notes"

    res = duplicate_slide(pkg, 0)
    new_part = res["part"]
    clone_notes = _rel_targets(pkg, new_part, RT_NOTES_SLIDE)
    assert clone_notes and clone_notes[0] != src_notes[0]
    assert clone_notes[0] in res["copied_parts"]
    # the clone's notesSlide back-rel targets the CLONE, not the source
    back = _rel_targets(pkg, clone_notes[0], RT_SLIDE)
    assert back == [new_part]
    # notesMaster stays shared
    src_master = _rel_targets(pkg, src_notes[0], RT_NOTES_SLIDE.replace("notesSlide", "notesMaster"))
    clone_master = _rel_targets(pkg, clone_notes[0], RT_NOTES_SLIDE.replace("notesSlide", "notesMaster"))
    assert src_master == clone_master
    pkg.save(do_backup=False)
    assert _oracle_slide_count(doc) == 7  # 6 + 1


def test_duplicate_chart_slide_deep_copies_chart_family(make_deck):
    """The corpus carries no chart-bearing deck (verified 2026-08-30), so the
    chart family topology is synthesized here: chart + colors + style +
    embedded xlsx, graphicFrame r:id kept via the cloned rels file."""
    doc = make_deck()
    pkg = PptxPackage(doc)
    slide_part = pkg.slide_parts()[1]
    _attach_chart(pkg, slide_part)
    pkg.save(do_backup=False)

    pkg = PptxPackage(doc)
    res = duplicate_slide(pkg, 1)
    new_part = res["part"]
    copied = set(res["copied_parts"])
    assert "ppt/charts/chart2.xml" in copied
    assert "ppt/charts/colors2.xml" in copied
    assert "ppt/charts/style2.xml" in copied
    assert "ppt/embeddings/Microsoft_Excel_Worksheet2.xlsx" in copied
    # original chart rels untouched, clone points at the copies
    assert _rel_targets(pkg, slide_part, RT_CHART) == ["ppt/charts/chart1.xml"]
    assert _rel_targets(pkg, new_part, RT_CHART) == ["ppt/charts/chart2.xml"]
    # clone chart's companions resolve to the cloned parts
    assert _rel_targets(pkg, "ppt/charts/chart2.xml", RT_PACKAGE) == [
        "ppt/embeddings/Microsoft_Excel_Worksheet2.xlsx"
    ]
    # graphicFrame r:id in the clone XML still resolves through its OWN rels
    root = pkg.root(new_part)
    cchart = root.find(f".//{{{_C}}}chart")
    rid = cchart.get(qn("r:id"))
    assert resolve_target(new_part, "../charts/chart2.xml") == pkg.relationship_target(new_part, rid)
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_duplicate_by_slide_id_and_position(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    entries = pkg.presentation().find(qn("p:sldIdLst"))
    sid = int(entries[2].get("id"))
    res = duplicate_slide(pkg, {"slide_id": sid}, position=0)
    assert res["index"] == 0
    assert res["source_slide_id"] == sid
    pkg.save(do_backup=False)
    pkg2 = PptxPackage(doc)
    assert pkg2.slide_parts()[0] == res["part"]
    with pytest.raises(TargetNotFound, match="no slide with slide_id"):
        duplicate_slide(pkg2, {"slide_id": 999999})


# ---------------------------------------------------------------- delete_slide


def test_delete_slide_gc_notes_and_media(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    # slide index 3 is the picture slide with speaker notes
    part = pkg.slide_parts()[3]
    notes = _rel_targets(pkg, part, RT_NOTES_SLIDE)
    images = [
        t for t in (
            resolve_target(part, rel.get("Target", ""))
            for rel in pkg.root(rels_name(part))
            if rel.get("TargetMode") != "External"
        )
        if t.startswith("ppt/media/")
    ]
    assert notes and images, "expected the picture+notes slide at index 3"
    n_before = len(pkg.slide_parts())

    res = delete_slide(pkg, 3)
    assert res["deleted_part"] == part
    assert notes[0] in res["gc_parts"]
    assert images[0] in res["gc_parts"]
    assert not pkg.has_part(part)
    assert not pkg.has_part(rels_name(part))
    assert not pkg.has_part(notes[0])
    assert not pkg.has_part(images[0])
    # content-type override for the deleted slide is gone
    ct = pkg.root("[Content_Types].xml")
    assert not any(
        o.get("PartName") == "/" + part for o in ct.findall(qn("ct:Override"))
    )
    pkg.save(do_backup=False)  # runs _validate_payload: no dangling rels
    assert _oracle_slide_count(doc) == n_before - 1


def test_delete_slide_keeps_shared_media(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    part = pkg.slide_parts()[3]
    images = [
        t for t in (
            resolve_target(part, rel.get("Target", ""))
            for rel in pkg.root(rels_name(part))
            if rel.get("TargetMode") != "External"
        )
        if t.startswith("ppt/media/")
    ]
    duplicate_slide(pkg, 3)  # media rel shared with the clone
    res = delete_slide(pkg, 3)
    assert images[0] not in res["gc_parts"]
    assert pkg.has_part(images[0]), "shared media must survive the refcount"
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_delete_slide_neuters_jump_hyperlinks(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    parts = pkg.slide_parts()
    rid = _add_jump_hyperlink(pkg, parts[0], parts[2])
    pkg.save(do_backup=False)

    pkg = PptxPackage(doc)
    res = delete_slide(pkg, 2)
    flagged = res["flagged_hyperlinks"]
    assert flagged == [
        {"part": parts[0], "rid": rid, "hyperlinks_removed": 1}
    ]
    # rel gone, element gone
    assert not any(
        rel.get("Id") == rid for rel in pkg.root(rels_name(parts[0]))
    )
    assert not any(
        el.get(qn("r:id")) == rid
        for el in pkg.root(parts[0]).iter(qn("a:hlinkClick"))
    )
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_delete_slide_prunes_custom_shows(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    rids = _slide_rids(pkg)
    _add_custom_show(pkg, [rids[0], rids[2]])
    pkg.save(do_backup=False)

    pkg = PptxPackage(doc)
    res = delete_slide(pkg, 0)
    assert res["custom_shows"]["entries_removed"] == 1
    assert res["custom_shows"]["shows_removed"] == []
    assert pkg.presentation().find(qn("p:custShowLst")) is not None
    # deleting the show's last member removes the emptied show and list
    res2 = delete_slide(pkg, {"slide_id": int(
        [e.get("id") for e in pkg.presentation().find(qn("p:sldIdLst"))
         if e.get(qn("r:id")) == rids[2]][0]
    )})
    assert res2["custom_shows"]["entries_removed"] == 1
    assert res2["custom_shows"]["shows_removed"] == ["Demo"]
    assert pkg.presentation().find(qn("p:custShowLst")) is None
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_delete_slide_updates_sections(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    _add_sections(pkg, 2)
    pkg.save(do_backup=False)

    pkg = PptxPackage(doc)
    res = delete_slide(pkg, 0)
    assert res["section_entries_removed"] == 1
    _assert_section_invariants(pkg)
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_delete_refuses_last_slide(make_deck):
    doc = make_deck(extra_slides=0)
    pkg = PptxPackage(doc)
    while len(pkg.slide_parts()) > 1:
        delete_slide(pkg, len(pkg.slide_parts()) - 1)
    with pytest.raises(UnsupportedStructure, match="last remaining slide"):
        delete_slide(pkg, 0)
    pkg.save(do_backup=False)
    assert _oracle_slide_count(doc) == 1


def test_delete_out_of_range_message(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    n = len(pkg.slide_parts())
    with pytest.raises(TargetNotFound, match=f"presentation has {n} slides"):
        delete_slide(pkg, n)


# --------------------------------------------------------- reorder, move, hide


def test_reorder_full_permutation(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    before = pkg.slide_parts()
    n = len(before)
    order = list(range(1, n)) + [0]  # rotate left
    res = reorder_slides(pkg, order)
    assert pkg.slide_parts() == [before[i] for i in order]
    assert len(res["slide_ids"]) == n
    pkg.save(do_backup=False)

    pkg2 = PptxPackage(doc)
    assert pkg2.slide_parts() == [before[i] for i in order]
    assert _oracle_slide_count(doc) == n


def test_reorder_rejects_non_permutation(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    with pytest.raises(PptMcpError, match="permutation"):
        reorder_slides(pkg, [0, 0, 1])
    with pytest.raises(PptMcpError, match="permutation"):
        reorder_slides(pkg, [0])


def test_reorder_keeps_sections_consistent(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    _add_sections(pkg, 3)
    n = len(pkg.slide_parts())
    order = [n - 1] + list(range(n - 1))  # last slide to the front
    res = reorder_slides(pkg, order)
    _assert_section_invariants(pkg)
    assert isinstance(res["section_moves"], list)
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_move_slide(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    before = pkg.slide_parts()
    res = move_slide(pkg, 0, 2)
    assert res["from"] == 0 and res["to"] == 2
    assert pkg.slide_parts() == [before[1], before[2], before[0]] + before[3:]
    with pytest.raises(TargetNotFound, match="out of range"):
        move_slide(pkg, 0, 99)
    pkg.save(do_backup=False)
    PptxPackage(doc)


def test_move_slide_joins_destination_section(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    _add_sections(pkg, 3)
    deck_ids = [e.get("id") for e in pkg.presentation().find(qn("p:sldIdLst"))]
    moved = deck_ids[0]
    move_slide(pkg, 0, len(deck_ids) - 1)  # front slide to the end -> section B
    state = _section_state(pkg)
    assert moved in state["B"] and moved not in state["A"]
    _assert_section_invariants(pkg)
    pkg.save(do_backup=False)


def test_hidden_roundtrip(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    part = pkg.slide_parts()[1]
    res = set_slide_hidden(pkg, 1, True)
    assert res["hidden"] is True
    pkg.save(do_backup=False)

    pkg2 = PptxPackage(doc)
    assert pkg2.root(part).get("show") == "0"
    set_slide_hidden(pkg2, 1, False)
    pkg2.save(do_backup=False)

    pkg3 = PptxPackage(doc)
    assert pkg3.root(part).get("show") is None


# --------------------------------------------------------- create_presentation


def _design_parts(path: Path) -> dict[str, int]:
    with zipfile.ZipFile(path) as zf:
        names = zf.namelist()
    return {
        "masters": sum(1 for n in names if n.startswith("ppt/slideMasters/") and n.endswith(".xml") and "_rels" not in n),
        "layouts": sum(1 for n in names if n.startswith("ppt/slideLayouts/") and n.endswith(".xml") and "_rels" not in n),
        "themes": sum(1 for n in names if n.startswith("ppt/theme/")),
    }


def test_create_presentation_from_potx(tmp_path):
    template = CORPUS / "conference_template.potx"
    if not template.exists():
        pytest.skip("corpus file missing: conference_template.potx")
    template_bytes = template.read_bytes()
    dest = tmp_path / "from_potx.pptx"
    res = create_presentation(dest, template)
    assert res["converted_content_type"] is True
    assert res["slides_kept"] == 0
    assert template.read_bytes() == template_bytes, "template was modified"

    pkg = PptxPackage(dest)
    assert pkg.slide_parts() == []
    ct = pkg.root("[Content_Types].xml")
    (override,) = [
        o for o in ct.findall(qn("ct:Override"))
        if o.get("PartName") == "/" + PRESENTATION_PART
    ]
    assert override.get("ContentType") == CT_PRESENTATION_MAIN
    assert _design_parts(dest) == _design_parts(template), (
        "masters/layouts/themes must survive slide stripping"
    )
    # a fresh deck accepts a new slide built from its own layouts
    r = insert_slide(pkg, 0)
    pkg.save(do_backup=False)
    assert _oracle_slide_count(dest) == 1
    assert r["slide_id"] >= 256


def test_create_presentation_from_pptx_strip_slides(tmp_path):
    src = CORPUS / "proposal_defense.pptx"
    if not src.exists():
        pytest.skip("corpus file missing: proposal_defense.pptx")
    src_bytes_hash = hash(src.read_bytes())
    dest = tmp_path / "fresh_from_deck.pptx"
    res = create_presentation(dest, src)
    assert res["slides_kept"] == 0
    assert res["converted_content_type"] is False
    assert hash(src.read_bytes()) == src_bytes_hash
    assert _design_parts(dest) == _design_parts(src)
    pkg = PptxPackage(dest)
    assert pkg.slide_parts() == []
    assert _oracle_slide_count(dest) == 0


def test_create_presentation_keep_slides(tmp_path):
    src = CORPUS / "proposal_defense.pptx"
    if not src.exists():
        pytest.skip("corpus file missing: proposal_defense.pptx")
    dest = tmp_path / "copy_with_slides.pptx"
    res = create_presentation(dest, src, keep_slides=True)
    assert res["slides_kept"] == _oracle_slide_count(dest) > 0


def test_create_presentation_refusals(tmp_path, make_deck):
    doc = make_deck()
    dest = tmp_path / "exists.pptx"
    dest.write_bytes(b"occupied")
    with pytest.raises(PptMcpError, match="already exists"):
        create_presentation(dest, doc)
    with pytest.raises(PptMcpError, match="must be a .pptx"):
        create_presentation(tmp_path / "out.potx", doc)
    with pytest.raises(Exception, match="no template file"):
        create_presentation(tmp_path / "out.pptx", tmp_path / "missing.potx")
    assert not (tmp_path / "out.pptx").exists()


def test_create_presentation_blank(tmp_path):
    dest = tmp_path / "blank.pptx"
    res = create_presentation(dest)
    assert res["slides_kept"] == 0
    pkg = PptxPackage(dest)
    insert_slide(pkg, 0)
    pkg.save(do_backup=False)
    assert _oracle_slide_count(dest) == 1
