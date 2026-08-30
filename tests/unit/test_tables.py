"""Phase 6 tables: real-deck reads (merge map), synthetic create/fill/merge/
border/style, merge-aware row/col surgery, widths, CSV/JSON round-trips.

Every mutated deck is saved, which runs pkg._validate_payload (dangling rels
and structure). Real military_brief tables (9 native tables, 25+ merge
continuations on the big ones) drive the read and surgery paths; synthetic
tables drive the build paths so CI passes on corpus stand-ins.
"""

from __future__ import annotations

import csv
import json
import shutil
from pathlib import Path

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import (
    AmbiguousTarget,
    PptMcpError,
    TargetNotFound,
    UnsupportedStructure,
)
from kitchensink4ppt.core.package import PptxPackage, qn
from kitchensink4ppt.ops import tables as tb
from kitchensink4ppt.ops.read import list_elements

CORPUS = Path(__file__).parent.parent / "corpus"


@pytest.fixture()
def mb(tmp_path):
    """Temp copy of military_brief (the 9-native-table deck) as a package."""
    p = tmp_path / "mb.pptx"
    shutil.copy(CORPUS / "military_brief.pptx", p)
    return PptxPackage(p)


@pytest.fixture()
def deck(make_deck):
    """(pkg, slide_index) with a fresh empty slide to build tables on."""
    from kitchensink4ppt.ops import slides as sl

    path = make_deck("tables.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    slide = sl.insert_slide(pkg, 0)["index"]
    return pkg, slide


def _corpus_tables(pkg):
    """[(slide_index, shape_id, rows, cols)] for every table in the deck."""
    items = list_elements(pkg, "tables")["items"]
    return [
        (i["slide_index"], i["id"], i["rows"], i["cols"]) for i in items
    ]


def _merged_corpus_table(pkg):
    """(slide_index, shape_id, info) of the corpus table with the most merge
    regions, or a skip when the corpus stand-in has none."""
    best = None
    for slide_index, shape_id, _r, _c in _corpus_tables(pkg):
        info = tb.get_table(pkg, slide_index, {"shape_id": shape_id})
        if best is None or len(info["merge_regions"]) > len(best[2]["merge_regions"]):
            best = (slide_index, shape_id, info)
    if best is None:
        pytest.skip("corpus deck has no tables (synthetic stand-in)")
    return best


# ------------------------------------------------------------- real-deck read


class TestGetTableReal:
    def test_reads_every_corpus_table(self, mb):
        tables = _corpus_tables(mb)
        if not tables:
            pytest.skip("corpus deck has no tables (synthetic stand-in)")
        for slide_index, shape_id, rows, cols in tables:
            info = tb.get_table(mb, slide_index, {"shape_id": shape_id})
            assert info["rows"] == rows
            assert info["cols"] == cols
            assert len(info["cells"]) == rows * cols
            assert len(info["column_widths_in"]) == cols
            assert len(info["row_heights_in"]) == rows

    def test_merge_map_consistency(self, mb):
        """Every region's continuation cells carry the h/vMerge flags and
        every flagged cell belongs to exactly one region: the merge map is a
        partition, not a guess."""
        slide_index, shape_id, info = _merged_corpus_table(mb)
        if not info["merge_regions"]:
            pytest.skip("no merged tables in this corpus")
        by_addr = {(c["row"], c["col"]): c for c in info["cells"]}
        covered = set()
        for reg in info["merge_regions"]:
            for r in range(reg["r1"], reg["r2"] + 1):
                for c in range(reg["c1"], reg["c2"] + 1):
                    assert (r, c) not in covered, "regions overlap"
                    covered.add((r, c))
                    cell = by_addr[(r, c)]
                    if (r, c) == (reg["r1"], reg["c1"]):
                        assert cell["merge"]["role"] == "origin"
                        assert cell["merge"]["rows"] == reg["r2"] - reg["r1"] + 1
                        assert cell["merge"]["cols"] == reg["c2"] - reg["c1"] + 1
                    else:
                        assert cell["merge"]["role"] == "continuation"
                        assert cell["merge"]["origin"] == [reg["r1"], reg["c1"]]
        # And no continuation exists OUTSIDE a mapped region.
        for (r, c), cell in by_addr.items():
            if cell.get("merge", {}).get("role") == "continuation":
                assert (r, c) in covered

    def test_style_guid_reported_raw_when_not_builtin(self, mb):
        tables = _corpus_tables(mb)
        if not tables:
            pytest.skip("corpus deck has no tables")
        slide_index, shape_id, _r, _c = tables[0]
        info = tb.get_table(mb, slide_index, {"shape_id": shape_id})
        if info["style_guid"] is None:
            pytest.skip("table carries no styleId")
        # Friendly name when built-in, raw braced GUID otherwise; never None
        # while a styleId exists.
        assert info["style"] is not None
        assert info["style_guid"].startswith("{")


# ---------------------------------------------------------------- addressing


class TestResolve:
    def test_no_table_refuses(self, deck):
        pkg, slide = deck
        with pytest.raises(TargetNotFound):
            tb.get_table(pkg, slide, None)

    def test_multi_table_without_selector_refuses_with_candidates(self, deck):
        pkg, slide = deck
        a = tb.create_table(pkg, slide, 2, 2, 0.5, 0.5, 3, 1)
        b = tb.create_table(pkg, slide, 2, 2, 0.5, 3.0, 3, 1)
        with pytest.raises(AmbiguousTarget) as exc:
            tb.get_table(pkg, slide, None)
        msg = str(exc.value)
        assert str(a["shape_id"]) in msg and str(b["shape_id"]) in msg
        # Index and shape id both resolve.
        assert tb.get_table(pkg, slide, 1)["shape_id"] == b["shape_id"]
        assert (
            tb.get_table(pkg, slide, {"shape_id": a["shape_id"]})["table_index"]
            == 0
        )

    def test_bad_selectors(self, deck):
        pkg, slide = deck
        tb.create_table(pkg, slide, 2, 2, 0.5, 0.5, 3, 1)
        with pytest.raises(TargetNotFound):
            tb.get_table(pkg, slide, 99)
        with pytest.raises(PptMcpError):
            tb.get_table(pkg, slide, "first")


# ------------------------------------------------------- create + fill + save


class TestCreate:
    def test_create_fill_and_validate(self, deck):
        pkg, slide = deck
        r = tb.create_table(
            pkg, slide, 3, 4, 1, 1, 8, 2,
            data=[["A", "B", "C", "D"], ["1", "2"]],
        )
        info = tb.get_table(pkg, slide, {"shape_id": r["shape_id"]})
        assert (info["rows"], info["cols"]) == (3, 4)
        texts = {(c["row"], c["col"]): c["text"] for c in info["cells"]}
        assert texts[(0, 0)] == "A" and texts[(0, 3)] == "D"
        assert texts[(1, 0)] == "1" and texts[(1, 2)] == ""
        # Grid widths sum to the frame width; both sum to 8 inches.
        assert abs(sum(info["column_widths_in"]) - 8.0) < 0.01
        assert abs(info["geometry"]["cx_in"] - 8.0) < 0.01
        pkg.save()  # runs payload validation

    def test_default_style_is_powerpoint_default(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        info = tb.get_table(pkg, slide, {"shape_id": r["shape_id"]})
        assert info["style_guid"] == tb.DEFAULT_STYLE_GUID
        assert info["flags"] == {"firstRow": True, "bandRow": True}

    def test_input_validation(self, deck):
        pkg, slide = deck
        with pytest.raises(PptMcpError):
            tb.create_table(pkg, slide, 0, 2, 1, 1, 4, 1)
        with pytest.raises(PptMcpError):
            tb.create_table(pkg, slide, 2, 2, 1, 1, -4, 1)
        with pytest.raises(PptMcpError):
            tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1, data=[["a"] * 5])

    def test_pptx_oracle_reads_created_table(self, deck):
        pkg, slide = deck
        tb.create_table(pkg, slide, 2, 3, 1, 1, 6, 1.5, data=[["x", "y", "z"]])
        path = pkg.save()
        from pptx import Presentation

        prs = Presentation(str(path))
        frames = [s for s in prs.slides[slide].shapes if s.has_table]
        assert len(frames) == 1
        t = frames[0].table
        assert len(t.rows) == 2 and len(t.columns) == 3
        assert t.cell(0, 0).text == "x" and t.cell(0, 2).text == "z"


# -------------------------------------------------------------------- merges


class TestMerge:
    def _table(self, pkg, slide):
        r = tb.create_table(
            pkg, slide, 4, 4, 1, 1, 8, 3,
            data=[["a", "b", "c", "d"], ["e", "f", "g", "h"]],
        )
        return {"shape_id": r["shape_id"]}

    def test_merge_writes_powerpoint_shape(self, deck):
        pkg, slide = deck
        sel = self._table(pkg, slide)
        out = tb.merge_cells(pkg, slide, sel, 0, 0, 1, 1)
        assert out["cells_absorbed"] == 3
        assert out["text_moved_from_cells"] == 3  # b, e, f moved into a
        info = tb.get_table(pkg, slide, sel)
        origin = next(c for c in info["cells"] if (c["row"], c["col"]) == (0, 0))
        assert origin["merge"] == {"role": "origin", "rows": 2, "cols": 2}
        assert "b" in origin["text"] and "f" in origin["text"]
        # XML-level continuation flags per the internals doc.
        tbl = tb.resolve_table(pkg, slide, sel)["tbl"]
        tr0 = tbl.findall(qn("a:tr"))[0]
        tr1 = tbl.findall(qn("a:tr"))[1]
        assert tr0.findall(qn("a:tc"))[1].get("hMerge") == "1"
        assert tr0.findall(qn("a:tc"))[1].get("rowSpan") == "2"
        assert tr1.findall(qn("a:tc"))[0].get("vMerge") == "1"
        assert tr1.findall(qn("a:tc"))[0].get("gridSpan") == "2"
        assert tr1.findall(qn("a:tc"))[1].get("hMerge") == "1"
        assert tr1.findall(qn("a:tc"))[1].get("vMerge") == "1"
        pkg.save()

    def test_overlap_refuses_with_region_named(self, deck):
        pkg, slide = deck
        sel = self._table(pkg, slide)
        tb.merge_cells(pkg, slide, sel, 0, 0, 1, 1)
        with pytest.raises(UnsupportedStructure) as exc:
            tb.merge_cells(pkg, slide, sel, 1, 1, 2, 2)
        assert "unmerge_cells" in str(exc.value)

    def test_editing_continuation_refuses(self, deck):
        pkg, slide = deck
        sel = self._table(pkg, slide)
        tb.merge_cells(pkg, slide, sel, 0, 0, 1, 1)
        with pytest.raises(UnsupportedStructure) as exc:
            tb.set_table_cells(pkg, slide, sel, [{"row": 0, "col": 1, "text": "x"}])
        assert "(0, 0)" in str(exc.value)

    def test_unmerge_roundtrip(self, deck):
        pkg, slide = deck
        sel = self._table(pkg, slide)
        tb.merge_cells(pkg, slide, sel, 0, 0, 1, 1)
        # Any cell of the region addresses it, continuation included.
        out = tb.unmerge_cells(pkg, slide, sel, 1, 1)
        assert out["cells_freed"] == 3
        info = tb.get_table(pkg, slide, sel)
        assert info["merge_regions"] == []
        assert all("merge" not in c for c in info["cells"])
        with pytest.raises(TargetNotFound):
            tb.unmerge_cells(pkg, slide, sel, 0, 0)

    def test_single_cell_refuses(self, deck):
        pkg, slide = deck
        sel = self._table(pkg, slide)
        with pytest.raises(PptMcpError):
            tb.merge_cells(pkg, slide, sel, 0, 0, 0, 0)


# ------------------------------------------------------------ row/col surgery


class TestRowColSurgery:
    def _merged(self, pkg, slide):
        r = tb.create_table(pkg, slide, 4, 4, 1, 1, 8, 3)
        sel = {"shape_id": r["shape_id"]}
        tb.merge_cells(pkg, slide, sel, 1, 1, 2, 2)  # 2x2 block mid-table
        return sel

    def test_insert_inside_vertical_span_refuses(self, deck):
        pkg, slide = deck
        sel = self._merged(pkg, slide)
        with pytest.raises(UnsupportedStructure) as exc:
            tb.insert_table_rows(pkg, slide, sel, 2)
        assert "split" in str(exc.value)
        # Boundaries are fine.
        assert tb.insert_table_rows(pkg, slide, sel, 1)["rows"] == 5
        assert tb.insert_table_rows(pkg, slide, sel, 4)["rows"] == 6

    def test_insert_inside_horizontal_span_refuses(self, deck):
        pkg, slide = deck
        sel = self._merged(pkg, slide)
        with pytest.raises(UnsupportedStructure):
            tb.insert_table_cols(pkg, slide, sel, 2)
        assert tb.insert_table_cols(pkg, slide, sel, 1)["cols"] == 5

    def test_delete_tail_shrinks_span(self, deck):
        pkg, slide = deck
        sel = self._merged(pkg, slide)
        out = tb.delete_table_rows(pkg, slide, sel, 2, 1)  # region rows 1..2
        assert out["rows"] == 3
        info = tb.get_table(pkg, slide, sel)
        assert info["merge_regions"] == [{"r1": 1, "c1": 1, "r2": 1, "c2": 2}]
        pkg.save()

    def test_delete_origin_with_surviving_tail_refuses(self, deck):
        pkg, slide = deck
        sel = self._merged(pkg, slide)
        with pytest.raises(UnsupportedStructure) as exc:
            tb.delete_table_rows(pkg, slide, sel, 1, 1)
        assert "unmerge" in str(exc.value)

    def test_delete_whole_region_with_rows(self, deck):
        pkg, slide = deck
        sel = self._merged(pkg, slide)
        out = tb.delete_table_rows(pkg, slide, sel, 1, 2)
        assert out["rows"] == 2
        assert tb.get_table(pkg, slide, sel)["merge_regions"] == []

    def test_col_widths_shift_vs_rescale(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 4, 1, 1, 8, 1.5)
        sel = {"shape_id": r["shape_id"]}
        before = tb.get_table(pkg, slide, sel)
        assert abs(sum(before["column_widths_in"]) - 8.0) < 0.01
        tb.insert_table_cols(pkg, slide, sel, 4, widths="rescale")
        mid = tb.get_table(pkg, slide, sel)
        assert mid["cols"] == 5
        assert abs(sum(mid["column_widths_in"]) - 8.0) < 0.01  # total kept
        tb.insert_table_cols(pkg, slide, sel, 0, widths="shift")
        after = tb.get_table(pkg, slide, sel)
        assert after["cols"] == 6
        assert sum(after["column_widths_in"]) > 8.5  # table grew
        assert abs(after["geometry"]["cx_in"] - sum(after["column_widths_in"])) < 0.02

    def test_delete_cols_rescale_keeps_total(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 4, 1, 1, 8, 1.5)
        sel = {"shape_id": r["shape_id"]}
        tb.delete_table_cols(pkg, slide, sel, 1, 2, widths="rescale")
        info = tb.get_table(pkg, slide, sel)
        assert info["cols"] == 2
        assert abs(sum(info["column_widths_in"]) - 8.0) < 0.01

    def test_delete_everything_refuses(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        sel = {"shape_id": r["shape_id"]}
        with pytest.raises(PptMcpError):
            tb.delete_table_rows(pkg, slide, sel, 0, 2)
        with pytest.raises(PptMcpError):
            tb.delete_table_cols(pkg, slide, sel, 0, 2)

    def test_surgery_on_real_merged_table(self, mb):
        """The merge-aware paths against a REAL PowerPoint-authored table:
        insert a row at the very end (always a legal seam), delete it again,
        and verify the merge map is unchanged."""
        slide_index, shape_id, info = _merged_corpus_table(mb)
        sel = {"shape_id": shape_id}
        regions_before = info["merge_regions"]
        rows_before = info["rows"]
        out = tb.insert_table_rows(mb, slide_index, sel, rows_before)
        assert out["rows"] == rows_before + 1
        out = tb.delete_table_rows(mb, slide_index, sel, rows_before)
        assert out["rows"] == rows_before
        after = tb.get_table(mb, slide_index, sel)
        assert after["merge_regions"] == regions_before
        mb.save()


# --------------------------------------------------------- borders and format


class TestFormat:
    def test_border_schema_order_and_content(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        sel = {"shape_id": r["shape_id"]}
        tb.format_table_cells(
            pkg, slide, sel,
            range={"r1": 0, "c1": 0, "r2": 0, "c2": 0},
            borders={"all": {"width": 2, "color": "C00000", "dash": "dash"}},
            fill="EEF3FA",
        )
        tbl = tb.resolve_table(pkg, slide, sel)["tbl"]
        tc = tbl.findall(qn("a:tr"))[0].findall(qn("a:tc"))[0]
        tcpr = tc.find(qn("a:tcPr"))
        tags = [etree.QName(c).localname for c in tcpr]
        # Schema order: lnL, lnR, lnT, lnB, then the fill.
        assert tags == ["lnL", "lnR", "lnT", "lnB", "solidFill"]
        lnb = tcpr.find(qn("a:lnB"))
        assert lnb.get("w") == "25400"  # 2 pt
        assert lnb.find(f"{qn('a:solidFill')}/{qn('a:srgbClr')}").get("val") == "C00000"
        assert lnb.find(qn("a:prstDash")).get("val") == "dash"
        pkg.save()

    def test_inner_borders_write_both_edges(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        sel = {"shape_id": r["shape_id"]}
        tb.format_table_cells(
            pkg, slide, sel, borders={"inner_h": {"width": 1, "color": "888888"}}
        )
        tbl = tb.resolve_table(pkg, slide, sel)["tbl"]
        rows = tbl.findall(qn("a:tr"))
        upper = rows[0].findall(qn("a:tc"))[0].find(qn("a:tcPr"))
        lower = rows[1].findall(qn("a:tc"))[0].find(qn("a:tcPr"))
        assert upper.find(qn("a:lnB")) is not None
        assert lower.find(qn("a:lnT")) is not None
        assert upper.find(qn("a:lnT")) is None  # outer edges untouched

    def test_outer_borders_only_on_range_edge(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 3, 3, 1, 1, 6, 2)
        sel = {"shape_id": r["shape_id"]}
        tb.format_table_cells(
            pkg, slide, sel,
            range={"r1": 0, "c1": 0, "r2": 2, "c2": 2},
            borders={"outer": {"width": 1.5, "color": "000000"}},
        )
        tbl = tb.resolve_table(pkg, slide, sel)["tbl"]
        center = tbl.findall(qn("a:tr"))[1].findall(qn("a:tc"))[1]
        tcpr = center.find(qn("a:tcPr"))
        assert tcpr is None or all(
            tcpr.find(qn(t)) is None for t in ("a:lnL", "a:lnR", "a:lnT", "a:lnB")
        )

    def test_border_none_and_margins_anchor(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        sel = {"shape_id": r["shape_id"]}
        tb.format_table_cells(
            pkg, slide, sel,
            borders={"bottom": "none"},
            margins={"left": 0.2, "top": 0.1},
            anchor="middle",
        )
        tbl = tb.resolve_table(pkg, slide, sel)["tbl"]
        tcpr = tbl.findall(qn("a:tr"))[0].findall(qn("a:tc"))[0].find(qn("a:tcPr"))
        assert tcpr.find(f"{qn('a:lnB')}/{qn('a:noFill')}") is not None
        assert tcpr.get("marL") == str(round(0.2 * 914400))
        assert tcpr.get("marT") == str(round(0.1 * 914400))
        assert tcpr.get("anchor") == "ctr"

    def test_nothing_to_change_refuses(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        with pytest.raises(PptMcpError):
            tb.format_table_cells(pkg, slide, {"shape_id": r["shape_id"]})

    def test_cell_text_and_style_edits(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1, data=[["keep", "x"]])
        sel = {"shape_id": r["shape_id"]}
        tb.set_table_cells(
            pkg, slide, sel,
            [
                {"row": 0, "col": 0, "bold": True, "size": 14},  # restyle only
                {"row": 0, "col": 1, "text": "new\ntext", "color": "FF0000"},
                {"row": 1, "col": 0, "fill": "none", "anchor": "bottom"},
            ],
        )
        info = tb.get_table(pkg, slide, sel)
        texts = {(c["row"], c["col"]): c["text"] for c in info["cells"]}
        assert texts[(0, 0)] == "keep"  # text survived the restyle
        assert texts[(0, 1)] == "new\ntext"
        tbl = tb.resolve_table(pkg, slide, sel)["tbl"]
        rpr = tbl.findall(qn("a:tr"))[0].findall(qn("a:tc"))[0].find(
            f"{qn('a:txBody')}/{qn('a:p')}/{qn('a:r')}/{qn('a:rPr')}"
        )
        assert rpr.get("b") == "1" and rpr.get("sz") == "1400"
        pkg.save()


# ------------------------------------------------------- widths, heights, style


class TestSizesAndStyle:
    def test_set_column_widths_and_row_heights(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 3, 1, 1, 6, 1.5)
        sel = {"shape_id": r["shape_id"]}
        out = tb.set_column_widths(pkg, slide, sel, [1.0, 2.0, 3.0])
        assert out["widths_in"] == [1.0, 2.0, 3.0]
        out = tb.set_row_heights(pkg, slide, sel, {1: 1.25})
        assert out["heights_in"][1] == 1.25
        info = tb.get_table(pkg, slide, sel)
        assert abs(info["geometry"]["cx_in"] - 6.0) < 0.01
        with pytest.raises(PptMcpError):
            tb.set_column_widths(pkg, slide, sel, [1.0])
        with pytest.raises(TargetNotFound):
            tb.set_row_heights(pkg, slide, sel, {7: 1.0})

    def test_apply_style_named_guid_and_flags(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        sel = {"shape_id": r["shape_id"]}
        out = tb.apply_table_style(pkg, slide, sel, "dark2_accent3")
        assert out["style"] == "dark2_accent3"
        info = tb.get_table(pkg, slide, sel)
        assert info["style_guid"] == tb.TABLE_STYLES["dark2_accent3"]
        # Un-braced GUID passthrough is normalized (issue #645: braces are
        # part of the value).
        raw = "5940675a-b579-460e-94d1-54222c63f5da"
        out = tb.apply_table_style(pkg, slide, sel, raw)
        assert out["style"] == "no_style_table_grid"
        # Unknown GUID passes through braced (graceful fallback in PowerPoint).
        foreign = "{11111111-2222-3333-4444-555555555555}"
        assert tb.apply_table_style(pkg, slide, sel, foreign)["style"] == foreign
        # Flags only.
        out = tb.apply_table_style(pkg, slide, sel, first_col=True, band_rows=False)
        assert out["flags"]["firstCol"] is True
        assert out["flags"]["bandRow"] is False
        # Style removal.
        assert tb.apply_table_style(pkg, slide, sel, "none")["style"] == "none"
        assert tb.get_table(pkg, slide, sel)["style_guid"] is None
        with pytest.raises(PptMcpError):
            tb.apply_table_style(pkg, slide, sel, "no_such_style")
        with pytest.raises(PptMcpError):
            tb.apply_table_style(pkg, slide, sel)
        pkg.save()

    def test_style_table_has_74_entries(self):
        assert len(tb.TABLE_STYLES) == 74
        assert all(v.startswith("{") and v.endswith("}") for v in tb.TABLE_STYLES.values())


# ------------------------------------------------------------ export / import


class TestExportImport:
    def test_csv_roundtrip_real_table(self, mb, make_deck, tmp_path):
        """Export a military_brief table to CSV, import it into a fresh
        deck, compare every cell text."""
        tables = _corpus_tables(mb)
        if not tables:
            pytest.skip("corpus deck has no tables")
        slide_index, shape_id, rows, cols = tables[0]
        csv_path = tmp_path / "table.csv"
        out = tb.export_table(mb, slide_index, {"shape_id": shape_id}, str(csv_path))
        assert out["path"] == str(csv_path)
        with open(csv_path, encoding="utf-8", newline="") as fh:
            exported = [list(row) for row in csv.reader(fh)]
        assert len(exported) == rows and len(exported[0]) == cols

        dest = PptxPackage(make_deck("import.pptx", extra_slides=0))
        r = tb.import_table(dest, 0, str(csv_path))
        assert (r["rows"], r["cols"]) == (rows, cols)
        reimported = tb.export_table(dest, 0, {"shape_id": r["shape_id"]})["data"]
        assert reimported == exported
        dest.save()

    def test_json_roundtrip_preserves_merges(self, deck, tmp_path):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 3, 3, 1, 1, 6, 2, data=[["a", "b", "c"]])
        sel = {"shape_id": r["shape_id"]}
        tb.merge_cells(pkg, slide, sel, 0, 0, 0, 2)
        json_path = tmp_path / "table.json"
        tb.export_table(pkg, slide, sel, str(json_path))
        payload = json.loads(json_path.read_text(encoding="utf-8"))
        assert payload["merge_regions"] == [{"r1": 0, "c1": 0, "r2": 0, "c2": 2}]

        r2 = tb.import_table(pkg, slide, str(json_path), x=1, y=4)
        assert r2["merge_regions_applied"] == 1
        info = tb.get_table(pkg, slide, {"shape_id": r2["shape_id"]})
        assert info["merge_regions"] == [{"r1": 0, "c1": 0, "r2": 0, "c2": 2}]
        pkg.save()

    def test_import_into_existing_table_dims_must_match(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 2, 2, 1, 1, 4, 1)
        sel = {"shape_id": r["shape_id"]}
        with pytest.raises(UnsupportedStructure) as exc:
            tb.import_table(pkg, slide, [["a", "b", "c"]], table=sel)
        assert "1 x 3" in str(exc.value)
        out = tb.import_table(pkg, slide, [["a", "b"], ["c", "d"]], table=sel)
        assert out["imported"] is True
        texts = tb.export_table(pkg, slide, sel)["data"]
        assert texts == [["a", "b"], ["c", "d"]]

    def test_inline_export_no_path(self, deck):
        pkg, slide = deck
        r = tb.create_table(pkg, slide, 1, 2, 1, 1, 4, 0.5, data=[["p", "q"]])
        out = tb.export_table(pkg, slide, {"shape_id": r["shape_id"]})
        assert out["data"] == [["p", "q"]]
        assert "path" not in out
