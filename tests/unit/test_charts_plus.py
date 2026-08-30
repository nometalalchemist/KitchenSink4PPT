"""Charts wave 2: combo charts (per-series type, secondary axis),
format_chart (title/legend/axis titles/number format/gridlines/data labels),
and update_chart_data over multiple plot groups.

Axis-wiring assertions mirror the PowerPoint 365 ground truth captured via
COM on 2026-08-30 (scratchpad gt_combo): each secondary group carries its
own axis pair, catAx id first in c:axId order; the secondary catAx is
delete=1, the secondary valAx sits right with crosses=max, and each axis's
crossAx points at its partner. python-pptx is the oracle where its reader
copes (multi-plot values); axis internals are asserted on the XML directly.
"""

from __future__ import annotations

import io
import zipfile

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import PptMcpError, UnsupportedStructure
from kitchensink4ppt.core.package import PptxPackage
from kitchensink4ppt.ops import charts as ch

_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_SML = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"


def _qc(name):
    return f"{{{_C}}}{name}"


@pytest.fixture()
def deck(make_deck):
    from kitchensink4ppt.ops import slides as sl

    path = make_deck("charts_plus.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    slide = sl.insert_slide(pkg, 0)["index"]
    return pkg, slide


CATS = ["Q1", "Q2", "Q3"]
COMBO_SERIES = [
    {"name": "Revenue", "values": [10, 12, 14], "type": "column"},
    {"name": "Cost", "values": [7, 8, 9], "type": "column"},
    {"name": "Margin", "values": [0.3, 0.33, 0.36], "type": "line",
     "axis": "secondary"},
]


def _plot_area(pkg, part):
    return pkg.root(part).find(f"{_qc('chart')}/{_qc('plotArea')}")


def _axis_map(plot_area):
    out = {}
    for ax in plot_area:
        tag = etree.QName(ax).localname
        if tag in ("catAx", "valAx"):
            out[ax.find(_qc("axId")).get("val")] = (tag, ax)
    return out


class TestComboCreate:
    def test_column_line_secondary_structure(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "combo", CATS, COMBO_SERIES, 1, 1, 8, 5,
            title="Combo",
        )
        pa = _plot_area(pkg, r["chart_part"])
        bar = pa.find(_qc("barChart"))
        line = pa.find(_qc("lineChart"))
        assert bar is not None and line is not None
        # Groups carry their own axis pair; catAx id first (ground truth).
        bar_ax = [a.get("val") for a in bar.findall(_qc("axId"))]
        line_ax = [a.get("val") for a in line.findall(_qc("axId"))]
        assert bar_ax == [ch._CAT_AX_ID, ch._VAL_AX_ID]
        assert line_ax == [ch._CAT2_AX_ID, ch._VAL2_AX_ID]
        axes = _axis_map(pa)
        assert set(axes) == {
            ch._CAT_AX_ID, ch._VAL_AX_ID, ch._CAT2_AX_ID, ch._VAL2_AX_ID
        }
        # Secondary pair wiring: hidden catAx2, right-side valAx2 at max.
        cat2_tag, cat2 = axes[ch._CAT2_AX_ID]
        val2_tag, val2 = axes[ch._VAL2_AX_ID]
        assert cat2_tag == "catAx" and val2_tag == "valAx"
        assert cat2.find(_qc("delete")).get("val") == "1"
        assert cat2.find(_qc("crossAx")).get("val") == ch._VAL2_AX_ID
        assert val2.find(_qc("axPos")).get("val") == "r"
        assert val2.find(_qc("crosses")).get("val") == "max"
        assert val2.find(_qc("crossAx")).get("val") == ch._CAT2_AX_ID
        # Primary axes untouched by the secondary pair.
        assert axes[ch._CAT_AX_ID][1].find(_qc("delete")).get("val") == "0"
        assert axes[ch._VAL_AX_ID][1].find(_qc("axPos")).get("val") == "l"

    def test_global_series_index_and_workbook(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "combo", CATS, COMBO_SERIES, 1, 1, 8, 5)
        pa = _plot_area(pkg, r["chart_part"])
        seen = {}
        for group in (pa.find(_qc("barChart")), pa.find(_qc("lineChart"))):
            for ser in group.findall(_qc("ser")):
                idx = ser.find(_qc("idx")).get("val")
                assert idx == ser.find(_qc("order")).get("val")
                name = ser.find(
                    f"{_qc('tx')}/{_qc('strRef')}/{_qc('strCache')}/"
                    f"{_qc('pt')}/{_qc('v')}"
                ).text
                seen[name] = idx
        # idx/order are GLOBAL across groups: 0,1 in bar, 2 in line.
        assert seen == {"Revenue": "0", "Cost": "1", "Margin": "2"}
        # The embedded workbook covers ALL series (col D = Margin).
        with zipfile.ZipFile(
            io.BytesIO(pkg.raw_part(r["embedded_workbook"]))
        ) as zf:
            sheet = etree.fromstring(zf.read("xl/worksheets/sheet1.xml"))
        cells = {c.get("r"): c for c in sheet.iter(f"{{{_SML}}}c")}
        assert cells["D1"].find(f"{{{_SML}}}is/{{{_SML}}}t").text == "Margin"
        assert cells["D3"].find(f"{{{_SML}}}v").text == "0.33"

    def test_single_type_secondary_axis_splits_groups(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "line", CATS,
            [
                {"name": "A", "values": [1, 2, 3]},
                {"name": "B", "values": [100, 200, 300], "axis": "secondary"},
            ],
            1, 1, 6, 4,
        )
        pa = _plot_area(pkg, r["chart_part"])
        lines = pa.findall(_qc("lineChart"))
        assert len(lines) == 2  # one group per axis
        assert [a.get("val") for a in lines[0].findall(_qc("axId"))] == [
            ch._CAT_AX_ID, ch._VAL_AX_ID
        ]
        assert [a.get("val") for a in lines[1].findall(_qc("axId"))] == [
            ch._CAT2_AX_ID, ch._VAL2_AX_ID
        ]

    def test_primary_only_chart_unchanged(self, deck):
        """No secondary series -> exactly the v1 two-axis emission."""
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "column", CATS,
            [{"name": "A", "values": [1, 2, 3]}], 1, 1, 6, 4,
        )
        pa = _plot_area(pkg, r["chart_part"])
        assert len(pa.findall(_qc("catAx"))) == 1
        assert len(pa.findall(_qc("valAx"))) == 1

    def test_combo_refusals(self, deck):
        pkg, slide = deck
        parts_before = set(pkg.part_names())
        with pytest.raises(PptMcpError, match="secondary"):
            ch.create_chart(
                pkg, slide, "combo", CATS,
                [{"name": "A", "values": [1, 2, 3], "axis": "secondary"}],
                1, 1, 6, 4,
            )
        with pytest.raises(PptMcpError, match="combo"):
            ch.create_chart(
                pkg, slide, "column", CATS,
                [{"name": "A", "values": [1, 2, 3], "type": "line"}],
                1, 1, 6, 4,
            )
        with pytest.raises(PptMcpError, match="not valid in a combo"):
            ch.create_chart(
                pkg, slide, "combo", CATS,
                [{"name": "A", "values": [1, 2, 3], "type": "pie"}],
                1, 1, 6, 4,
            )
        with pytest.raises(PptMcpError, match="chartex"):
            ch.create_chart(
                pkg, slide, "waterfall", CATS,
                [{"name": "A", "values": [1, 2, 3]}], 1, 1, 6, 4,
            )
        with pytest.raises(PptMcpError, match="axis"):
            ch.create_chart(
                pkg, slide, "pie", CATS,
                [{"name": "A", "values": [1, 2, 3], "axis": "secondary"}],
                1, 1, 6, 4,
            )
        assert set(pkg.part_names()) == parts_before

    def test_oracle_reads_combo_values(self, deck):
        pkg, slide = deck
        ch.create_chart(pkg, slide, "combo", CATS, COMBO_SERIES, 1, 1, 8, 5)
        path = pkg.save()
        from pptx import Presentation

        prs = Presentation(str(path))
        frames = [s for s in prs.slides[slide].shapes if s.has_chart]
        assert len(frames) == 1
        plots = frames[0].chart.plots
        assert len(plots) == 2
        values = [list(s.values) for p in plots for s in p.series]
        assert values == [
            [10.0, 12.0, 14.0], [7.0, 8.0, 9.0], [0.3, 0.33, 0.36]
        ]


class TestComboUpdate:
    def test_update_combo_roundtrip(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "combo", CATS, COMBO_SERIES, 1, 1, 8, 5)
        out = ch.update_chart_data(
            pkg, slide, {"shape_id": r["shape_id"]},
            ["J", "F", "M", "A"],
            [
                {"name": "Rev2", "values": [1, 2, 3, 4]},
                {"name": "Cost2", "values": [4, 3, 2, 1]},
                {"name": "Mgn2", "values": [0.1, 0.2, 0.3, 0.4]},
            ],
        )
        assert out["type"] == "combo"
        assert out["points"] == 4
        assert out["embedded_workbook"] == "regenerated"
        pa = _plot_area(pkg, r["chart_part"])
        line_ser = pa.find(_qc("lineChart")).find(_qc("ser"))
        vals = [
            pt.find(_qc("v")).text
            for pt in line_ser.find(
                f"{_qc('val')}/{_qc('numRef')}/{_qc('numCache')}"
            ).findall(_qc("pt"))
        ]
        assert vals == ["0.1", "0.2", "0.3", "0.4"]
        name = line_ser.find(
            f"{_qc('tx')}/{_qc('strRef')}/{_qc('strCache')}/{_qc('pt')}/{_qc('v')}"
        ).text
        assert name == "Mgn2"
        # The line series still reads from column D of the workbook.
        f = line_ser.find(f"{_qc('val')}/{_qc('numRef')}/{_qc('f')}").text
        assert f.startswith("Sheet1!$D$")
        pkg.save()

    def test_update_combo_series_count_refuses(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "combo", CATS, COMBO_SERIES, 1, 1, 8, 5)
        with pytest.raises(UnsupportedStructure, match="series"):
            ch.update_chart_data(
                pkg, slide, {"shape_id": r["shape_id"]}, CATS,
                [{"name": "only", "values": [1, 2, 3]}],
            )


class TestFormatChart:
    def test_title_legend_axis_number_gridlines_labels(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "column", CATS,
            [{"name": "A", "values": [1, 2, 3]}], 1, 1, 6, 4,
        )
        out = ch.format_chart(
            pkg, slide, r["shape_id"] if False else None,
            title="Quarterly",
            legend_pos="r",
            cat_axis_title="Quarter",
            val_axis_title="Units",
            number_format="#,##0.0",
            gridlines=True,
            data_labels=True,
        )
        assert set(out["changed"]) == {
            "title", "legend", "cat_axis_title", "val_axis_title",
            "number_format", "gridlines", "data_labels",
        }
        chart_el = pkg.root(r["chart_part"]).find(_qc("chart"))
        children = [etree.QName(c).localname for c in chart_el]
        # Schema order: title before plotArea, legend after it.
        assert children.index("title") < children.index("plotArea")
        assert children.index("plotArea") < children.index("legend")
        leg = chart_el.find(_qc("legend"))
        assert leg.find(_qc("legendPos")).get("val") == "r"
        pa = chart_el.find(_qc("plotArea"))
        cat_ax = pa.find(_qc("catAx"))
        val_ax = pa.find(_qc("valAx"))
        assert cat_ax.find(
            f"{_qc('title')}/{_qc('tx')}/{_qc('rich')}"
        ) is not None
        fmt = val_ax.find(_qc("numFmt"))
        assert fmt.get("formatCode") == "#,##0.0"
        assert fmt.get("sourceLinked") == "0"
        # numFmt sits at its schema position (after title, before crossAx).
        ax_children = [etree.QName(c).localname for c in val_ax]
        assert ax_children.index("numFmt") < ax_children.index("crossAx")
        assert ax_children.index("majorGridlines") < ax_children.index("crossAx")
        group = pa.find(_qc("barChart"))
        dlbls = group.find(_qc("dLbls"))
        assert dlbls is not None
        assert dlbls.find(_qc("showVal")).get("val") == "1"
        # dLbls follows the last c:ser.
        g_children = [etree.QName(c).localname for c in group]
        assert g_children.index("dLbls") == (
            len(g_children) - 1 - g_children[::-1].index("ser")
        ) + 1
        pkg.save()

    def test_title_and_labels_removal(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "column", CATS,
            [{"name": "A", "values": [1, 2, 3]}], 1, 1, 6, 4, title="T",
        )
        ch.format_chart(pkg, slide, None, data_labels=True, gridlines=True)
        out = ch.format_chart(
            pkg, slide, None, title="", data_labels=False, gridlines=False,
            legend=False,
        )
        assert "title_removed" in out["changed"]
        chart_el = pkg.root(r["chart_part"]).find(_qc("chart"))
        assert chart_el.find(_qc("title")) is None
        assert chart_el.find(_qc("autoTitleDeleted")).get("val") == "1"
        assert chart_el.find(_qc("legend")) is None
        pa = chart_el.find(_qc("plotArea"))
        assert pa.find(_qc("barChart")).find(_qc("dLbls")) is None
        assert pa.find(_qc("valAx")).find(_qc("majorGridlines")) is None

    def test_secondary_axis_title(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "combo", CATS, COMBO_SERIES, 1, 1, 8, 5)
        ch.format_chart(pkg, slide, None, secondary_val_axis_title="Margin %")
        pa = _plot_area(pkg, r["chart_part"])
        axes = _axis_map(pa)
        val2 = axes[ch._VAL2_AX_ID][1]
        t = val2.find(f"{_qc('title')}/{_qc('tx')}/{_qc('rich')}")
        assert t is not None
        # And the primary valAx did NOT get the title.
        assert axes[ch._VAL_AX_ID][1].find(_qc("title")) is None

    def test_format_refusals(self, deck):
        pkg, slide = deck
        ch.create_chart(
            pkg, slide, "column", CATS,
            [{"name": "A", "values": [1, 2, 3]}], 1, 1, 6, 4,
        )
        with pytest.raises(PptMcpError, match="nothing to change"):
            ch.format_chart(pkg, slide, None)
        with pytest.raises(PptMcpError, match="legend_pos"):
            ch.format_chart(pkg, slide, None, legend_pos="bottom")
        with pytest.raises(UnsupportedStructure, match="secondary"):
            ch.format_chart(pkg, slide, None, secondary_val_axis_title="X")

    def test_format_pie_axis_refusal(self, deck):
        pkg, slide = deck
        from kitchensink4ppt.ops import slides as sl

        s2 = sl.insert_slide(pkg, 0)["index"]
        ch.create_chart(
            pkg, s2, "pie", CATS, [{"name": "A", "values": [1, 2, 3]}],
            1, 1, 4, 3,
        )
        with pytest.raises(UnsupportedStructure, match="no value axis"):
            ch.format_chart(pkg, s2, None, gridlines=True)
        # Non-axis formatting still applies to pies.
        out = ch.format_chart(pkg, s2, None, title="Pie", data_labels=True)
        assert "data_labels" in out["changed"]

    def test_oracle_reads_formatted_title(self, deck):
        pkg, slide = deck
        ch.create_chart(
            pkg, slide, "column", CATS,
            [{"name": "A", "values": [1, 2, 3]}], 1, 1, 6, 4,
        )
        ch.format_chart(pkg, slide, None, title="Oracle Title", legend=True)
        path = pkg.save()
        from pptx import Presentation

        prs = Presentation(str(path))
        chart = next(
            s for s in prs.slides[slide].shapes if s.has_chart
        ).chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "Oracle Title"
        assert chart.has_legend


def test_com_validates_combo_and_formatted_deck(make_deck, tmp_path):
    """PowerPoint itself opens the combo + formatted deck clean."""
    import com_validate

    com_validate.com_gate()
    from kitchensink4ppt.ops import slides as sl

    deck = make_deck("charts_plus_com.pptx", extra_slides=0)
    pkg = PptxPackage(deck)
    slide = sl.insert_slide(pkg, 0)["index"]
    ch.create_chart(
        pkg, slide, "combo", CATS, COMBO_SERIES, 0.5, 0.5, 9, 5,
        title="Combo COM",
    )
    ch.format_chart(
        pkg, slide, None, legend_pos="b", cat_axis_title="Quarter",
        val_axis_title="USD", secondary_val_axis_title="Margin",
        number_format="#,##0", gridlines=True, data_labels=True,
    )
    s2 = sl.insert_slide(pkg, 0)["index"]
    ch.create_chart(
        pkg, s2, "line", CATS,
        [
            {"name": "A", "values": [1, 2, 3]},
            {"name": "B", "values": [100, 150, 90], "axis": "secondary"},
        ],
        0.5, 0.5, 9, 5,
    )
    pkg.save()

    out = com_validate.validate_files(tmp_path, [str(deck)])
    verdict = out["files"][str(deck)]
    assert verdict["opens_clean"] is True, verdict
    assert out["new_zombies"] == []  # PID-precise (com_validate)
