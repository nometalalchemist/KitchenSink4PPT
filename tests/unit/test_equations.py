"""LaTeX equations as native PowerPoint math (ops/equations.py).

Structure oracle: the a14 wrapper captured from PowerPoint 365's own
serialization (COM round-trip ground truth, 2026-08-31; the full COM proof
that PowerPoint's math engine ingests these decks lives in
tests/com_gates/equations_gate.py — run it with PowerPoint closed).

Checks here: wrapper shape (mc:AlternateContent / mc:Choice[@Requires="a14"]
/ a14:m / m:oMathPara / m:oMath), the Fallback run, conversion refusals that
leave the package untouched, invisibility to the plain-text index (by
design), the list_equations read path, and save/reload survival with
python-pptx as the opens-at-all oracle.
"""

from __future__ import annotations

import pytest
from lxml import etree

from kitchensink4ppt.core.package import PptxPackage, qn
from kitchensink4ppt.ops import equations as eq
from kitchensink4ppt.ops import read as rd
from kitchensink4ppt.ops.equations import (
    A14_NS,
    M_NS,
    MC_NS,
    EquationConversionError,
)
from kitchensink4ppt.core.errors import PptMcpError, UnsupportedStructure


def _math_paras(pkg, slide_part="ppt/slides/slide1.xml"):
    """Every a:p on the slide that carries an mc:AlternateContent math."""
    out = []
    for p in pkg.root(slide_part).iter(qn("a:p")):
        ac = p.find(f"{{{MC_NS}}}AlternateContent")
        if ac is not None:
            out.append((p, ac))
    return out


# ------------------------------------------------------------ wrapper shape


def test_simple_equation_wrapper_structure(make_deck):
    path = make_deck("eq_simple.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    result = eq.insert_equation(pkg, 0, r"E = mc^2", 1.0, 1.0)
    assert result["equation_added"] is True
    assert result["text"] == "E=mc2"

    paras = _math_paras(pkg)
    assert len(paras) == 1
    _p, ac = paras[0]
    choice = ac.find(f"{{{MC_NS}}}Choice")
    assert choice is not None
    # Requires="a14" with the a14 prefix declared in scope (MCE contract).
    assert choice.get("Requires") == "a14"
    assert choice.nsmap.get("a14") == A14_NS
    a14m = choice.find(f"{{{A14_NS}}}m")
    assert a14m is not None
    opara = a14m.find(f"{{{M_NS}}}oMathPara")
    assert opara is not None
    omath = opara.find(f"{{{M_NS}}}oMath")
    assert omath is not None
    # Content is pure math namespace (no Word w: leakage).
    for el in omath.iter():
        assert etree.QName(el).namespace in (M_NS,), etree.QName(el).text
    # Fallback carries the linearized run for pre-2010 consumers.
    fb = ac.find(f"{{{MC_NS}}}Fallback")
    assert fb is not None
    t = fb.find(f"{qn('a:r')}/{qn('a:t')}")
    assert t is not None and t.text == "E=mc2"

    # The box is a real text box shape with geometry.
    info = rd.get_slide_info(pkg, 0)
    box = next(s for s in info["shapes"] if s["id"] == result["shape_id"])
    assert box["type"] == "textbox"
    assert box["geometry"]["x"] == 914400


def test_fraction_and_greek_convert(make_deck):
    path = make_deck("eq_kinds.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    r_frac = eq.insert_equation(pkg, 0, r"\frac{a+b}{c}", 1.0, 1.0)
    r_greek = eq.insert_equation(
        pkg, 0, r"\alpha + \beta \geq \gamma", 1.0, 2.0
    )
    assert r_frac["text"] == "a+bc"  # linear approximation loses the bar
    assert r_greek["text"] == "α+β≥γ"
    # The fraction is structural OMML (m:f with num/den), not flat text.
    slide = pkg.root("ppt/slides/slide1.xml")
    fracs = slide.findall(f".//{{{M_NS}}}f")
    assert len(fracs) == 1
    assert fracs[0].find(f"{{{M_NS}}}num") is not None
    assert fracs[0].find(f"{{{M_NS}}}den") is not None


def test_equation_survives_save_and_reload(make_deck):
    path = make_deck("eq_roundtrip.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    eq.insert_equation(pkg, 0, r"\sum_{i=1}^{n} x_i", 1.0, 1.0)
    pkg.save()  # atomic save runs full payload validation

    pkg2 = PptxPackage(path)
    listed = eq.list_equations(pkg2)
    assert listed["equation_count"] == 1
    assert listed["equations"][0]["text"] == "i=1nxi"

    # Oracle: python-pptx opens the deck and sees the text box shape.
    from pptx import Presentation

    prs = Presentation(str(path))
    names = [s.name for s in prs.slides[0].shapes]
    assert any(n.startswith("Equation") for n in names)


# ---------------------------------------------------------------- refusals


def test_bad_latex_refuses_and_leaves_package_untouched(make_deck):
    path = make_deck("eq_bad.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    before = etree.tostring(pkg.root("ppt/slides/slide1.xml"))
    with pytest.raises(EquationConversionError) as exc:
        eq.insert_equation(pkg, 0, r"\frac{1}{", 1.0, 1.0)
    assert "not changed" in str(exc.value)
    assert r"\frac{1}{" in str(exc.value)
    assert etree.tostring(pkg.root("ppt/slides/slide1.xml")) == before
    assert not pkg._dirty  # nothing marked for re-serialization


def test_file_macro_latex_refused_by_name(make_deck):
    path = make_deck("eq_macro.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    with pytest.raises(PptMcpError, match="file/preamble macro"):
        eq.insert_equation(pkg, 0, r"\input{secrets.tex}", 1.0, 1.0)
    with pytest.raises(PptMcpError, match="non-empty"):
        eq.insert_equation(pkg, 0, "   ", 1.0, 1.0)


def test_equation_refuses_textless_shapes(make_deck):
    from kitchensink4ppt.ops import charts as ch

    path = make_deck("eq_target.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    r = ch.create_chart(
        pkg, 0, "pie", ["a", "b"], [{"name": "S", "values": [1, 2]}],
        1, 1, 4, 3,
    )
    with pytest.raises(UnsupportedStructure, match="chart"):
        eq.add_equation_to_shape(pkg, 0, r["shape_id"], r"x^2")


# ------------------------------------------------------- into existing shape


def test_add_equation_to_shape_appends_and_positions(make_deck):
    from kitchensink4ppt.ops import text as tx

    path = make_deck("eq_shape.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    box = tx.insert_textbox(pkg, 0, "before\nafter", 1.0, 1.0, 3.0, 1.5)
    sid = box["shape_id"]

    appended = eq.add_equation_to_shape(pkg, 0, sid, r"x^2")
    assert appended["paragraph"] == 2  # after the two text paragraphs

    inserted = eq.add_equation_to_shape(pkg, 0, sid, r"y^2", position=1)
    assert inserted["paragraph"] == 1  # between "before" and "after"

    listed = eq.list_equations(pkg, scope=0)
    assert [e["paragraph"] for e in listed["equations"]] == [1, 3]
    assert {e["text"] for e in listed["equations"]} == {"x2", "y2"}

    # Surrounding plain text is untouched.
    text = rd.get_text(pkg, scope=0)["slides"][0]["text"]
    assert "before" in text and "after" in text


# ------------------------------------------------------- read-side contract


def test_equations_invisible_to_text_index_by_design(make_deck):
    path = make_deck("eq_invisible.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    eq.insert_equation(pkg, 0, r"E = mc^2", 1.0, 1.0)

    # get_text and find_text never see math (m:t is outside the a:r index),
    # which also means search_and_replace can never corrupt an equation.
    assert "E=mc2" not in rd.get_text(pkg, scope=0)["slides"][0]["text"]
    assert rd.find_text(pkg, "mc2")["count"] == 0
    # list_equations is the read story.
    assert eq.list_equations(pkg)["equation_count"] == 1
