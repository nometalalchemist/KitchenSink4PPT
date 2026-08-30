"""Cross-deck slide copy (ops.interdeck.copy_slide_between).

Every test asserts the SOURCE deck is byte-identical (md5) after the copy,
and every mutated destination is saved, which runs pkg._validate_payload
(dangling rels, unresolvable sldIds, coordinate ceiling). A COM opens-clean
round (subprocess, tasklist gate, honest skip) validates the link-mode and
import-mode outputs against real PowerPoint at the end.
"""

from __future__ import annotations

import hashlib
import shutil
import zipfile
from pathlib import Path

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import (
    DocumentNotFound,
    PptMcpError,
    TargetNotFound,
)
from kitchensink4ppt.core.package import (
    PptxPackage,
    RT_SLIDE,
    RT_SLIDE_LAYOUT,
    qn,
    rels_name,
    resolve_target,
)
from kitchensink4ppt.ops.interdeck import (
    MASTER_ID_MIN,
    RT_SLIDE_MASTER,
    RT_THEME,
    copy_slide_between,
)
from kitchensink4ppt.ops.slides import (
    RT_CHART,
    RT_NOTES_SLIDE,
    create_presentation,
)

CORPUS = Path(__file__).parents[1] / "corpus"


def _md5(path: Path) -> str:
    return hashlib.md5(path.read_bytes()).hexdigest()


def _work_copy(name: str, tmp_path: Path) -> Path:
    src = CORPUS / name
    if not src.exists():
        pytest.skip(f"corpus file missing: {name}")
    work = tmp_path / name
    shutil.copy2(src, work)
    return work


def _copy_between(dest: Path, src: Path, slide, **kwargs) -> dict:
    """One cross-deck copy with the standing guarantees asserted: the source
    file's bytes never change, and the destination save passes payload
    validation."""
    before = _md5(src)
    pkg = PptxPackage(dest)
    res = copy_slide_between(pkg, str(src), slide, **kwargs)
    pkg.save(do_backup=False)  # runs _validate_payload on the output
    assert _md5(src) == before, "SOURCE deck was modified by the copy"
    return res


def _rel_targets(pkg: PptxPackage, part: str, rel_type: str) -> list[str]:
    name = rels_name(part)
    if not pkg.has_part(name):
        return []
    return [
        resolve_target(part, rel.get("Target", ""))
        for rel in pkg.root(name)
        if rel.get("Type") == rel_type and rel.get("TargetMode") != "External"
    ]


def _dest_layout_parts(pkg: PptxPackage) -> set[str]:
    from kitchensink4ppt.ops.slides import _layouts

    return {part for part, _name in _layouts(pkg)}


def _media_parts(pkg: PptxPackage) -> list[str]:
    return [n for n in pkg.part_names() if n.startswith("ppt/media/")]


def _fresh_dest(tmp_path: Path, name: str = "dest.pptx", template=None) -> Path:
    dest = tmp_path / name
    create_presentation(dest, template=template)
    return dest


# ----------------------------------------------------------------- link mode


def test_link_mode_proposal_slide_to_fresh_template_deck(tmp_path):
    src = _work_copy("proposal_defense.pptx", tmp_path)
    dest = _fresh_dest(
        tmp_path, template=CORPUS / "conference_template.potx"
    )
    res = _copy_between(dest, src, 1)  # the bullets slide

    assert res["design"] == "link"
    assert res["layout_match"] in ("name", "signature", "fallback")
    if res["layout_match"] == "fallback":
        # No match must come with an explicit warning listing what changed.
        assert any("first layout" in w for w in res["warnings"])
    assert res["index"] == 0 and res["creation_ids_regenerated"] >= 0

    pkg = PptxPackage(dest)
    parts = pkg.slide_parts()
    assert len(parts) == 1 and parts[0] == res["part"]
    # The layout rel points at a DESTINATION layout, never a dragged-in one.
    layout_targets = _rel_targets(pkg, parts[0], RT_SLIDE_LAYOUT)
    assert len(layout_targets) == 1
    assert layout_targets[0] in _dest_layout_parts(pkg)

    # Visual formatting is EXPLICIT on the copied slide: every shape either
    # brought its own geometry or had the source layout's carried inline.
    root = pkg.root(parts[0])
    ph_path = f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}"
    shapes = root.find(f"{qn('p:cSld')}/{qn('p:spTree')}").findall(qn("p:sp"))
    assert shapes, "copied slide lost its shapes"
    with_xfrm = [
        sp
        for sp in shapes
        if sp.find(f"{qn('p:spPr')}/{qn('a:xfrm')}") is not None
    ]
    assert with_xfrm, "no inline spPr geometry on the copied slide"
    # Placeholders that inherited geometry from the source layout must have
    # had it carried inline (real corpus decks may have no placeholders).
    spkg = PptxPackage(src)
    s_root = spkg.root(spkg.slide_parts()[1])
    inherited_ph = [
        sp
        for sp in s_root.iter(qn("p:sp"))
        if sp.find(ph_path) is not None
        and sp.find(f"{qn('p:spPr')}/{qn('a:xfrm')}") is None
    ]
    if inherited_ph:
        assert res["carried"].get("xfrm_carried", 0) >= 1

    from pptx import Presentation

    assert len(Presentation(str(dest)).slides) == 1


def test_link_mode_bakes_theme_references(tmp_path):
    src = _work_copy("proposal_defense.pptx", tmp_path)
    dest = _fresh_dest(tmp_path)
    res = _copy_between(dest, src, 0)  # title slide
    pkg = PptxPackage(dest)
    root = pkg.root(pkg.slide_parts()[0])
    # After baking against the SOURCE theme nothing theme-dependent remains
    # for the destination theme to re-resolve.
    assert root.find(f".//{qn('a:schemeClr')}") is None
    for el in root.iter(qn("a:latin"), qn("a:ea"), qn("a:cs")):
        assert not (el.get("typeface") or "").startswith("+"), (
            "theme font token survived baking"
        )
    assert res["design"] == "link" and "carried" in res


def _first_table_slide(pkg: PptxPackage) -> int | None:
    for i, part in enumerate(pkg.slide_parts()):
        if pkg.root(part).find(f".//{qn('a:tbl')}") is not None:
            return i
    return None


def test_link_mode_table_slide_military_to_proposal_derived(tmp_path):
    # The real military_brief may keep its tables on different slides than
    # the synthetic stand-in; find one (falling back to the tables deck).
    src = _work_copy("military_brief.pptx", tmp_path)
    index = _first_table_slide(PptxPackage(src))
    if index is None:
        src = _work_copy("pmr_tables.pptx", tmp_path)
        index = _first_table_slide(PptxPackage(src))
    if index is None:
        pytest.skip("no corpus slide with a native table")
    n_tables_src = len(
        PptxPackage(src)
        .root(PptxPackage(src).slide_parts()[index])
        .findall(f".//{qn('a:tbl')}")
    )

    dest = _fresh_dest(
        tmp_path, "proposal_derived.pptx",
        template=CORPUS / "proposal_defense.pptx",
    )
    res = _copy_between(dest, src, index)
    pkg = PptxPackage(dest)
    root = pkg.root(pkg.slide_parts()[0])
    assert len(root.findall(f".//{qn('a:tbl')}")) == n_tables_src, (
        "native table lost in transit"
    )
    assert res["design"] == "link"

    from pptx import Presentation

    prs = Presentation(str(dest))
    assert any(sh.has_table for sh in prs.slides[0].shapes)


def test_notes_slide_travels_and_notes_master_is_imported(tmp_path, make_deck):
    src = make_deck("src_notes.pptx")
    dest = _fresh_dest(tmp_path)  # default template: NO notesMaster
    assert not any(
        "notesMaster" in n for n in PptxPackage(dest).part_names()
    )
    res = _copy_between(dest, src, 0)  # slide 0 carries speaker notes

    pkg = PptxPackage(dest)
    slide_part = pkg.slide_parts()[0]
    notes = _rel_targets(pkg, slide_part, RT_NOTES_SLIDE)
    assert len(notes) == 1
    # Back-rel points at the NEW slide; the notesMaster rel resolves.
    back = _rel_targets(pkg, notes[0], RT_SLIDE)
    assert back == [slide_part]
    assert any("notesMaster" in n for n in pkg.part_names())
    assert any("notesMaster" in p for p in res["imported_design_parts"])
    lst = pkg.presentation().find(qn("p:notesMasterIdLst"))
    assert lst is not None and len(lst.findall(qn("p:notesMasterId"))) == 1

    from pptx import Presentation

    prs = Presentation(str(dest))
    assert (
        "Synthetic speaker notes"
        in prs.slides[0].notes_slide.notes_text_frame.text
    )


def test_chart_family_travels_with_intact_workbook(tmp_path, make_deck):
    from kitchensink4ppt.ops.charts import create_chart

    src = make_deck("src_chart.pptx")
    pkg = PptxPackage(src)
    made = create_chart(
        pkg, 0, "column", ["A", "B", "C"],
        [{"name": "Series 1", "values": [1, 2, 3]}],
        1, 1, 5, 3,
    )
    pkg.save(do_backup=False)
    src_xlsx_bytes = PptxPackage(src).part_bytes(made["embedded_workbook"])

    dest = _fresh_dest(tmp_path)
    res = _copy_between(dest, src, 0)

    pkg = PptxPackage(dest)
    slide_part = pkg.slide_parts()[0]
    charts = _rel_targets(pkg, slide_part, RT_CHART)
    assert len(charts) == 1
    assert charts[0] in res["copied_parts"]
    # The chart's own rels resolve and the embedded xlsx traveled intact.
    embedded = [
        t for t in _rel_targets(pkg, charts[0], _RT_PACKAGE)
        if t.endswith(".xlsx")
    ]
    assert len(embedded) == 1
    assert pkg.part_bytes(embedded[0]) == src_xlsx_bytes, (
        "embedded workbook corrupted in transit"
    )

    from pptx import Presentation

    prs = Presentation(str(dest))
    assert any(sh.has_chart for sh in prs.slides[0].shapes)


_RT_PACKAGE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package"
)


# --------------------------------------------------------------- import mode


def test_import_mode_brings_design_family_without_collisions(tmp_path):
    src = _work_copy("proposal_defense.pptx", tmp_path)
    # Destination already HAS slideMaster1.xml / slideLayout1..N / theme1.xml:
    # the exact partname-collision trap of the #1036 recipe.
    dest = _fresh_dest(tmp_path, template=CORPUS / "nsu_pcsj.pptx")
    before = PptxPackage(dest)
    masters_before = len(
        before.presentation()
        .find(qn("p:sldMasterIdLst"))
        .findall(qn("p:sldMasterId"))
    )
    themes_before = sum(
        1 for n in before.part_names() if n.startswith("ppt/theme/")
    )

    res = _copy_between(dest, src, 1, design="import")
    assert res["design"] == "import"
    assert res["imported_design_parts"], "import mode reported no imports"
    assert any("slideMaster" in p for p in res["imported_design_parts"])
    assert any("slideLayout" in p for p in res["imported_design_parts"])
    assert any("theme" in p for p in res["imported_design_parts"])

    pkg = PptxPackage(dest)
    # No duplicate zip entry names (the #1036 corruption).
    with zipfile.ZipFile(dest) as zf:
        names = zf.namelist()
    assert len(names) == len(set(names)), "duplicate partnames in package"

    m_lst = pkg.presentation().find(qn("p:sldMasterIdLst"))
    entries = m_lst.findall(qn("p:sldMasterId"))
    assert len(entries) == masters_before + 1
    # Fresh ids in the >= 2147483648 space, unique across masters + layouts.
    all_ids = [int(e.get("id")) for e in entries]
    for e in entries:
        master_part = pkg.relationship_target(
            "ppt/presentation.xml", e.get(qn("r:id"))
        )
        lst = pkg.root(master_part).find(qn("p:sldLayoutIdLst"))
        all_ids.extend(
            int(lid.get("id")) for lid in lst.findall(qn("p:sldLayoutId"))
        )
    assert len(all_ids) == len(set(all_ids)), "master/layout id collision"
    assert all(i >= MASTER_ID_MIN for i in all_ids)
    assert (
        sum(1 for n in pkg.part_names() if n.startswith("ppt/theme/"))
        == themes_before + 1
    )

    # The copied slide binds to the IMPORTED layout, whose master rel points
    # at the imported master.
    slide_part = pkg.slide_parts()[-1]
    layout = _rel_targets(pkg, slide_part, RT_SLIDE_LAYOUT)[0]
    assert layout in res["imported_design_parts"]
    master = _rel_targets(pkg, layout, RT_SLIDE_MASTER)[0]
    assert master in res["imported_design_parts"]
    theme = _rel_targets(pkg, master, RT_THEME)[0]
    assert theme in res["imported_design_parts"]

    from pptx import Presentation

    Presentation(str(dest))  # oracle opens


# ------------------------------------------------------- media deduplication


def test_media_dedup_two_slides_sharing_one_image(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    import make_corpus

    png = make_corpus._png(tmp_path / "shared.png", rgb=(10, 200, 30))
    src = tmp_path / "src_media.pptx"
    prs = Presentation()
    for _ in range(2):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_picture(str(png), Inches(1), Inches(1))
    prs.save(str(src))

    dest = _fresh_dest(tmp_path)
    media_start = len(_media_parts(PptxPackage(dest)))

    res1 = _copy_between(dest, src, 0)
    assert len(res1["media_added"]) == 1 and not res1["media_reused"]
    res2 = _copy_between(dest, src, 1)
    assert not res2["media_added"], "shared image was duplicated"
    assert res2["media_reused"], "dedup did not report the reused part"

    pkg = PptxPackage(dest)
    assert len(_media_parts(pkg)) == media_start + 1
    # Both copied slides reference the ONE deduplicated media part.
    targets = set()
    for part in pkg.slide_parts():
        for rel in pkg.root(rels_name(part)):
            if rel.get("TargetMode") != "External" and "media" in rel.get(
                "Target", ""
            ):
                targets.add(resolve_target(part, rel.get("Target", "")))
    assert len(targets) == 1


# ------------------------------------------------- warnings and refusals


def test_slide_size_mismatch_warns_with_both_dimensions(tmp_path):
    from pptx import Presentation
    from pptx.util import Emu

    src = tmp_path / "wide.pptx"
    prs = Presentation()  # default template is 4:3 (9144000 x 6858000)
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    prs.slides.add_slide(prs.slide_layouts[1])
    prs.save(str(src))

    dest = _fresh_dest(tmp_path)  # 4:3
    res = _copy_between(dest, src, 0)
    hits = [w for w in res["warnings"] if "slide size mismatch" in w]
    assert len(hits) == 1
    assert "12192000" in hits[0] and "9144000" in hits[0]
    assert "without rescaling" in hits[0]


def test_refuses_same_file_with_duplicate_slide_hint(tmp_path, make_deck):
    deck = make_deck("same.pptx")
    pkg = PptxPackage(deck)
    with pytest.raises(PptMcpError, match="duplicate_slide"):
        copy_slide_between(pkg, str(deck), 0)


def test_refuses_missing_source_and_missing_slide(tmp_path, make_deck):
    src = make_deck("src_ref.pptx")
    dest = _fresh_dest(tmp_path)
    pkg = PptxPackage(dest)
    with pytest.raises(DocumentNotFound):
        copy_slide_between(pkg, str(tmp_path / "nope.pptx"), 0)
    with pytest.raises(TargetNotFound, match="out of range"):
        copy_slide_between(pkg, str(src), 99)
    with pytest.raises(TargetNotFound, match="slide_id"):
        copy_slide_between(pkg, str(src), {"slide_id": 424242})
    with pytest.raises(PptMcpError, match="design"):
        copy_slide_between(pkg, str(src), 0, design="clone")


def test_cross_slide_jump_hyperlink_neutered_with_warning(tmp_path, make_deck):
    import posixpath

    src = make_deck("src_jump.pptx")
    pkg = PptxPackage(src)
    parts = pkg.slide_parts()
    target = posixpath.relpath(parts[2], posixpath.dirname(parts[0]))
    rid = pkg.add_relationship(parts[0], RT_SLIDE, target)
    sp = pkg.root(parts[0]).find(
        f"{qn('p:cSld')}/{qn('p:spTree')}/{qn('p:sp')}"
    )
    cnv = sp.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}")
    hl = etree.SubElement(cnv, qn("a:hlinkClick"))
    hl.set(qn("r:id"), rid)
    hl.set("action", "ppaction://hlinksldjump")
    pkg.mark_dirty(parts[0])
    pkg.save(do_backup=False)

    dest = _fresh_dest(tmp_path)
    res = _copy_between(dest, src, 0)
    assert res["neutered_hyperlinks"]
    assert res["neutered_hyperlinks"][0]["hyperlinks_removed"] == 1
    assert any("jump hyperlink" in w for w in res["warnings"])
    pkg = PptxPackage(dest)
    slide_part = pkg.slide_parts()[0]
    assert not _rel_targets(pkg, slide_part, RT_SLIDE)
    assert not list(pkg.root(slide_part).iter(qn("a:hlinkClick")))


def test_position_insert_and_section_membership(tmp_path, make_deck):
    src = make_deck("src_pos.pptx")
    dest = tmp_path / "dest_pos.pptx"
    create_presentation(
        dest, template=CORPUS / "nsu_pcsj.pptx", keep_slides=True
    )
    n = len(PptxPackage(dest).slide_parts())
    res = _copy_between(dest, src, 1, position=0)
    assert res["index"] == 0
    pkg = PptxPackage(dest)
    assert len(pkg.slide_parts()) == n + 1
    assert pkg.slide_parts()[0] == res["part"]

    from pptx import Presentation

    Presentation(str(dest))


# -------------------------------------------------------------- COM round


def test_com_validates_link_and_import_outputs(tmp_path):
    """One PowerPoint opens-clean round over both design modes (subprocess,
    tasklist gate, honest skip when the user's PowerPoint is open)."""
    import sys

    sys.path.insert(0, str(Path(__file__).parents[1]))
    import com_validate

    com_validate.com_gate()

    src = _work_copy("proposal_defense.pptx", tmp_path)
    linked = _fresh_dest(
        tmp_path, "com_link.pptx", template=CORPUS / "conference_template.potx"
    )
    _copy_between(linked, src, 0)
    _copy_between(linked, src, 1)
    imported = _fresh_dest(
        tmp_path, "com_import.pptx", template=CORPUS / "military_brief.pptx"
    )
    _copy_between(imported, src, 1, design="import")

    out = com_validate.validate_files(tmp_path, [str(linked), str(imported)])
    for path, verdict in out["files"].items():
        assert verdict["opens_clean"] is True, (
            f"PowerPoint did not open {path} clean: {verdict}"
        )
    assert out["post_powerpnt"] == 0
    assert out["new_zombies"] == []  # PID-precise (com_validate)
