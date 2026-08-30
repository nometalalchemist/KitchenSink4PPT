"""Phase 6 charts: the ported word-mcp chart builder against pptx packaging.

Checks the two render-critical invariants from the research: literal caches
(what renders) and the embedded workbook (what Edit Data opens) are ALWAYS
written together and always agree, and every r:id in the chart part resolves
(the dangling-rel repair trigger). python-pptx is the independent oracle for
reading created charts back.
"""

from __future__ import annotations

import io
import zipfile

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import (
    AmbiguousTarget,
    PptMcpError,
    TargetNotFound,
    UnsupportedStructure,
)
from kitchensink4ppt.core.package import PptxPackage, qn
from kitchensink4ppt.ops import charts as ch

_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_SML = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"


def _qc(name):
    return f"{{{_C}}}{name}"


@pytest.fixture()
def deck(make_deck):
    from kitchensink4ppt.ops import slides as sl

    path = make_deck("charts.pptx", extra_slides=0)
    pkg = PptxPackage(path)
    slide = sl.insert_slide(pkg, 0)["index"]
    return pkg, slide


CATS = ["Alpha", "Beta", "Gamma"]
SERIES = [
    {"name": "One", "values": [4, 2.5, 3]},
    {"name": "Two", "values": [1, 2, 3]},
]


def _num_cache_values(ser):
    cache = ser.find(f"{_qc('val')}/{_qc('numRef')}/{_qc('numCache')}")
    return [pt.find(_qc("v")).text for pt in cache.findall(_qc("pt"))]


def _workbook_sheet(pkg, workbook_part):
    with zipfile.ZipFile(io.BytesIO(pkg.raw_part(workbook_part))) as zf:
        names = set(zf.namelist())
        assert {
            "[Content_Types].xml",
            "_rels/.rels",
            "xl/workbook.xml",
            "xl/_rels/workbook.xml.rels",
            "xl/worksheets/sheet1.xml",
        } <= names
        return etree.fromstring(zf.read("xl/worksheets/sheet1.xml"))


class TestCreateChart:
    def test_column_chart_parts_and_caches(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "column", CATS, SERIES, 1, 1, 6, 4, title="T"
        )
        assert pkg.has_part(r["chart_part"])
        assert pkg.has_part(r["embedded_workbook"])
        root = pkg.root(r["chart_part"])
        group = root.find(f"{_qc('chart')}/{_qc('plotArea')}/{_qc('barChart')}")
        assert group.find(_qc("barDir")).get("val") == "col"
        assert group.find(_qc("grouping")).get("val") == "clustered"
        sers = group.findall(_qc("ser"))
        assert len(sers) == 2
        assert _num_cache_values(sers[0]) == ["4", "2.5", "3"]
        cat_cache = sers[0].find(
            f"{_qc('cat')}/{_qc('strRef')}/{_qc('strCache')}"
        )
        assert [p.find(_qc("v")).text for p in cat_cache.findall(_qc("pt"))] == CATS
        # Deliberate spPr absence: series follow the theme accent cycle.
        assert all(s.find(_qc("spPr")) is None for s in sers)
        # externalData is wired and its target exists.
        ext = root.find(_qc("externalData"))
        assert ext is not None
        pkg.save()

    def test_stacked_gets_overlap_100(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "bar_stacked", CATS, [SERIES[0]], 1, 1, 5, 3)
        group = pkg.root(r["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('barChart')}"
        )
        assert group.find(_qc("grouping")).get("val") == "stacked"
        assert group.find(_qc("overlap")).get("val") == "100"
        assert group.find(_qc("barDir")).get("val") == "bar"

    def test_pie_and_line_shapes(self, deck):
        pkg, slide = deck
        r1 = ch.create_chart(pkg, slide, "pie", CATS, [SERIES[0]], 0.5, 0.5, 4, 3)
        r2 = ch.create_chart(pkg, slide, "line", CATS, SERIES, 5, 0.5, 4, 3)
        pie = pkg.root(r1["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('pieChart')}"
        )
        assert pie is not None and pie.find(_qc("varyColors")).get("val") == "1"
        line = pkg.root(r2["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('lineChart')}"
        )
        assert line is not None
        assert len(line.findall(_qc("ser"))) == 2
        pkg.save()

    def test_workbook_matches_caches(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "column", CATS, SERIES, 1, 1, 6, 4)
        sheet = _workbook_sheet(pkg, r["embedded_workbook"])
        cells = {
            c.get("r"): c for c in sheet.iter(f"{{{_SML}}}c")
        }
        # A1 blank; B1/C1 series names; A2.. categories; B2.. values.
        assert cells["B1"].find(f"{{{_SML}}}is/{{{_SML}}}t").text == "One"
        assert cells["A2"].find(f"{{{_SML}}}is/{{{_SML}}}t").text == "Alpha"
        assert cells["B3"].find(f"{{{_SML}}}v").text == "2.5"
        assert cells["C4"].find(f"{{{_SML}}}v").text == "3"

    def test_data_refusals_fire_before_mutation(self, deck):
        pkg, slide = deck
        parts_before = set(pkg.part_names())
        with pytest.raises(PptMcpError):
            ch.create_chart(pkg, slide, "column", CATS, [{"name": "x", "values": [1]}], 1, 1, 4, 3)
        with pytest.raises(PptMcpError):
            ch.create_chart(pkg, slide, "pie", CATS, SERIES, 1, 1, 4, 3)
        with pytest.raises(PptMcpError):
            ch.create_chart(pkg, slide, "donut3d", CATS, SERIES, 1, 1, 4, 3)
        with pytest.raises(PptMcpError):
            ch.create_chart(
                pkg, slide, "line", CATS,
                [{"name": "x", "values": [1, float("nan"), 2]}], 1, 1, 4, 3,
            )
        assert set(pkg.part_names()) == parts_before

    def test_pptx_oracle_reads_chart(self, deck):
        pkg, slide = deck
        ch.create_chart(pkg, slide, "column", CATS, SERIES, 1, 1, 6, 4, title="Oracle")
        path = pkg.save()
        from pptx import Presentation

        prs = Presentation(str(path))
        frames = [s for s in prs.slides[slide].shapes if s.has_chart]
        assert len(frames) == 1
        chart = frames[0].chart
        assert list(chart.plots[0].categories) == CATS
        values = [list(s.values) for s in chart.plots[0].series]
        assert values == [[4.0, 2.5, 3.0], [1.0, 2.0, 3.0]]
        assert chart.has_title and chart.chart_title.text_frame.text == "Oracle"


class TestUpdateChartData:
    def test_roundtrip_caches_and_workbook(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "column", CATS, SERIES, 1, 1, 6, 4)
        out = ch.update_chart_data(
            pkg, slide, {"shape_id": r["shape_id"]},
            ["X", "Y", "Z", "W"],
            [
                {"name": "One2", "values": [9, 8, 7, 6]},
                {"name": "Two2", "values": [1, 1, 1, 1]},
            ],
        )
        assert out["points"] == 4
        assert out["embedded_workbook"] == "regenerated"
        root = pkg.root(r["chart_part"])
        group = root.find(f"{_qc('chart')}/{_qc('plotArea')}/{_qc('barChart')}")
        sers = group.findall(_qc("ser"))
        assert _num_cache_values(sers[0]) == ["9", "8", "7", "6"]
        name_cache = sers[1].find(f"{_qc('tx')}/{_qc('strRef')}/{_qc('strCache')}")
        assert name_cache.find(f"{_qc('pt')}/{_qc('v')}").text == "Two2"
        # Workbook agrees with the caches after the rewrite.
        sheet = _workbook_sheet(pkg, r["embedded_workbook"])
        cells = {c.get("r"): c for c in sheet.iter(f"{{{_SML}}}c")}
        assert cells["A5"].find(f"{{{_SML}}}is/{{{_SML}}}t").text == "W"
        assert cells["B2"].find(f"{{{_SML}}}v").text == "9"
        pkg.save()

    def test_series_count_change_refuses(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "column", CATS, SERIES, 1, 1, 6, 4)
        with pytest.raises(UnsupportedStructure) as exc:
            ch.update_chart_data(
                pkg, slide, {"shape_id": r["shape_id"]}, CATS,
                [{"name": "only", "values": [1, 2, 3]}],
            )
        assert "series count" in str(exc.value)

    def test_type_preserved_on_update(self, deck):
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "bar_stacked", CATS, SERIES, 1, 1, 6, 4)
        ch.update_chart_data(
            pkg, slide, {"shape_id": r["shape_id"]}, CATS, SERIES
        )
        group = pkg.root(r["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('barChart')}"
        )
        assert group.find(_qc("grouping")).get("val") == "stacked"
        assert group.find(_qc("overlap")).get("val") == "100"


class TestChartAddressing:
    def test_no_chart_refuses(self, deck):
        pkg, slide = deck
        with pytest.raises(TargetNotFound):
            ch.update_chart_data(pkg, slide, None, CATS, SERIES)

    def test_multi_chart_ambiguity(self, deck):
        pkg, slide = deck
        a = ch.create_chart(pkg, slide, "pie", CATS, [SERIES[0]], 0.5, 0.5, 4, 3)
        b = ch.create_chart(pkg, slide, "line", CATS, SERIES, 5, 0.5, 4, 3)
        with pytest.raises(AmbiguousTarget) as exc:
            ch.update_chart_data(pkg, slide, None, CATS, SERIES)
        assert str(a["shape_id"]) in str(exc.value)
        # Index addressing picks document order.
        out = ch.update_chart_data(pkg, slide, 1, CATS, SERIES)
        assert out["shape_id"] == b["shape_id"]

    def test_bad_selector(self, deck):
        pkg, slide = deck
        ch.create_chart(pkg, slide, "pie", CATS, [SERIES[0]], 1, 1, 4, 3)
        with pytest.raises(TargetNotFound):
            ch.update_chart_data(pkg, slide, 42, CATS, [SERIES[0]])
        with pytest.raises(PptMcpError):
            ch.update_chart_data(pkg, slide, "first", CATS, [SERIES[0]])


class TestClosure:
    def test_chart_rels_resolve_after_save(self, deck):
        """The saved package passes the dangling-rel validation with charts
        in it, and the chart part's own rels point at real parts."""
        pkg, slide = deck
        r = ch.create_chart(pkg, slide, "column", CATS, SERIES, 1, 1, 6, 4)
        path = pkg.save()
        with zipfile.ZipFile(path) as zf:
            names = set(zf.namelist())
            assert r["chart_part"] in names
            assert r["embedded_workbook"] in names
            rels = etree.fromstring(
                zf.read("ppt/charts/_rels/" + r["chart_part"].rsplit("/", 1)[1] + ".rels")
            )
            targets = [rel.get("Target") for rel in rels]
            assert any("embeddings" in t for t in targets)
            ct = zf.read("[Content_Types].xml").decode()
            assert "/" + r["chart_part"] in ct
            assert 'Extension="xlsx"' in ct


_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

SCATTER_SERIES = [
    {"name": "Alpha", "x": [1, 2, 3], "y": [2.5, 1, 4]},
    {"name": "Beta", "x": [0.5, 1.5], "y": [3, 3.5]},  # shorter on purpose
]


class TestScatter:
    """Wave-final scatter charts: c:scatterChart with xVal/yVal numRefs,
    marker-only styling (a:ln noFill), smooth=0, two value axes, and an
    embedded workbook laying out an X/Y column pair per series."""

    def test_scatter_structure_and_caches(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "scatter", None, SCATTER_SERIES, 1, 1, 6, 4,
            title="XY",
        )
        assert r["type"] == "scatter"
        assert r["points"] == 5
        root = pkg.root(r["chart_part"])
        plot = root.find(f"{_qc('chart')}/{_qc('plotArea')}")
        group = plot.find(_qc("scatterChart"))
        assert group is not None
        assert group.find(_qc("scatterStyle")).get("val") == "lineMarker"
        sers = group.findall(_qc("ser"))
        assert len(sers) == 2
        for ser, spec in zip(sers, SCATTER_SERIES):
            # Marker-only: the connecting line is suppressed per-series.
            ln = ser.find(f"{_qc('spPr')}/{{{_A_NS}}}ln")
            assert ln is not None and ln.find(f"{{{_A_NS}}}noFill") is not None
            assert ser.find(
                f"{_qc('marker')}/{_qc('symbol')}"
            ).get("val") == "circle"
            xc = ser.find(f"{_qc('xVal')}/{_qc('numRef')}/{_qc('numCache')}")
            yc = ser.find(f"{_qc('yVal')}/{_qc('numRef')}/{_qc('numCache')}")
            assert int(xc.find(_qc("ptCount")).get("val")) == len(spec["x"])
            assert int(yc.find(_qc("ptCount")).get("val")) == len(spec["y"])
            assert [float(p.find(_qc("v")).text)
                    for p in yc.findall(_qc("pt"))] == [
                float(v) for v in spec["y"]
            ]
            assert ser.find(_qc("smooth")).get("val") == "0"
            assert ser.find(_qc("cat")) is None
            assert ser.find(_qc("val")) is None
        # Two VALUE axes, no category axis, cross-wired ids.
        assert plot.find(_qc("catAx")) is None
        vals = plot.findall(_qc("valAx"))
        assert len(vals) == 2
        ids = [ax.find(_qc("axId")).get("val") for ax in vals]
        crosses = [ax.find(_qc("crossAx")).get("val") for ax in vals]
        assert ids == crosses[::-1]
        # Workbook: X/Y column pair per series (A/B then C/D).
        sheet = _workbook_sheet(pkg, r["embedded_workbook"])
        sml = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
        refs = [c.get("r") for c in sheet.iter(f"{sml}c")]
        assert {"A1", "B1", "C1", "D1"} <= set(refs)
        assert "A4" in refs and "C4" not in refs  # series lengths differ
        pkg.save()

    def test_scatter_refusals_fire_before_mutation(self, deck):
        pkg, slide = deck
        cases = [
            (["c1"], SCATTER_SERIES, "no categories"),
            (None, [{"name": "A", "x": [1, 2], "y": [1]}], "x values but"),
            (None, [{"name": "A", "values": [1]}], "x/y value lists"),
            (None, [{"name": "A", "x": [], "y": []}], "no data points"),
            (None, [{"name": "A", "x": [1], "y": [1], "axis": "secondary"}],
             "primary-axis only"),
            (None, [{"name": "A", "x": [1], "y": [1], "type": "line"}],
             "do not join combo"),
        ]
        for cats, series, msg in cases:
            with pytest.raises(PptMcpError, match=msg):
                ch.create_chart(
                    pkg, slide, "scatter", cats, series, 1, 1, 4, 3
                )
        # no chart/workbook parts were created by any refused call
        assert not any(
            n.startswith(("ppt/charts/", "ppt/embeddings/"))
            for n in pkg.part_names()
        )

    def test_scatter_update_roundtrip(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "scatter", None, SCATTER_SERIES, 1, 1, 6, 4
        )
        out = ch.update_chart_data(
            pkg, slide, None, None,
            [{"name": "A2", "x": [9, 8, 7, 6], "y": [1, 2, 3, 4]},
             {"name": "B2", "x": [5], "y": [5]}],
        )
        assert out["type"] == "scatter"
        assert out["embedded_workbook"] == "regenerated"
        group = pkg.root(r["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('scatterChart')}"
        )
        sers = group.findall(_qc("ser"))
        xc = sers[0].find(f"{_qc('xVal')}/{_qc('numRef')}/{_qc('numCache')}")
        assert [p.find(_qc("v")).text for p in xc.findall(_qc("pt"))] == [
            "9", "8", "7", "6"
        ]
        tx0 = sers[0].find(f"{_qc('tx')}/{_qc('strRef')}/{_qc('strCache')}")
        assert tx0.find(_qc("pt")).find(_qc("v")).text == "A2"
        # Update refuses the category data model on a scatter chart.
        with pytest.raises(PptMcpError, match="no categories"):
            ch.update_chart_data(
                pkg, slide, None, ["c"],
                [{"name": "A", "values": [1]}, {"name": "B", "values": [2]}],
            )
        pkg.save()

    def test_scatter_pptx_oracle(self, deck):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE

        pkg, slide = deck
        ch.create_chart(
            pkg, slide, "scatter", None, SCATTER_SERIES, 1, 1, 6, 4
        )
        path = pkg.save()
        prs = Presentation(str(path))
        chart = next(
            s.chart for sl in prs.slides for s in sl.shapes if s.has_chart
        )
        assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER
        series = list(chart.series)
        assert len(series) == 2
        # python-pptx reads scatter series .values as the Y values.
        assert list(series[0].values) == [2.5, 1.0, 4.0]
        assert list(series[1].values) == [3.0, 3.5]


class TestSeriesColors:
    """Per-series color OPT-IN: solidFill lands exactly where asked (bar/pie
    fill, line stroke, scatter marker) and the deliberate-spPr-absence
    default is untouched for series that did not ask."""

    def test_bar_color_optin_and_default_absence(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "column", CATS,
            [{"name": "One", "values": [1, 2, 3], "color": "1F4E79"},
             {"name": "Two", "values": [4, 5, 6]}],
            1, 1, 6, 4,
        )
        sers = pkg.root(r["chart_part"]).findall(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('barChart')}/{_qc('ser')}"
        )
        fill = sers[0].find(
            f"{_qc('spPr')}/{{{_A_NS}}}solidFill/{{{_A_NS}}}srgbClr"
        )
        assert fill is not None and fill.get("val") == "1F4E79"
        assert sers[1].find(_qc("spPr")) is None  # default stays deliberate
        # spPr must sit between c:tx and c:cat (schema order).
        kids = [etree.QName(c).localname for c in sers[0]]
        assert kids.index("tx") < kids.index("spPr") < kids.index("cat")
        pkg.save()

    def test_line_color_is_stroke_and_theme_token_works(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "line", CATS,
            [{"name": "One", "values": [1, 2, 3], "color": "accent2"}],
            1, 1, 6, 4,
        )
        ser = pkg.root(r["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('lineChart')}/{_qc('ser')}"
        )
        clr = ser.find(
            f"{_qc('spPr')}/{{{_A_NS}}}ln/{{{_A_NS}}}solidFill"
            f"/{{{_A_NS}}}schemeClr"
        )
        assert clr is not None and clr.get("val") == "accent2"
        # Fill itself untouched: only the stroke was asked for.
        assert ser.find(f"{_qc('spPr')}/{{{_A_NS}}}solidFill") is None

    def test_scatter_color_lands_on_marker(self, deck):
        pkg, slide = deck
        r = ch.create_chart(
            pkg, slide, "scatter", None,
            [{"name": "A", "x": [1, 2], "y": [3, 4], "color": "#FF0000"}],
            1, 1, 6, 4,
        )
        ser = pkg.root(r["chart_part"]).find(
            f"{_qc('chart')}/{_qc('plotArea')}/{_qc('scatterChart')}"
            f"/{_qc('ser')}"
        )
        mfill = ser.find(
            f"{_qc('marker')}/{_qc('spPr')}/{{{_A_NS}}}solidFill"
            f"/{{{_A_NS}}}srgbClr"
        )
        assert mfill is not None and mfill.get("val") == "FF0000"
        # The connecting line stays suppressed regardless of color.
        assert ser.find(
            f"{_qc('spPr')}/{{{_A_NS}}}ln/{{{_A_NS}}}noFill"
        ) is not None

    def test_invalid_color_refuses_before_mutation(self, deck):
        pkg, slide = deck
        with pytest.raises(PptMcpError, match="color must be"):
            ch.create_chart(
                pkg, slide, "pie", CATS,
                [{"name": "S", "values": [1, 2, 3], "color": "red"}],
                1, 1, 4, 3,
            )
        # no chart/workbook parts were created by the refused call
        assert not any(
            n.startswith(("ppt/charts/", "ppt/embeddings/"))
            for n in pkg.part_names()
        )
