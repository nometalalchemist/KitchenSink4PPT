"""PptxPackage: byte-identical round-trips on the real corpus, pre-open
refusals, presentation-order slide resolution against the python-pptx oracle,
payload validation (dangling rels), the per-slide shape-id allocator, and the
atomic add_slide_part machinery.

Every corpus file is copied to tmp before ANY open; the corpus copies are
never opened for saving, and the user's original decks are never touched at
all (tests/corpus already holds byte copies).
"""

from __future__ import annotations

import io
import shutil
import sys
import zipfile
from pathlib import Path

import pytest
from lxml import etree

from kitchensink4ppt.core.errors import (
    DocumentCorrupt,
    DocumentLocked,
    DocumentNotFound,
    DocumentProtected,
    ValidationFailed,
)
from kitchensink4ppt.core.package import (
    PRESENTATION_PART,
    RT_SLIDE,
    RT_SLIDE_LAYOUT,
    PptxPackage,
    qn,
    resolve_target,
)

CORPUS = Path(__file__).parents[1] / "corpus"

CORPUS_NAMES = [
    "proposal_defense.pptx",
    "nsu_pcsj.pptx",
    "unitar_final.pptx",
    "conference_template.potx",
    "military_brief.pptx",
    "pmr_tables.pptx",
]
PPTX_NAMES = [n for n in CORPUS_NAMES if n.endswith(".pptx")]

# python-pptx's own new-slide XML shape: the minimal slide PowerPoint accepts.
MINIMAL_SLIDE_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\r\n'
    '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
    "<p:cSld><p:spTree><p:nvGrpSpPr>"
    '<p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/>'
    "</p:nvGrpSpPr><p:grpSpPr><a:xfrm>"
    '<a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
    '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/>'
    "</a:xfrm></p:grpSpPr></p:spTree></p:cSld>"
    "<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>"
).encode("utf-8")


def _work_copy(name: str, tmp_path: Path) -> Path:
    src = CORPUS / name
    if not src.exists():
        pytest.skip(f"corpus file missing: {name}")
    work = tmp_path / name
    shutil.copy2(src, work)
    return work


def _oracle_slide_parts(path: Path) -> list[str]:
    """python-pptx as an independent oracle for presentation-order slide
    parts. Deliberately avoids the Presentation.slides property: that
    property calls rename_slide_parts(), which cosmetically renumbers slide
    partnames into presentation order and would make the comparison
    vacuous. This resolves sldIdLst rIds through python-pptx's own rels
    machinery instead."""
    from pptx import Presentation

    part = Presentation(str(path)).part
    sld_id_lst = part._element.get_or_add_sldIdLst()
    return [str(part.related_part(s.rId).partname)[1:] for s in sld_id_lst]


def _layout_of(pkg: PptxPackage, slide_part: str) -> str:
    rels = pkg.rels_for(slide_part)
    for rel in rels.getroot():
        if rel.get("Type") == RT_SLIDE_LAYOUT:
            return resolve_target(slide_part, rel.get("Target"))
    raise AssertionError(f"{slide_part} has no layout relationship")


# ----------------------------------------------------------------- round-trip


@pytest.mark.parametrize("name", CORPUS_NAMES)
def test_roundtrip_byte_identical_when_nothing_dirty(name, tmp_path):
    """Open a real deck, read it (parsing parts lazily), save with nothing
    marked dirty: every part's bytes and the entry order must be identical."""
    work = _work_copy(name, tmp_path)
    pkg = PptxPackage(work)
    # Reading parses trees; it must NOT cause re-serialization on save.
    pkg.tree(PRESENTATION_PART)
    pkg.slide_parts()
    out = tmp_path / ("out_" + name)
    pkg.save(out, do_backup=False)
    with zipfile.ZipFile(work) as za, zipfile.ZipFile(out) as zb:
        assert za.namelist() == zb.namelist(), "entry set/order changed"
        for n in za.namelist():
            assert za.read(n) == zb.read(n), f"part bytes differ: {n}"


# ------------------------------------------------------------ pre-open checks


def test_missing_file_refused(tmp_path):
    with pytest.raises(DocumentNotFound):
        PptxPackage(tmp_path / "nope.pptx")


def test_ole_magic_refused_as_protected(tmp_path):
    f = tmp_path / "legacy.pptx"
    f.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 128)
    with pytest.raises(DocumentProtected):
        PptxPackage(f)


def test_bad_zip_refused_as_corrupt(tmp_path):
    f = tmp_path / "junk.pptx"
    f.write_bytes(b"this is not a zip archive at all, not even close")
    with pytest.raises(DocumentCorrupt):
        PptxPackage(f)


def test_zip_without_presentation_refused(tmp_path):
    f = tmp_path / "notppt.pptx"
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("hello.xml", "<hello/>")
    f.write_bytes(buf.getvalue())
    with pytest.raises(DocumentCorrupt):
        PptxPackage(f)


@pytest.mark.skipif(sys.platform != "win32", reason="share-mode locking is Windows-only")
def test_exclusively_locked_file_detected(make_deck):
    doc = make_deck()
    import win32con
    import win32file

    handle = win32file.CreateFile(
        str(doc),
        win32con.GENERIC_READ,
        0,  # no sharing: same effect as PowerPoint's exclusive open
        None,
        win32con.OPEN_EXISTING,
        0,
        None,
    )
    try:
        with pytest.raises(DocumentLocked):
            PptxPackage(doc)
    finally:
        handle.Close()


# ------------------------------------------------------- slide order (oracle)


@pytest.mark.parametrize("name", PPTX_NAMES)
def test_slide_parts_order_matches_python_pptx_oracle(name, tmp_path):
    work = _work_copy(name, tmp_path)
    ours = PptxPackage(work).slide_parts()
    assert ours == _oracle_slide_parts(work)
    assert ours, f"{name} reported zero slides"


def test_potx_slide_parts_resolve(tmp_path):
    work = _work_copy("conference_template.potx", tmp_path)
    pkg = PptxPackage(work)
    for part in pkg.slide_parts():
        assert pkg.has_part(part)


def test_slide_parts_follows_sldidlst_not_zip_order(make_deck):
    """Reorder sldIdLst without renaming any part: slide_parts must follow
    the new list order while the ZIP layout stays as it was."""
    doc = make_deck(extra_slides=2)
    pkg = PptxPackage(doc)
    before = pkg.slide_parts()
    assert len(before) >= 3
    zip_order_before = pkg.part_names()
    lst = pkg.presentation().find(qn("p:sldIdLst"))
    entries = lst.findall(qn("p:sldId"))
    lst.remove(entries[0])
    lst.append(entries[0])  # first slide moves to the end
    pkg.mark_dirty(PRESENTATION_PART)
    pkg.save(do_backup=False)

    pkg2 = PptxPackage(doc)
    assert pkg2.slide_parts() == before[1:] + [before[0]]
    assert pkg2.part_names() == zip_order_before
    assert pkg2.slide_parts() == _oracle_slide_parts(doc)


# --------------------------------------------------------- payload validation


def test_validate_payload_catches_dangling_rel(make_deck):
    doc = make_deck()
    original = doc.read_bytes()
    pkg = PptxPackage(doc)
    pkg.add_relationship(PRESENTATION_PART, RT_SLIDE, "slides/slide999.xml")
    with pytest.raises(ValidationFailed, match="dangling"):
        pkg.save(do_backup=False)
    assert doc.read_bytes() == original, "failed save touched the original"


def test_validate_payload_catches_unresolvable_sldid(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    lst = pkg.presentation().find(qn("p:sldIdLst"))
    sld = etree.SubElement(lst, qn("p:sldId"))
    sld.set("id", "99999")
    sld.set(qn("r:id"), "rId9999")
    pkg.mark_dirty(PRESENTATION_PART)
    with pytest.raises(ValidationFailed):
        pkg.save(do_backup=False)


def test_validate_payload_catches_malformed_xml(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    pkg.set_raw_part("ppt/broken.xml", b"<unclosed")
    pkg.add_content_type_override(
        "ppt/broken.xml", "application/xml"
    )
    with pytest.raises(ValidationFailed, match="well-formed"):
        pkg.save(do_backup=False)


# ------------------------------------------------------------------ shape ids


def test_next_shape_id_is_max_plus_one_per_slide(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    for part in pkg.slide_parts():
        ids = [
            int(n.get("id"))
            for n in pkg.root(part).iter(qn("p:cNvPr"))
        ]
        assert ids, f"{part} has no cNvPr elements"
        nid = pkg.next_shape_id(part)
        assert nid == max(ids) + 1
        assert nid not in ids


def test_shape_id_scope_is_per_slide_not_global(make_deck):
    """Two slides may carry overlapping cNvPr ids; the allocator must key on
    the single slide, never on package-wide maxima."""
    doc = make_deck()
    pkg = PptxPackage(doc)
    parts = pkg.slide_parts()
    per_slide = {
        part: max(
            int(n.get("id")) for n in pkg.root(part).iter(qn("p:cNvPr"))
        )
        for part in parts
    }
    global_max = max(per_slide.values())
    low = min(per_slide, key=per_slide.get)
    if per_slide[low] == global_max:
        pytest.skip("synthetic deck has uniform shape counts")
    assert pkg.next_shape_id(low) == per_slide[low] + 1 != global_max + 1


# -------------------------------------------------------------- add_slide_part


def test_add_slide_part_end_to_end(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    first = pkg.slide_parts()[0]
    layout = _layout_of(pkg, first)
    before = pkg.slide_parts()
    info = pkg.add_slide_part(MINIMAL_SLIDE_XML, layout_part=layout)
    assert info["slide_id"] >= 256
    pkg.save(do_backup=False)

    pkg2 = PptxPackage(doc)
    parts = pkg2.slide_parts()
    assert len(parts) == len(before) + 1
    assert parts[-1] == info["part"]
    # Content-type override landed.
    ct = pkg2.root("[Content_Types].xml")
    assert any(
        o.get("PartName") == "/" + info["part"]
        for o in ct.findall(qn("ct:Override"))
    )
    # The new slide's layout relationship resolves.
    assert _layout_of(pkg2, info["part"]) == layout
    # Independent oracle agrees the deck grew by one slide.
    from pptx import Presentation

    assert len(list(Presentation(str(doc)).slides)) == len(before) + 1


def test_add_slide_part_refuses_bad_inputs_without_mutating(make_deck):
    doc = make_deck()
    pkg = PptxPackage(doc)
    names_before = pkg.part_names()
    with pytest.raises(Exception):
        pkg.add_slide_part(b"<not-well-formed", layout_part="ppt/slideLayouts/slideLayout1.xml")
    with pytest.raises(Exception):
        pkg.add_slide_part(MINIMAL_SLIDE_XML, layout_part="ppt/slideLayouts/slideLayout999.xml")
    assert pkg.part_names() == names_before, "failed add left partial state"
    pkg.save(do_backup=False)  # still saves clean
    PptxPackage(doc)
