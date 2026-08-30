"""Chart read-back (ops/chartread.py): round trips against charts created
by OUR create_chart, against a FOREIGN specimen built by python-pptx (the
independent oracle; the corpus decks carry no charts), scatter x/y, the
update_chart_data round trip, and the input contract."""

from __future__ import annotations

from pathlib import Path

import pytest

from kitchensink4ppt.core.errors import PptMcpError, TargetNotFound
from kitchensink4ppt.core.package import PptxPackage
from kitchensink4ppt.ops import charts, chartread, slides

CORPUS = Path(__file__).resolve().parents[1] / "corpus"

CATS = ["Q1", "Q2", "Q3"]
SERIES = [
    {"name": "Revenue", "values": [10, 20, 30]},
    {"name": "Cost", "values": [4, 8, 12]},
]


@pytest.fixture()
def deck(tmp_path):
    path = tmp_path / "charts.pptx"
    slides.create_presentation(path)
    pkg = PptxPackage(path)
    slides.insert_slide(pkg, 6)
    return pkg


# ==================================================== our own charts


def test_roundtrip_column_chart(deck):
    made = charts.create_chart(
        deck, 0, "column", CATS, SERIES, 1, 1, 8, 5, title="Quarterly"
    )
    res = chartread.get_chart_data(deck, 0)
    assert res["supported"] is True
    assert res["shape_id"] == made["shape_id"]
    assert res["chart_part"] == made["chart_part"]
    assert res["title"] == "Quarterly"
    assert res["group_count"] == 1
    grp = res["groups"][0]
    assert grp["type"] == "bar" and grp["direction"] == "col"
    assert grp["grouping"] == "clustered"
    assert [s["name"] for s in grp["series"]] == ["Revenue", "Cost"]
    for s, spec in zip(grp["series"], SERIES):
        assert s["categories"] == CATS
        assert s["values"] == spec["values"]
        assert s["values_ref"]  # the workbook formula is reported
    assert res["series_count"] == 2
    assert res["embedded_workbook"] == made["embedded_workbook"]
    assert "caches" in res["caveat"] or "cache" in res["caveat"]


def test_roundtrip_pie_and_line(deck):
    charts.create_chart(
        deck, 0, "pie", ["A", "B"], [{"name": "Share", "values": [70, 30]}],
        0.5, 0.5, 4, 3,
    )
    res = chartread.get_chart_data(deck, 0)
    assert res["groups"][0]["type"] == "pie"
    assert res["groups"][0]["series"][0]["values"] == [70, 30]

    slides.insert_slide(deck, 6)
    charts.create_chart(deck, 1, "line", CATS, SERIES, 0.5, 0.5, 8, 4)
    res2 = chartread.get_chart_data(deck, 1)
    assert res2["groups"][0]["type"] == "line"
    assert res2["groups"][0]["series"][0]["values"] == [10, 20, 30]


def test_roundtrip_scatter_xy(deck):
    """Wave 7's scatter (probed generically: xVal/yVal caches)."""
    made = charts.create_chart(
        deck, 0, "scatter", None,
        [
            {"name": "Trial", "x": [1, 2, 3], "y": [2.5, 4, 8]},
            {"name": "Control", "x": [1, 2], "y": [1, 1.5]},
        ],
        1, 1, 7, 5,
    )
    res = chartread.get_chart_data(deck, 0, made["shape_id"])
    grp = res["groups"][0]
    assert grp["type"] == "scatter"
    s1, s2 = grp["series"]
    assert s1["x_values"] == [1, 2, 3] and s1["y_values"] == [2.5, 4, 8]
    assert s2["x_values"] == [1, 2] and s2["y_values"] == [1, 1.5]
    assert "categories" not in s1


def test_read_reflects_update_chart_data(deck):
    made = charts.create_chart(deck, 0, "column", CATS, SERIES, 1, 1, 8, 5)
    charts.update_chart_data(
        deck, 0, made["shape_id"],
        categories=CATS,
        series=[
            {"name": "Revenue", "values": [11, 21, 31]},
            {"name": "Cost", "values": [5, 9, 13]},
        ],
    )
    res = chartread.get_chart_data(deck, 0)
    assert res["groups"][0]["series"][0]["values"] == [11, 21, 31]
    assert res["groups"][0]["series"][1]["values"] == [5, 9, 13]


def test_axis_titles_reported(deck):
    made = charts.create_chart(deck, 0, "column", CATS, SERIES, 1, 1, 8, 5)
    charts.format_chart(
        deck, 0, made["shape_id"],
        cat_axis_title="Quarter", val_axis_title="USD",
    )
    res = chartread.get_chart_data(deck, 0)
    titles = {a["kind"]: a["title"] for a in res["axes"]}
    assert titles.get("catAx") == "Quarter"
    assert titles.get("valAx") == "USD"
    positions = {a["kind"]: a.get("position") for a in res["axes"]}
    assert positions["catAx"] in ("b", "l")


def test_combo_multi_group(deck):
    charts.create_chart(
        deck, 0, "combo", CATS,
        [
            {"name": "Bars", "values": [1, 2, 3], "type": "column"},
            {"name": "Line", "values": [10, 20, 15], "type": "line"},
        ],
        1, 1, 8, 5,
    )
    res = chartread.get_chart_data(deck, 0)
    assert res["group_count"] == 2
    types = sorted(g["type"] for g in res["groups"])
    assert types == ["bar", "line"]
    assert res["series_count"] == 2


# ==================================================== foreign specimens


def _foreign_deck(tmp_path):
    """A chart deck written by python-pptx: the independent oracle standing
    in for legacy decks (corpus military_brief/pmr carry no charts)."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["East", "West", "North"]
    data.add_series("2025", (19.2, 21.4, 16.7))
    data.add_series("2026", (22.3, 28.6, 15.2))
    s1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(8), Inches(5), data,
    )
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    xy = XyChartData()
    ser = xy.add_series("Model")
    ser.add_data_point(0.7, 2.7)
    ser.add_data_point(1.8, 3.2)
    s2.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(1), Inches(1), Inches(8), Inches(5), xy,
    )
    path = tmp_path / "foreign.pptx"
    prs.save(path)
    return path


def test_reads_foreign_pptx_chart(tmp_path):
    path = _foreign_deck(tmp_path)
    pkg = PptxPackage(path)
    res = chartread.get_chart_data(pkg, 0)
    assert res["supported"] is True
    grp = res["groups"][0]
    assert grp["type"] == "bar" and grp["direction"] == "col"
    assert [s["name"] for s in grp["series"]] == ["2025", "2026"]
    assert grp["series"][0]["categories"] == ["East", "West", "North"]
    assert grp["series"][0]["values"] == [19.2, 21.4, 16.7]
    assert grp["series"][1]["values"] == [22.3, 28.6, 15.2]


def test_reads_foreign_scatter_xy(tmp_path):
    path = _foreign_deck(tmp_path)
    pkg = PptxPackage(path)
    res = chartread.get_chart_data(pkg, 1)
    grp = res["groups"][0]
    assert grp["type"] == "scatter"
    assert grp["series"][0]["x_values"] == [0.7, 1.8]
    assert grp["series"][0]["y_values"] == [2.7, 3.2]


# ==================================================== the input contract


def test_no_chart_and_selection_errors(deck):
    with pytest.raises(TargetNotFound, match="no charts"):
        chartread.get_chart_data(deck, 0)
    a = charts.create_chart(deck, 0, "column", CATS, SERIES, 0.5, 0.5, 4, 3)
    b = charts.create_chart(deck, 0, "pie", ["A"], [{"name": "s", "values": [1]}], 5, 0.5, 4, 3)
    with pytest.raises(PptMcpError, match="pass `chart`"):
        chartread.get_chart_data(deck, 0)
    assert chartread.get_chart_data(deck, 0, a["shape_id"])["groups"][0]["type"] == "bar"
    assert chartread.get_chart_data(deck, 0, b["shape_id"])["groups"][0]["type"] == "pie"
    with pytest.raises(TargetNotFound, match="chart shape ids there"):
        chartread.get_chart_data(deck, 0, 9999)
    with pytest.raises(PptMcpError, match="shape id"):
        chartread.get_chart_data(deck, 0, "first")


def test_read_only_nothing_dirtied(deck):
    charts.create_chart(deck, 0, "column", CATS, SERIES, 1, 1, 8, 5)
    deck.save()
    before = set(deck._dirty)
    chartread.get_chart_data(deck, 0)
    assert set(deck._dirty) == before == set()


def test_corpus_decks_have_no_charts_probe():
    """Corpus probe pinned (2026-08-31): no corpus deck carries a chart
    part, which is WHY the foreign specimen above is python-pptx-built."""
    for name in ("military_brief.pptx", "pmr_tables.pptx"):
        pkg = PptxPackage(CORPUS / name)
        assert not [
            n for n in pkg.part_names() if n.startswith("ppt/charts/")
        ]
        with pytest.raises(TargetNotFound):
            chartread.get_chart_data(pkg, 0)
