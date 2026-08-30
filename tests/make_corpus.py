"""Generate a synthetic test corpus, structurally equivalent to the private
real-world corpus the suite is developed against.

Any file already present in tests/corpus/ is left alone, so a local real
corpus always wins; CI and contributors get generated stand-ins with the same
structural properties (title slide, bullet levels, a table, a picture,
speaker notes) but zero real content.

Run directly:  python -X utf8 tests/make_corpus.py
Or automatically via conftest when corpus files are missing.
"""

from __future__ import annotations

import random
import struct
import sys
import tempfile
import zlib
from pathlib import Path

CORPUS = Path(__file__).parent / "corpus"

_WORDS = (
    "the framework extends across cases while compliance rises and the "
    "institution adapts its mandate toward renewed cooperation under "
    "conditions of uncertainty where actors weigh obligations against "
    "expected consequences and prior commitments shape present choices"
).split()

#: Content type stamped on ppt/presentation.xml for a .potx stand-in.
_CT_TEMPLATE_MAIN = (
    "application/vnd.openxmlformats-officedocument.presentationml"
    ".template.main+xml"
)


def _sentence(rng: random.Random, n: int = 8) -> str:
    ws = [rng.choice(_WORDS) for _ in range(n)]
    return (" ".join(ws)).capitalize() + "."


def _png(path: Path, w: int = 60, h: int = 30, rgb=(120, 60, 40)) -> Path:
    """Minimal hand-rolled truecolor PNG; no imaging library needed."""

    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(
            ">I", zlib.crc32(c) & 0xFFFFFFFF
        )

    raw = b"".join(b"\x00" + bytes(rgb) * w for _ in range(h))
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )
    return path


def build_deck(
    path: Path,
    *,
    seed: int = 0,
    extra_slides: int = 2,
    potx: bool = False,
) -> Path:
    """One synthetic deck: title slide, bullet slide, table slide, picture
    slide, speaker notes on two slides, plus `extra_slides` filler slides.
    Structure only, zero real content."""
    from pptx import Presentation
    from pptx.util import Inches

    rng = random.Random(seed)
    prs = Presentation()

    s = prs.slides.add_slide(prs.slide_layouts[0])  # title
    s.shapes.title.text = "Synthetic Deck"
    s.placeholders[1].text = "Structural stand-in, zero real content"
    s.notes_slide.notes_text_frame.text = "Synthetic speaker notes, slide one."

    s = prs.slides.add_slide(prs.slide_layouts[1])  # title + content bullets
    s.shapes.title.text = "Bullets"
    tf = s.placeholders[1].text_frame
    tf.text = _sentence(rng, 5)
    for level in (0, 1, 1, 2):
        para = tf.add_paragraph()
        para.text = _sentence(rng, 4)
        para.level = level

    s = prs.slides.add_slide(prs.slide_layouts[5])  # title only + table
    s.shapes.title.text = "Table"
    rows, cols = 3, 3
    table = s.shapes.add_table(
        rows, cols, Inches(1), Inches(1.8), Inches(8), Inches(3)
    ).table
    for r in range(rows):
        for c in range(cols):
            table.cell(r, c).text = f"r{r}c{c}"

    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank + picture
    with tempfile.TemporaryDirectory() as td:
        png = _png(Path(td) / "syn.png", rgb=(60, 90, 140))
        s.shapes.add_picture(str(png), Inches(1), Inches(1))
    s.notes_slide.notes_text_frame.text = "Synthetic speaker notes, picture slide."

    for _ in range(extra_slides):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = _sentence(rng, 3)[:-1]
        s.placeholders[1].text_frame.text = _sentence(rng, 6)

    prs.save(str(path))

    if potx:
        # Restamp the presentation part's content type so the stand-in is a
        # real template package, not a renamed .pptx.
        sys.path.insert(0, str(Path(__file__).parents[1] / "src"))
        from kitchensink4ppt.core.package import PptxPackage

        pkg = PptxPackage(path)
        pkg.add_content_type_override("ppt/presentation.xml", _CT_TEMPLATE_MAIN)
        pkg.save(do_backup=False)
    return path


GENERATORS = {
    "proposal_defense.pptx": lambda p: build_deck(p, seed=26, extra_slides=22),
    "nsu_pcsj.pptx": lambda p: build_deck(p, seed=13, extra_slides=6),
    "unitar_final.pptx": lambda p: build_deck(p, seed=3, extra_slides=8),
    "conference_template.potx": lambda p: build_deck(p, seed=7, potx=True),
    "military_brief.pptx": lambda p: build_deck(p, seed=18, extra_slides=30),
    "pmr_tables.pptx": lambda p: build_deck(p, seed=23, extra_slides=4),
}


def generate_missing(verbose: bool = True) -> list[str]:
    CORPUS.mkdir(exist_ok=True)
    made = []
    for name, gen in GENERATORS.items():
        target = CORPUS / name
        if target.exists():
            continue
        if verbose:
            print(f"generating synthetic {name} ...")
        gen(target)
        made.append(name)
    return made


if __name__ == "__main__":
    made = generate_missing()
    print(f"generated {len(made)} synthetic corpus file(s): {made or 'none needed'}")
