"""Slide CRUD and create-from-template operations.

Contract (all ops modules): every function takes the open PptxPackage first,
mutates only the in-memory package, calls pkg.mark_dirty() on every part it
touches, and returns a summary dict describing what changed. Nothing here
writes to disk; the caller decides when to save. create_presentation is the
one sanctioned exception: it creates a NEW file (byte copy of the template,
then a PptxPackage edit and an atomic validated save) and never modifies the
template source.

Slide addressing: a 0-based presentation-order index (int) or
{"slide_id": N} (the durable p:sldId id; survives reordering, indices do
not). Errors come from the core.errors taxonomy with actionable messages.

Duplication follows the package-level algorithm (research doc Part II, the
python-pptx issue #132 ground truth): the cloned slide KEEPS its rIds so the
slide XML needs no rewriting; only deep-copied targets (notesSlide, chart
plus its colors/style/embedded xlsx, oleObject/package embeddings) get their
rel Targets retargeted; media and layout rels stay shared by design; external
rels keep TargetMode; a16/p14 creationId GUIDs are REGENERATED on every clone
(duplicate GUIDs across slides are a confirmed corruption source).

Deletion garbage-collects: parts reachable only from the deleted slide,
custom-show references, section membership, and jump hyperlinks on OTHER
slides targeting the deleted one (rel dropped, a:hlinkClick neutered,
flagged in the result).
"""

from __future__ import annotations

import copy
import posixpath
import re
import uuid
from pathlib import Path

from lxml import etree

from ..core.errors import (
    AmbiguousTarget,
    DocumentNotFound,
    PptMcpError,
    TargetNotFound,
    UnsupportedStructure,
)
from ..core.package import (
    CT_SLIDE,
    NSMAP,
    PRESENTATION_PART,
    PptxPackage,
    RT_SLIDE,
    RT_SLIDE_LAYOUT,
    qn,
    rels_name,
    rels_source,
    resolve_target,
)
from ..core.sandbox import check_path

_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
RT_NOTES_SLIDE = _RT + "notesSlide"
RT_NOTES_MASTER = _RT + "notesMaster"
RT_CHART = _RT + "chart"
RT_OLE_OBJECT = _RT + "oleObject"
RT_PACKAGE = _RT + "package"
RT_TAGS = _RT + "tags"
RT_CHART_COLORS = "http://schemas.microsoft.com/office/2011/relationships/chartColorStyle"
RT_CHART_STYLE = "http://schemas.microsoft.com/office/2011/relationships/chartStyle"

CT_PRESENTATION_MAIN = (
    "application/vnd.openxmlformats-officedocument.presentationml"
    ".presentation.main+xml"
)
CT_NOTES_SLIDE = (
    "application/vnd.openxmlformats-officedocument.presentationml"
    ".notesSlide+xml"
)

#: Deep-copy-on-duplicate reltypes without special companion handling.
_DEEP_COPY_SIMPLE = (RT_OLE_OBJECT, RT_PACKAGE, RT_TAGS)

# Modern threaded comments (mirrors ops/comments.py; kept local so the
# slide layer does not import the comments module).
RT_MODERN_COMMENTS = (
    "http://schemas.microsoft.com/office/2018/10/relationships/comments"
)
RT_LEGACY_COMMENTS = _RT + "comments"
#: Slide extLst extension that makes the modern comments rel render.
EXT_URI_COMMENT_REL = "{6950BFC3-D8DA-4A85-94F7-54DA5524770B}"

#: Latent placeholder types (date, footer, slide number) are never cloned to
#: a new slide; they render from the layout/master when enabled.
_LATENT_PH_TYPES = {"dt", "ftr", "sldNum"}

#: Placeholder types that carry a text frame on a fresh slide.
_PH_TEXT_TYPES = {"title", "ctrTitle", "subTitle", "body", "obj"}

_PH_BASENAMES = {
    "title": "Title",
    "ctrTitle": "Title",
    "subTitle": "Subtitle",
    "body": "Text Placeholder",
    "obj": "Content Placeholder",
    "chart": "Chart Placeholder",
    "tbl": "Table Placeholder",
    "pic": "Picture Placeholder",
    "media": "Media Placeholder",
    "clipArt": "ClipArt Placeholder",
    "dgm": "Diagram Placeholder",
    "hdr": "Header Placeholder",
    "sldImg": "Slide Image Placeholder",
}

#: p14:sectionLst lives in p:extLst under this extension uri.
_SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


# ------------------------------------------------------------- addressing


def _sld_id_lst(pkg: PptxPackage) -> etree._Element | None:
    return pkg.presentation().find(qn("p:sldIdLst"))


def _slide_entries(pkg: PptxPackage) -> list[etree._Element]:
    lst = _sld_id_lst(pkg)
    return [] if lst is None else lst.findall(qn("p:sldId"))


def _resolve_slide(pkg: PptxPackage, slide) -> tuple[int, str, etree._Element, int, str]:
    """Resolve a slide selector to (index, part name, p:sldId element,
    slide id, presentation rId)."""
    entries = _slide_entries(pkg)
    n = len(entries)
    if isinstance(slide, dict):
        if set(slide) != {"slide_id"}:
            raise PptMcpError(
                "slide selector must be a 0-based index (int) or "
                '{"slide_id": N}; got ' + repr(slide)
            )
        sid = slide["slide_id"]
        for index, entry in enumerate(entries):
            if entry.get("id") == str(sid):
                break
        else:
            present = ", ".join(e.get("id", "?") for e in entries)
            raise TargetNotFound(
                f"no slide with slide_id {sid}; ids present: {present or 'none'}"
            )
    elif isinstance(slide, int) and not isinstance(slide, bool):
        if not 0 <= slide < n:
            raise TargetNotFound(
                f"slide index {slide} out of range, presentation has {n} "
                f"slide{'s' if n != 1 else ''}"
            )
        index = slide
    else:
        raise PptMcpError(
            "slide selector must be a 0-based index (int) or "
            '{"slide_id": N}; got ' + repr(slide)
        )
    entry = entries[index]
    rid = entry.get(qn("r:id"))
    part = pkg.relationship_target(PRESENTATION_PART, rid)
    return index, part, entry, int(entry.get("id")), rid


# --------------------------------------------------------------- layouts


def _masters(pkg: PptxPackage) -> list[str]:
    lst = pkg.presentation().find(qn("p:sldMasterIdLst"))
    if lst is None:
        return []
    return [
        pkg.relationship_target(PRESENTATION_PART, m.get(qn("r:id")))
        for m in lst.findall(qn("p:sldMasterId"))
    ]


def _layouts(pkg: PptxPackage) -> list[tuple[str, str]]:
    """[(layout part, layout name)] in master order, then each master's
    p:sldLayoutIdLst order. Index into this list is the global layout index."""
    out: list[tuple[str, str]] = []
    for master in _masters(pkg):
        lst = pkg.root(master).find(qn("p:sldLayoutIdLst"))
        if lst is None:
            continue
        for lid in lst.findall(qn("p:sldLayoutId")):
            part = pkg.relationship_target(master, lid.get(qn("r:id")))
            csld = pkg.root(part).find(qn("p:cSld"))
            name = (csld.get("name") if csld is not None else None) or ""
            out.append((part, name))
    return out


def _resolve_layout(pkg: PptxPackage, layout) -> str:
    layouts = _layouts(pkg)
    if not layouts:
        raise UnsupportedStructure(
            "presentation has no slide layouts; cannot build a slide"
        )
    if isinstance(layout, int) and not isinstance(layout, bool):
        if not 0 <= layout < len(layouts):
            raise TargetNotFound(
                f"layout index {layout} out of range, presentation has "
                f"{len(layouts)} layouts across {len(_masters(pkg))} master(s)"
            )
        return layouts[layout][0]
    if isinstance(layout, str):
        exact = [(p, nm) for p, nm in layouts if nm == layout]
        if not exact:
            exact = [(p, nm) for p, nm in layouts if nm.lower() == layout.lower()]
        if len(exact) == 1:
            return exact[0][0]
        if len(exact) > 1:
            hits = ", ".join(f"{i}:{nm!r}" for i, (p, nm) in enumerate(layouts) if (p, nm) in exact)
            raise AmbiguousTarget(
                f"{len(exact)} layouts are named {layout!r} (multiple masters); "
                f"use a layout index instead: {hits}"
            )
        names = ", ".join(repr(nm) for _, nm in layouts)
        raise TargetNotFound(
            f"no layout named {layout!r}; available layouts: {names}"
        )
    raise PptMcpError(
        "layout must be a layout name (str) or 0-based global index (int)"
    )


# ------------------------------------------------------- XML construction


def _serialize(root: etree._Element) -> bytes:
    return etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )


def _rel_target(source_part: str, target_part: str) -> str:
    """Part-relative Target string, as PowerPoint writes them."""
    return posixpath.relpath(target_part, posixpath.dirname(source_part))


def _new_slide_root() -> tuple[etree._Element, etree._Element]:
    """Minimal p:sld skeleton (the shape PowerPoint itself accepts); returns
    (root, spTree)."""
    nsmap = {k: NSMAP[k] for k in ("a", "r", "p")}
    sld = etree.Element(qn("p:sld"), nsmap=nsmap)
    csld = etree.SubElement(sld, qn("p:cSld"))
    sp_tree = etree.SubElement(csld, qn("p:spTree"))
    nv = etree.SubElement(sp_tree, qn("p:nvGrpSpPr"))
    cnv = etree.SubElement(nv, qn("p:cNvPr"))
    cnv.set("id", "1")
    cnv.set("name", "")
    etree.SubElement(nv, qn("p:cNvGrpSpPr"))
    etree.SubElement(nv, qn("p:nvPr"))
    grp = etree.SubElement(sp_tree, qn("p:grpSpPr"))
    xfrm = etree.SubElement(grp, qn("a:xfrm"))
    for tag, attrs in (
        ("a:off", ("x", "y")),
        ("a:ext", ("cx", "cy")),
        ("a:chOff", ("x", "y")),
        ("a:chExt", ("cx", "cy")),
    ):
        el = etree.SubElement(xfrm, qn(tag))
        for a in attrs:
            el.set(a, "0")
    clr = etree.SubElement(sld, qn("p:clrMapOvr"))
    etree.SubElement(clr, qn("a:masterClrMapping"))
    return sld, sp_tree


def _clone_placeholder_skeleton(
    pkg: PptxPackage, layout_part: str, sp_tree: etree._Element
) -> list[dict]:
    """Empty placeholder shapes on the new slide, one per cloneable layout
    placeholder (same type/idx/orient/sz so inheritance binds), skipping the
    latent date/footer/slide-number placeholders. Geometry and formatting are
    NOT copied: an empty spPr inherits position and style from the layout."""
    added: list[dict] = []
    next_id = 2  # spTree's own cNvPr is id=1 by convention
    csld = pkg.root(layout_part).find(qn("p:cSld"))
    lt_tree = csld.find(qn("p:spTree")) if csld is not None else None
    if lt_tree is None:
        return added
    ph_path = f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}"
    for sp in lt_tree.findall(qn("p:sp")):
        ph = sp.find(ph_path)
        if ph is None:
            continue
        ph_type = ph.get("type", "obj")
        if ph_type in _LATENT_PH_TYPES:
            continue
        new_sp = etree.SubElement(sp_tree, qn("p:sp"))
        nvsp = etree.SubElement(new_sp, qn("p:nvSpPr"))
        cnv = etree.SubElement(nvsp, qn("p:cNvPr"))
        cnv.set("id", str(next_id))
        basename = _PH_BASENAMES.get(ph_type, "Content Placeholder")
        if ph.get("orient") == "vert":
            basename = "Vertical " + basename
        cnv.set("name", f"{basename} {next_id - 1}")
        cnvsp = etree.SubElement(nvsp, qn("p:cNvSpPr"))
        locks = etree.SubElement(cnvsp, qn("a:spLocks"))
        locks.set("noGrp", "1")
        nvpr = etree.SubElement(nvsp, qn("p:nvPr"))
        new_ph = etree.SubElement(nvpr, qn("p:ph"))
        for attr in ("type", "orient", "sz", "idx"):
            value = ph.get(attr)
            if value is not None:
                new_ph.set(attr, value)
        etree.SubElement(new_sp, qn("p:spPr"))
        if ph_type in _PH_TEXT_TYPES:
            tx = etree.SubElement(new_sp, qn("p:txBody"))
            etree.SubElement(tx, qn("a:bodyPr"))
            etree.SubElement(tx, qn("a:lstStyle"))
            etree.SubElement(tx, qn("a:p"))
        added.append(
            {
                "shape_id": next_id,
                "type": ph_type,
                "idx": int(ph.get("idx", "0")),
            }
        )
        next_id += 1
    return added


def _regenerate_creation_ids(root: etree._Element) -> int:
    """Fresh GUIDs for every a16:creationId (shapes) and p14:creationId
    (slide). Duplicate creationIds across slides after repeated duplication
    are a confirmed PowerPoint repair-prompt source."""
    count = 0
    for el in root.iter(qn("a16:creationId"), qn("p14:creationId")):
        el.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        count += 1
    return count


# ------------------------------------------------------------- sections


def _section_lst(pkg: PptxPackage) -> etree._Element | None:
    ext_lst = pkg.presentation().find(qn("p:extLst"))
    if ext_lst is None:
        return None
    for ext in ext_lst.findall(qn("p:ext")):
        sec = ext.find(qn("p14:sectionLst"))
        if sec is not None:
            return sec
    return None


def _section_sld_id_lst(section: etree._Element) -> etree._Element:
    lst = section.find(qn("p14:sldIdLst"))
    if lst is None:
        lst = etree.Element(qn("p14:sldIdLst"))
        section.insert(0, lst)  # p14:section children: sldIdLst then extLst
    return lst


def _drop_section_entry(pkg: PptxPackage, slide_id: int) -> int:
    sec_lst = _section_lst(pkg)
    if sec_lst is None:
        return 0
    removed = 0
    for section in sec_lst.findall(qn("p14:section")):
        lst = section.find(qn("p14:sldIdLst"))
        if lst is None:
            continue
        for entry in lst.findall(qn("p14:sldId")):
            if entry.get("id") == str(slide_id):
                lst.remove(entry)
                removed += 1
    if removed:
        pkg.mark_dirty(PRESENTATION_PART)
    return removed


def _assign_section_membership(
    pkg: PptxPackage, slide_id: int, index: int
) -> bool:
    """Put a (new or moved) slide id into exactly one section when a
    sectionLst exists: the section of its predecessor in deck order,
    inserted right after the predecessor's entry; a slide at position 0
    joins the front of the first section."""
    sec_lst = _section_lst(pkg)
    if sec_lst is None:
        return False
    _drop_section_entry(pkg, slide_id)  # idempotent
    entry = etree.Element(qn("p14:sldId"))
    entry.set("id", str(slide_id))
    sections = sec_lst.findall(qn("p14:section"))
    if not sections:
        return False
    entries = _slide_entries(pkg)
    pred_id = entries[index - 1].get("id") if index > 0 else None
    if pred_id is not None:
        for section in sections:
            lst = section.find(qn("p14:sldIdLst"))
            if lst is None:
                continue
            for existing in lst.findall(qn("p14:sldId")):
                if existing.get("id") == pred_id:
                    existing.addnext(entry)
                    pkg.mark_dirty(PRESENTATION_PART)
                    return True
    _section_sld_id_lst(sections[0]).insert(0, entry)
    pkg.mark_dirty(PRESENTATION_PART)
    return True


def _normalize_sections(pkg: PptxPackage) -> list[int]:
    """Rebuild section membership after reordering so the invariants hold:
    every slide in exactly one section, each section's p14:sldIdLst in deck
    order, sections contiguous in deck order. A slide whose original section
    would become non-contiguous is reassigned to the section open at its new
    position; reassigned slide ids are returned for the caller's report."""
    sec_lst = _section_lst(pkg)
    if sec_lst is None:
        return []
    sections = sec_lst.findall(qn("p14:section"))
    if not sections:
        return []
    membership: dict[str, etree._Element] = {}
    for section in sections:
        lst = section.find(qn("p14:sldIdLst"))
        if lst is None:
            continue
        for entry in lst.findall(qn("p14:sldId")):
            membership[entry.get("id")] = section
    deck = [e.get("id") for e in _slide_entries(pkg)]
    moves: list[int] = []
    assigned: dict[str, etree._Element] = {}
    open_section: etree._Element | None = None
    closed: set[int] = set()
    for sid in deck:
        want = membership.get(sid)
        if want is None:
            want = open_section if open_section is not None else sections[0]
        if want is not open_section:
            if id(want) in closed:
                moves.append(int(sid))
                want = open_section
            else:
                if open_section is not None:
                    closed.add(id(open_section))
                open_section = want
        assigned[sid] = want
    for section in sections:
        lst = _section_sld_id_lst(section)
        for child in list(lst):
            lst.remove(child)
    ordered: list[etree._Element] = []
    for sid in deck:
        section = assigned[sid]
        entry = etree.SubElement(_section_sld_id_lst(section), qn("p14:sldId"))
        entry.set("id", sid)
        if section not in ordered:
            ordered.append(section)
    for section in ordered:  # section listing order mirrors deck order
        sec_lst.append(section)
    pkg.mark_dirty(PRESENTATION_PART)
    return moves


# ---------------------------------------------------------- insert_slide


def insert_slide(pkg: PptxPackage, layout, position: int | None = None) -> dict:
    """New slide built from a layout (name, or 0-based global index across
    all masters), carrying the layout's placeholder skeleton: empty
    placeholders with the layout's type/idx/orient/sz so inheritance binds,
    minus the latent date/footer/slide-number placeholders. Inserted at
    `position` (0-based final index) or at the end."""
    layout_part = _resolve_layout(pkg, layout)
    sld, sp_tree = _new_slide_root()
    placeholders = _clone_placeholder_skeleton(pkg, layout_part, sp_tree)
    info = pkg.add_slide_part(_serialize(sld), layout_part=layout_part)
    lst = _sld_id_lst(pkg)
    entries = lst.findall(qn("p:sldId"))
    index = len(entries) - 1
    if position is not None:
        if not 0 <= position < len(entries):
            raise TargetNotFound(
                f"position {position} out of range; the deck now has "
                f"{len(entries)} slides (valid: 0..{len(entries) - 1})"
            )
        entry = entries[index]
        lst.remove(entry)
        lst.insert(position, entry)
        pkg.mark_dirty(PRESENTATION_PART)
        index = position
    in_section = _assign_section_membership(pkg, info["slide_id"], index)
    return {
        "slide_id": info["slide_id"],
        "part": info["part"],
        "index": index,
        "layout": layout_part,
        "placeholders": placeholders,
        "in_section": in_section,
    }


# ------------------------------------------------------- duplicate_slide


def _copy_content_type(pkg: PptxPackage, source: str, dest: str) -> None:
    """Mirror the source part's explicit Override onto the clone; when the
    source is covered by a Default (same extension), the clone is too and
    nothing is added."""
    ct_root = pkg.root("[Content_Types].xml")
    for node in ct_root.findall(qn("ct:Override")):
        if node.get("PartName") == "/" + source:
            pkg.add_content_type_override(dest, node.get("ContentType"))
            return


def _clone_partname(pkg: PptxPackage, source_part: str) -> str:
    """Collision-safe partname for a clone, in the source's folder with the
    source's naming stem: chart3.xml -> chart<max+1>.xml."""
    folder, fname = posixpath.split(source_part)
    stem, ext = posixpath.splitext(fname)
    stem = re.sub(r"\d+$", "", stem)
    return pkg.next_partname(f"{folder}/{stem}{{}}{ext}")


def _clone_simple_part(pkg: PptxPackage, source_part: str) -> str:
    """Byte-copy a leaf part (embedded xlsx, oleObject, tags) under a new
    collision-safe name, mirroring its content-type registration. Any rels
    of the source are copied as-is (shared targets)."""
    new_part = _clone_partname(pkg, source_part)
    pkg.set_raw_part(new_part, pkg.part_bytes(source_part))
    _copy_content_type(pkg, source_part, new_part)
    src_rels = rels_name(source_part)
    if pkg.has_part(src_rels):
        pkg.set_raw_part(rels_name(new_part), pkg.part_bytes(src_rels))
    return new_part


def _clone_notes_slide(
    pkg: PptxPackage, notes_part: str, new_slide_part: str
) -> str:
    """Deep-copy a notesSlide for a cloned slide. The notesMaster rel stays
    shared; the back-rel to the owning slide is retargeted at the clone
    (sharing a notesSlide between two slides shows the same notes on both
    and dangles when either is deleted)."""
    new_notes = pkg.next_partname("ppt/notesSlides/notesSlide{}.xml")
    root = copy.deepcopy(pkg.root(notes_part))
    _regenerate_creation_ids(root)
    pkg.add_part_with_content_type(new_notes, _serialize(root), CT_NOTES_SLIDE)
    src_rels = rels_name(notes_part)
    if pkg.has_part(src_rels):
        rels_root = copy.deepcopy(pkg.root(src_rels))
        for rel in rels_root:
            if rel.get("TargetMode") == "External":
                continue
            if rel.get("Type") == RT_SLIDE:
                rel.set("Target", _rel_target(new_notes, new_slide_part))
        pkg.set_raw_part(rels_name(new_notes), _serialize(rels_root))
    return new_notes


def _clone_chart(pkg: PptxPackage, chart_part: str) -> tuple[str, list[str]]:
    """Deep-copy a chart part with its companion colors/style parts and the
    embedded xlsx (two slides must never share a chart part: edits bleed and
    PowerPoint offers Repair). The chart XML is copied verbatim; its rels
    keep their rIds and only the deep-copied Targets are rewritten, so
    c:externalData r:id references stay valid without touching the XML."""
    new_chart = _clone_partname(pkg, chart_part)
    extras: list[str] = []
    rels_root = None
    src_rels = rels_name(chart_part)
    if pkg.has_part(src_rels):
        rels_root = copy.deepcopy(pkg.root(src_rels))
        for rel in rels_root:
            if rel.get("TargetMode") == "External":
                continue
            if rel.get("Type") in (RT_CHART_COLORS, RT_CHART_STYLE, RT_PACKAGE):
                target = resolve_target(chart_part, rel.get("Target", ""))
                new_target = _clone_simple_part(pkg, target)
                rel.set("Target", _rel_target(new_chart, new_target))
                extras.append(new_target)
    pkg.set_raw_part(new_chart, pkg.part_bytes(chart_part))
    _copy_content_type(pkg, chart_part, new_chart)
    if rels_root is not None:
        pkg.set_raw_part(rels_name(new_chart), _serialize(rels_root))
    return new_chart, extras


def duplicate_slide(pkg: PptxPackage, slide, position: int | None = None) -> dict:
    """Fully independent copy of a slide (the python-pptx issue #132
    machinery). Placed right after the source by default, or at `position`
    (0-based final index). Layout and media rels stay shared; notesSlide,
    chart (with colors/style/xlsx), and oleObject/package parts are
    deep-copied; hyperlink rels are duplicated with TargetMode preserved;
    creationId GUIDs are regenerated; partnames are collision-safe.
    Comments (modern threaded AND classic) are STRIPPED from the copy with
    a warning, never shared or silently carried: modern comment anchors
    reference the source slide's creation ids (which this copy
    regenerates), so a shared or naively-copied comments part would render
    mis-anchored review threads on the clone."""
    src_index, src_part, _entry, src_id, _rid = _resolve_slide(pkg, slide)
    n_after = len(_slide_entries(pkg)) + 1
    final = src_index + 1 if position is None else position
    if not 0 <= final < n_after:
        raise TargetNotFound(
            f"position {position} out of range; the deck will have "
            f"{n_after} slides (valid: 0..{n_after - 1})"
        )
    new_part = pkg.next_partname("ppt/slides/slide{}.xml")

    new_root = copy.deepcopy(pkg.root(src_part))
    guids = _regenerate_creation_ids(new_root)

    copied: list[str] = []
    shared: list[str] = []
    stripped_comments: list[str] = []
    rels_root = None
    src_rels = rels_name(src_part)
    if pkg.has_part(src_rels):
        rels_root = copy.deepcopy(pkg.root(src_rels))
        for rel in list(rels_root):
            if rel.get("TargetMode") == "External":
                continue  # hyperlink/linked-media rels: duplicated as-is
            rel_type = rel.get("Type")
            target = resolve_target(src_part, rel.get("Target", ""))
            if rel_type in (RT_MODERN_COMMENTS, RT_LEGACY_COMMENTS):
                # Comments never travel with the clone: modern anchors
                # reference creation ids this copy regenerates, and a
                # SHARED comments part would show the same threads
                # mis-anchored on both slides. Strip rel + extLst wiring;
                # the caller is warned below.
                rels_root.remove(rel)
                stripped_comments.append(target)
                continue
            if rel_type == RT_NOTES_SLIDE:
                new_target = _clone_notes_slide(pkg, target, new_part)
            elif rel_type == RT_CHART:
                new_target, extras = _clone_chart(pkg, target)
                copied.extend(extras)
            elif rel_type in _DEEP_COPY_SIMPLE:
                new_target = _clone_simple_part(pkg, target)
            else:
                shared.append(target)  # layout, media, jump-hyperlink slides
                continue
            rel.set("Target", _rel_target(new_part, new_target))
            copied.append(new_target)
    if stripped_comments:
        # Drop the {6950BFC3-...} commentRel extension from the clone's
        # extLst; a wiring stub pointing at a removed rel corrupts.
        for ext in list(new_root.iter(qn("p:ext"))):
            if ext.get("uri") == EXT_URI_COMMENT_REL:
                parent = ext.getparent()
                parent.remove(ext)
                if len(parent) == 0 and parent.tag == qn("p:extLst"):
                    parent.getparent().remove(parent)

    pkg.add_part_with_content_type(new_part, _serialize(new_root), CT_SLIDE)
    if rels_root is not None:
        pkg.set_raw_part(rels_name(new_part), _serialize(rels_root))

    reg = pkg.register_slide_entry(new_part, position=final)
    in_section = _assign_section_membership(pkg, reg["slide_id"], final)
    result = {
        "slide_id": reg["slide_id"],
        "part": new_part,
        "index": final,
        "source_slide_id": src_id,
        "source_part": src_part,
        "copied_parts": copied,
        "shared_parts": sorted(set(shared)),
        "creation_ids_regenerated": guids,
        "in_section": in_section,
    }
    if stripped_comments:
        result["comments_stripped"] = sorted(set(stripped_comments))
        result["warnings"] = [
            "the source slide carries comments; they were NOT copied to "
            "the duplicate (comment anchors reference the original "
            "slide's identity and would mis-anchor on the clone). The "
            "original slide keeps its comments; add fresh ones to the "
            "copy with add_comment if needed."
        ]
    return result


# ---------------------------------------------------------- delete_slide


def _remove_relationship(pkg: PptxPackage, source_part: str, rid: str) -> None:
    rels = pkg.rels_for(source_part)
    for rel in rels.getroot():
        if rel.get("Id") == rid:
            rels.getroot().remove(rel)
            pkg.mark_dirty(rels_name(source_part) if source_part else "_rels/.rels")
            return
    raise TargetNotFound(f"{source_part} has no relationship {rid}")


def _remove_part_and_rels(pkg: PptxPackage, part: str) -> None:
    pkg.remove_part(part)
    part_rels = rels_name(part)
    if pkg.has_part(part_rels):
        pkg.remove_part(part_rels)
    pkg.remove_content_type_override(part)


def _drop_custom_show_refs(pkg: PptxPackage, rid: str) -> dict:
    """Remove p:custShowLst entries referencing a deleted slide's rId.
    Shows left empty are removed (PowerPoint refuses empty custom shows);
    an emptied custShowLst is removed entirely."""
    pres = pkg.presentation()
    cust_lst = pres.find(qn("p:custShowLst"))
    result = {"entries_removed": 0, "shows_removed": []}
    if cust_lst is None:
        return result
    changed = False
    for show in cust_lst.findall(qn("p:custShow")):
        sld_lst = show.find(qn("p:sldLst"))
        if sld_lst is None:
            continue
        for sld in sld_lst.findall(qn("p:sld")):
            if sld.get(qn("r:id")) == rid:
                sld_lst.remove(sld)
                result["entries_removed"] += 1
                changed = True
        if len(sld_lst.findall(qn("p:sld"))) == 0:
            result["shows_removed"].append(show.get("name", ""))
            cust_lst.remove(show)
            changed = True
    if len(cust_lst.findall(qn("p:custShow"))) == 0:
        pres.remove(cust_lst)
        changed = True
    if changed:
        pkg.mark_dirty(PRESENTATION_PART)
    return result


def _neuter_jump_hyperlinks(pkg: PptxPackage, deleted_part: str) -> list[dict]:
    """Other parts (slides, layouts, masters, notesSlides) can carry jump
    hyperlinks (slide-reltype rels) targeting the deleted slide. Drop the
    rel and remove the referencing a:hlinkClick/a:hlinkHover elements so the
    click becomes a no-op; every neutered link is flagged in the result."""
    flagged: list[dict] = []
    pres_rels = rels_name(PRESENTATION_PART)
    for name in list(pkg.part_names()):
        if not name.endswith(".rels") or name in (pres_rels, "_rels/.rels"):
            continue
        source = rels_source(name)
        root = pkg.root(name)
        changed = False
        for rel in list(root):
            if (
                rel.get("Type") == RT_SLIDE
                and rel.get("TargetMode") != "External"
                and resolve_target(source, rel.get("Target", "")) == deleted_part
            ):
                rid = rel.get("Id")
                root.remove(rel)
                changed = True
                removed = 0
                src_root = pkg.root(source)
                for el in list(src_root.iter(qn("a:hlinkClick"), qn("a:hlinkHover"))):
                    if el.get(qn("r:id")) == rid:
                        el.getparent().remove(el)
                        removed += 1
                if removed:
                    pkg.mark_dirty(source)
                flagged.append(
                    {"part": source, "rid": rid, "hyperlinks_removed": removed}
                )
        if changed:
            pkg.mark_dirty(name)
    return flagged


def _is_referenced(pkg: PptxPackage, target: str) -> bool:
    for name in pkg.part_names():
        if not name.endswith(".rels"):
            continue
        source = rels_source(name)
        for rel in pkg.root(name):
            if rel.get("TargetMode") == "External":
                continue
            if resolve_target(source, rel.get("Target", "")) == target:
                return True
    return False


def _gc_parts(pkg: PptxPackage, candidates: list[str]) -> list[str]:
    """Remove parts no longer referenced by ANY remaining rels file,
    starting from the deleted slide's former targets and cascading (a gc'd
    chart releases its colors/style/xlsx; a gc'd notesSlide never releases
    the notesMaster, which the presentation rels still reference; media
    shared with other slides survives the reference count)."""
    removed: list[str] = []
    queue = list(candidates)
    while queue:
        target = queue.pop(0)
        if target in removed or not pkg.has_part(target):
            continue
        if _is_referenced(pkg, target):
            continue
        target_rels = rels_name(target)
        if pkg.has_part(target_rels):
            for rel in pkg.root(target_rels):
                if rel.get("TargetMode") != "External":
                    queue.append(resolve_target(target, rel.get("Target", "")))
        _remove_part_and_rels(pkg, target)
        removed.append(target)
    return removed


def _delete_slide_impl(pkg: PptxPackage, slide, *, allow_last: bool) -> dict:
    index, part, entry, slide_id, rid = _resolve_slide(pkg, slide)
    if len(_slide_entries(pkg)) <= 1 and not allow_last:
        raise UnsupportedStructure(
            "refusing to delete the last remaining slide; PowerPoint decks "
            "should keep at least one slide (create_presentation with "
            "keep_slides=False is the sanctioned zero-slide path)"
        )
    candidates: list[str] = []
    part_rels = rels_name(part)
    if pkg.has_part(part_rels):
        for rel in pkg.root(part_rels):
            if rel.get("TargetMode") != "External":
                candidates.append(resolve_target(part, rel.get("Target", "")))

    lst = _sld_id_lst(pkg)
    lst.remove(entry)
    pkg.mark_dirty(PRESENTATION_PART)
    _remove_relationship(pkg, PRESENTATION_PART, rid)
    custom_shows = _drop_custom_show_refs(pkg, rid)
    section_entries = _drop_section_entry(pkg, slide_id)
    _remove_part_and_rels(pkg, part)
    flagged = _neuter_jump_hyperlinks(pkg, part)
    gc_removed = _gc_parts(pkg, candidates)
    return {
        "deleted_part": part,
        "slide_id": slide_id,
        "index": index,
        "gc_parts": gc_removed,
        "custom_shows": custom_shows,
        "section_entries_removed": section_entries,
        "flagged_hyperlinks": flagged,
    }


def delete_slide(pkg: PptxPackage, slide) -> dict:
    """Delete a slide and garbage-collect everything only it referenced:
    its notesSlide twin, charts with their colors/style/xlsx companions,
    embeddings, and media used nowhere else. Custom-show references, section
    membership, and jump hyperlinks on other slides targeting the deleted
    slide are cleaned up (neutered hyperlinks are flagged in the result).
    Refuses to delete the last remaining slide."""
    return _delete_slide_impl(pkg, slide, allow_last=False)


# ------------------------------------------------------ reorder and hide


def reorder_slides(pkg: PptxPackage, order: list[int]) -> dict:
    """Reorder the whole deck. `order` is a full permutation of the current
    0-based indices, listed in their new sequence (order[i] = current index
    of the slide that lands at position i). Pure p:sldIdLst surgery, then
    section membership is renormalized (contiguity enforced; reassigned
    slide ids reported)."""
    entries = _slide_entries(pkg)
    n = len(entries)
    if sorted(order) != list(range(n)):
        raise PptMcpError(
            f"order must be a full permutation of 0..{n - 1} listing current "
            f"indices in their new sequence; got {order!r}"
        )
    lst = _sld_id_lst(pkg)
    for current in order:
        lst.append(entries[current])  # lxml append MOVES a parented element
    pkg.mark_dirty(PRESENTATION_PART)
    section_moves = _normalize_sections(pkg)
    return {
        "order": list(order),
        "slide_ids": [int(e.get("id")) for e in _slide_entries(pkg)],
        "section_moves": section_moves,
    }


def move_slide(pkg: PptxPackage, slide, to: int) -> dict:
    """Move one slide to a new 0-based index. The slide joins the section of
    its new neighborhood (predecessor's section; first section when moved to
    the front), matching PowerPoint's drag behavior."""
    index, part, entry, slide_id, _rid = _resolve_slide(pkg, slide)
    n = len(_slide_entries(pkg))
    if not 0 <= to < n:
        raise TargetNotFound(
            f"destination index {to} out of range, presentation has {n} "
            f"slide{'s' if n != 1 else ''} (valid: 0..{n - 1})"
        )
    lst = _sld_id_lst(pkg)
    lst.remove(entry)
    lst.insert(to, entry)
    pkg.mark_dirty(PRESENTATION_PART)
    in_section = _assign_section_membership(pkg, slide_id, to)
    section_moves = _normalize_sections(pkg)
    return {
        "slide_id": slide_id,
        "part": part,
        "from": index,
        "to": to,
        "in_section": in_section,
        "section_moves": section_moves,
    }


def set_slide_hidden(pkg: PptxPackage, slide, hidden: bool) -> dict:
    """Hide or unhide a slide in slide-show mode: p:sld/@show="0" hides,
    removing the attribute restores the default (shown)."""
    index, part, _entry, slide_id, _rid = _resolve_slide(pkg, slide)
    root = pkg.root(part)
    if hidden:
        root.set("show", "0")
    else:
        root.attrib.pop("show", None)
    pkg.mark_dirty(part)
    return {"slide_id": slide_id, "index": index, "hidden": bool(hidden)}


# --------------------------------------------------- create_presentation


def create_presentation(
    path: str | Path,
    template: str | Path | None = None,
    *,
    keep_slides: bool = False,
) -> dict:
    """Create a NEW .pptx at `path` from a template (.pptx or .potx). The
    template's bytes are copied, a .potx/.ppsx main content type is restamped
    to the presentation type, and with keep_slides=False (the default) every
    slide is removed through the delete machinery, leaving masters, layouts,
    and themes intact. The template source file is never modified. With
    template=None a blank 4:3 deck is built from python-pptx's bundled
    default template bytes (16:9 blank creation is a Phase 7 item).

    This is the one ops function that writes to disk: byte copy, then a
    PptxPackage edit and an atomic validated save on the NEW file only."""
    dest = Path(path)
    check_path(dest, "create presentation")
    if dest.suffix.lower() != ".pptx":
        raise PptMcpError(
            f"output must be a .pptx file, got {dest.name!r}"
        )
    if dest.exists():
        raise PptMcpError(
            f"{dest} already exists; refusing to overwrite. Delete it first "
            "or choose another name."
        )
    if template is None:
        import pptx  # bundled default template used as raw bytes only

        src = Path(pptx.__file__).parent / "templates" / "default.pptx"
        if not src.exists():  # pragma: no cover
            raise DocumentNotFound(
                "python-pptx default template not found; pass an explicit "
                "template .pptx/.potx"
            )
    else:
        src = Path(template)
        check_path(src, "read template")
        if not src.exists():
            raise DocumentNotFound(f"no template file at {src}")
        if src.suffix.lower() not in (".pptx", ".potx", ".ppsx"):
            raise PptMcpError(
                f"template must be .pptx, .potx, or .ppsx, got {src.name!r}"
            )
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(src.read_bytes())
    try:
        pkg = PptxPackage(dest)
        converted = False
        ct_root = pkg.root("[Content_Types].xml")
        for node in ct_root.findall(qn("ct:Override")):
            if (
                node.get("PartName") == "/" + PRESENTATION_PART
                and node.get("ContentType") != CT_PRESENTATION_MAIN
            ):
                pkg.add_content_type_override(
                    PRESENTATION_PART, CT_PRESENTATION_MAIN
                )
                converted = True
        removed: list[str] = []
        gc_total: list[str] = []
        if not keep_slides:
            for i in range(len(pkg.slide_parts()) - 1, -1, -1):
                res = _delete_slide_impl(pkg, i, allow_last=True)
                removed.append(res["deleted_part"])
                gc_total.extend(res["gc_parts"])
        pkg.save(do_backup=False)
    except BaseException:
        dest.unlink(missing_ok=True)  # never leave a half-built file behind
        raise
    return {
        "path": str(dest),
        "template": str(src),
        "converted_content_type": converted,
        "slides_removed": len(removed),
        "gc_parts": gc_total,
        "slides_kept": len(pkg.slide_parts()),
    }


# ------------------------------------------------------------ apply_layout


def _ph_key(ph: etree._Element) -> tuple[str, int]:
    """Placeholder matching key: (normalized type, idx). PowerPoint binds a
    slide placeholder to a layout placeholder by idx (and by the title
    family regardless of ctrTitle vs title), so both title spellings
    normalize to one key."""
    ph_type = ph.get("type", "obj")
    if ph_type in ("title", "ctrTitle"):
        return ("title", 0)
    return (ph_type, int(ph.get("idx", "0")))


def _ph_key_from_record(record: dict) -> tuple[str, int]:
    if record["type"] in ("title", "ctrTitle"):
        return ("title", 0)
    return (record["type"], record["idx"])


def _layout_ph_keys(pkg: PptxPackage, layout_part: str) -> set[tuple[str, int]]:
    keys: set[tuple[str, int]] = set()
    csld = pkg.root(layout_part).find(qn("p:cSld"))
    lt_tree = csld.find(qn("p:spTree")) if csld is not None else None
    if lt_tree is None:
        return keys
    ph_path = f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}"
    for sp in lt_tree.findall(qn("p:sp")):
        ph = sp.find(ph_path)
        if ph is not None:
            keys.add(_ph_key(ph))
    return keys


def apply_layout(pkg: PptxPackage, slide, layout) -> dict:
    """Re-link a slide's layout relationship to a different layout (name or
    0-based global index; ambiguous names refuse).

    Placeholder reconciliation is MINIMAL and honest: nothing on the slide
    is added, removed, or moved. A slide placeholder whose (type, idx)
    exists on the new layout re-inherits position and styling from it; one
    with no match is an ORPHAN: it keeps its content and any explicit
    geometry but no longer inherits anything meaningful, and is reported
    with a warning so the caller can restyle or delete it. Layout
    placeholders with no slide counterpart are reported as unused
    (insert_slide clones them on fresh slides; this tool never fabricates
    shapes)."""
    index, part, _entry, slide_id, _rid = _resolve_slide(pkg, slide)
    layout_part = _resolve_layout(pkg, layout)

    rels = pkg.rels_for(part)
    layout_rel = None
    for rel in rels.getroot():
        if rel.get("Type") == RT_SLIDE_LAYOUT:
            layout_rel = rel
            break
    if layout_rel is None:
        raise UnsupportedStructure(
            f"slide {index} has no slideLayout relationship; refusing to "
            "invent one"
        )
    old_layout = resolve_target(part, layout_rel.get("Target", ""))
    if old_layout == layout_part:
        return {
            "slide_id": slide_id,
            "index": index,
            "layout": layout_part,
            "changed": False,
            "note": "slide already uses that layout",
        }
    layout_rel.set("Target", _rel_target(part, layout_part))
    pkg.mark_dirty(rels_name(part))

    layout_keys = _layout_ph_keys(pkg, layout_part)
    matched: list[dict] = []
    orphaned: list[dict] = []
    warnings: list[str] = []
    csld = pkg.root(part).find(qn("p:cSld"))
    sp_tree = csld.find(qn("p:spTree")) if csld is not None else None
    ph_path = f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}"
    if sp_tree is not None:
        for sp in sp_tree.findall(qn("p:sp")):
            ph = sp.find(ph_path)
            if ph is None:
                continue
            cnvpr = sp.find(f"{qn('p:nvSpPr')}/{qn('p:cNvPr')}")
            record = {
                "shape_id": (
                    int(cnvpr.get("id")) if cnvpr is not None else None
                ),
                "type": ph.get("type", "obj"),
                "idx": int(ph.get("idx", "0")),
            }
            if _ph_key(ph) in layout_keys:
                matched.append(record)
            else:
                orphaned.append(record)
    for rec_ph in orphaned:
        warnings.append(
            f"placeholder shape {rec_ph['shape_id']} (type "
            f"{rec_ph['type']!r}, idx {rec_ph['idx']}) has no counterpart "
            f"on the new layout; content kept, inheritance broken"
        )
    unused = sorted(
        layout_keys - {_ph_key_from_record(r) for r in matched}
    )
    return {
        "slide_id": slide_id,
        "index": index,
        "layout": layout_part,
        "old_layout": old_layout,
        "changed": True,
        "placeholders_matched": matched,
        "placeholders_orphaned": orphaned,
        "layout_placeholders_unused": [
            {"type": t, "idx": i} for t, i in unused
        ],
        "warnings": warnings,
    }


# ----------------------------------------------------------- manage_section


_SECTION_ACTIONS = ("create", "rename", "delete", "move_slide_into")


def _sections_summary(pkg: PptxPackage) -> list[dict]:
    from .read import _sections

    return [
        {"name": s["name"], "slide_ids": s["slide_ids"]}
        for s in _sections(pkg)
    ]


def _resolve_section(
    sec_lst: etree._Element, section
) -> tuple[int, etree._Element]:
    sections = sec_lst.findall(qn("p14:section"))
    if isinstance(section, int) and not isinstance(section, bool):
        if not 0 <= section < len(sections):
            raise TargetNotFound(
                f"section index {section} out of range; the deck has "
                f"{len(sections)} section(s)"
            )
        return section, sections[section]
    if isinstance(section, str):
        hits = [
            (i, s) for i, s in enumerate(sections)
            if s.get("name", "") == section
        ]
        if len(hits) == 1:
            return hits[0]
        if len(hits) > 1:
            raise AmbiguousTarget(
                f"{len(hits)} sections are named {section!r}; use a 0-based "
                "section index instead"
            )
        names = ", ".join(repr(s.get("name", "")) for s in sections)
        raise TargetNotFound(
            f"no section named {section!r}; sections present: "
            f"{names or 'none'}"
        )
    raise PptMcpError(
        "section must be a section name (str) or 0-based index (int)"
    )


def _new_section_el(name: str) -> etree._Element:
    el = etree.Element(qn("p14:section"))
    el.set("name", name)
    el.set("id", "{" + str(uuid.uuid4()).upper() + "}")
    etree.SubElement(el, qn("p14:sldIdLst"))
    return el


def _ensure_section_lst(pkg: PptxPackage) -> etree._Element:
    """The p14:sectionLst, created (with its p:extLst/p:ext wrapper) when
    the deck has none yet."""
    existing = _section_lst(pkg)
    if existing is not None:
        return existing
    pres = pkg.presentation()
    ext_lst = pres.find(qn("p:extLst"))
    if ext_lst is None:
        ext_lst = etree.Element(qn("p:extLst"))
        pkg._insert_presentation_child(ext_lst)
    ext = etree.SubElement(ext_lst, qn("p:ext"))
    ext.set("uri", _SECTION_EXT_URI)
    sec_lst = etree.SubElement(
        ext, qn("p14:sectionLst"), nsmap={"p14": NSMAP["p14"]}
    )
    pkg.mark_dirty(PRESENTATION_PART)
    return sec_lst


def _section_names(sec_lst: etree._Element) -> list[str]:
    return [s.get("name", "") for s in sec_lst.findall(qn("p14:section"))]


def _create_section(pkg: PptxPackage, name: str, slide) -> dict:
    entries = _slide_entries(pkg)
    at_index = at_id = None
    if slide is not None:
        at_index, _part, _entry, at_id, _rid = _resolve_slide(pkg, slide)
    existing_lst = _section_lst(pkg)
    if existing_lst is not None and name in _section_names(existing_lst):
        raise PptMcpError(
            f"a section named {name!r} already exists; section names must "
            "stay unique so they remain addressable"
        )
    sec_lst = _ensure_section_lst(pkg)
    sections = sec_lst.findall(qn("p14:section"))
    new_sec = _new_section_el(name)
    if not sections:
        # First sectioning of the deck: cover EVERY slide (the invariant).
        if at_index in (None, 0) or not entries:
            sec_lst.append(new_sec)
            lst = new_sec.find(qn("p14:sldIdLst"))
            for entry in entries:
                e = etree.SubElement(lst, qn("p14:sldId"))
                e.set("id", entry.get("id"))
        else:
            default = _new_section_el("Default Section")
            sec_lst.append(default)
            sec_lst.append(new_sec)
            dlst = default.find(qn("p14:sldIdLst"))
            nlst = new_sec.find(qn("p14:sldIdLst"))
            for i, entry in enumerate(entries):
                target = dlst if i < at_index else nlst
                e = etree.SubElement(target, qn("p14:sldId"))
                e.set("id", entry.get("id"))
    elif at_index is None:
        sec_lst.append(new_sec)  # empty section at the end
    else:
        # Split the section containing the slide at that slide.
        container = None
        for sec in sections:
            lst = sec.find(qn("p14:sldIdLst"))
            if lst is None:
                continue
            if any(
                entry.get("id") == str(at_id)
                for entry in lst.findall(qn("p14:sldId"))
            ):
                container = (sec, lst)
                break
        if container is None:
            # Invariant breach in the source deck; adopt the slide.
            sec_lst.append(new_sec)
            lst = new_sec.find(qn("p14:sldIdLst"))
            e = etree.SubElement(lst, qn("p14:sldId"))
            e.set("id", str(at_id))
        else:
            sec, lst = container
            sec.addnext(new_sec)
            nlst = new_sec.find(qn("p14:sldIdLst"))
            take = False
            for entry in list(lst.findall(qn("p14:sldId"))):
                if entry.get("id") == str(at_id):
                    take = True
                if take:
                    lst.remove(entry)
                    nlst.append(entry)
    pkg.mark_dirty(PRESENTATION_PART)
    moves = _normalize_sections(pkg)
    return {
        "action": "create",
        "name": name,
        "section_moves": moves,
        "sections": _sections_summary(pkg),
    }


def _delete_section(pkg: PptxPackage, section) -> dict:
    sec_lst = _section_lst(pkg)
    if sec_lst is None:
        raise TargetNotFound("the deck has no sections")
    idx, sec = _resolve_section(sec_lst, section)
    sections = sec_lst.findall(qn("p14:section"))
    deleted_name = sec.get("name", "")
    merged_into = None
    if len(sections) == 1:
        # Last section: remove sectioning entirely (invariant-clean).
        ext = sec_lst.getparent()
        ext_lst = ext.getparent()
        ext_lst.remove(ext)
        if len(ext_lst) == 0:
            ext_lst.getparent().remove(ext_lst)
    else:
        neighbor = sections[idx - 1] if idx > 0 else sections[idx + 1]
        merged_into = neighbor.get("name", "")
        nlst = _section_sld_id_lst(neighbor)
        lst = sec.find(qn("p14:sldIdLst"))
        if lst is not None:
            moved_entries = list(lst.findall(qn("p14:sldId")))
            if idx > 0:
                for entry in moved_entries:
                    lst.remove(entry)
                    nlst.append(entry)
            else:  # deleting the first section: slides PREPEND to the next
                for entry in reversed(moved_entries):
                    lst.remove(entry)
                    nlst.insert(0, entry)
        sec_lst.remove(sec)
    pkg.mark_dirty(PRESENTATION_PART)
    moves = _normalize_sections(pkg)
    return {
        "action": "delete",
        "deleted": deleted_name,
        "merged_into": merged_into,
        "section_moves": moves,
        "sections": _sections_summary(pkg),
    }


def _move_slide_into_section(pkg: PptxPackage, slide, section) -> dict:
    sec_lst = _section_lst(pkg)
    if sec_lst is None:
        raise TargetNotFound(
            "the deck has no sections; create one first (action='create')"
        )
    sec_idx, sec = _resolve_section(sec_lst, section)
    index, part, entry, slide_id, _rid = _resolve_slide(pkg, slide)

    # Destination deck position: right after the target section's current
    # last slide; for an empty section, after the last slide of the nearest
    # preceding non-empty section (deck start when none).
    sections = sec_lst.findall(qn("p14:section"))
    by_id = {e.get("id"): i for i, e in enumerate(_slide_entries(pkg))}

    def _last_index(sec_el) -> int | None:
        lst = sec_el.find(qn("p14:sldIdLst"))
        if lst is None:
            return None
        indexes = [
            by_id[e.get("id")]
            for e in lst.findall(qn("p14:sldId"))
            if e.get("id") in by_id and e.get("id") != str(slide_id)
        ]
        return max(indexes) if indexes else None

    dest_after = _last_index(sec)
    if dest_after is None:
        for prev in range(sec_idx - 1, -1, -1):
            dest_after = _last_index(sections[prev])
            if dest_after is not None:
                break
    to = 0 if dest_after is None else dest_after + 1
    if index < to:
        to -= 1  # removing the slide first shifts later positions left

    lst = _sld_id_lst(pkg)
    lst.remove(entry)
    lst.insert(to, entry)
    # Explicit membership: out of every section, into the target's tail.
    _drop_section_entry(pkg, slide_id)
    e = etree.SubElement(_section_sld_id_lst(sec), qn("p14:sldId"))
    e.set("id", str(slide_id))
    pkg.mark_dirty(PRESENTATION_PART)
    moves = _normalize_sections(pkg)
    return {
        "action": "move_slide_into",
        "slide_id": slide_id,
        "part": part,
        "from": index,
        "to": to,
        "section": sec.get("name", ""),
        "section_moves": moves,
        "sections": _sections_summary(pkg),
    }


def manage_section(
    pkg: PptxPackage,
    action: str,
    *,
    section=None,
    name: str | None = None,
    slide=None,
) -> dict:
    """Section lifecycle: create, rename, delete, move_slide_into. The
    Phase 2 invariants are renormalized after every mutation: when sections
    exist, every slide id lives in exactly one section, each section's list
    follows deck order, and sections are contiguous in deck order.

    - create: `name` required (unique). With `slide`: the new section
      starts at that slide, splitting the section containing it; on a deck
      with NO sections it covers from that slide to the end, with a
      "Default Section" ahead of it when needed. Without `slide`: the
      first section covers the whole deck; with existing sections an empty
      section is appended.
    - rename: `section` + new `name` (unique).
    - delete: `section`; its slides merge into the previous section (next
      when deleting the first); deleting the only section removes
      sectioning entirely. Slides are never deleted.
    - move_slide_into: `slide` + `section`; the slide MOVES in deck order
      to the section's end (sections are contiguous ranges, so membership
      cannot change without moving the slide)."""
    if action not in _SECTION_ACTIONS:
        raise PptMcpError(
            f"unknown action {action!r}; one of: "
            f"{', '.join(_SECTION_ACTIONS)}"
        )
    if action == "create":
        if not name or not str(name).strip():
            raise PptMcpError("create needs a non-empty section name")
        return _create_section(pkg, str(name).strip(), slide)
    if action == "rename":
        if section is None:
            raise PptMcpError("rename needs `section` (name or index)")
        if not name or not str(name).strip():
            raise PptMcpError("rename needs a non-empty new `name`")
        new_name = str(name).strip()
        sec_lst = _section_lst(pkg)
        if sec_lst is None:
            raise TargetNotFound("the deck has no sections")
        idx, sec = _resolve_section(sec_lst, section)
        old = sec.get("name", "")
        if new_name != old and new_name in _section_names(sec_lst):
            raise PptMcpError(
                f"a section named {new_name!r} already exists; names must "
                "stay unique"
            )
        sec.set("name", new_name)
        pkg.mark_dirty(PRESENTATION_PART)
        return {
            "action": "rename",
            "index": idx,
            "old_name": old,
            "name": new_name,
            "sections": _sections_summary(pkg),
        }
    if action == "delete":
        if section is None:
            raise PptMcpError("delete needs `section` (name or index)")
        return _delete_section(pkg, section)
    if section is None or slide is None:
        raise PptMcpError("move_slide_into needs both `slide` and `section`")
    return _move_slide_into_section(pkg, slide, section)


# =====================================================================
# TIER-1B APPEND-ONLY BLOCK: per-slide background (functionality audit
# item 8). Master/layout backgrounds belong to ops/masters.py (wave 8a);
# this is the PER-SLIDE p:bg override and its reset-to-inherit, the half
# agents used to fake with a full-bleed rectangle (breaking placeholder
# z-order and the reset story).
# =====================================================================


def _resolve_slide_rec(pkg: PptxPackage, slide) -> dict:
    """Local selector resolution (index or {"slide_id": N}) mirroring
    ops/read.py, kept self-contained inside the append-only band."""
    from .read import resolve_slide

    return resolve_slide(pkg, slide)


def set_slide_background(pkg: PptxPackage, slide, fill) -> dict:
    """Set or clear ONE slide's background (the p:bg override).

    fill:
    - "inherit" or None: remove the slide's own p:bg so the layout/master
      background shows through again (PowerPoint's Reset Background).
    - "RRGGBB" / "#RGB" / scheme name ("accent1") or
      {"type": "solid", "color": ..., "alpha": 0..1}: solid fill.
    - {"type": "gradient", "stops": [{"pos", "color", "alpha"?}...],
      "angle": deg, "radial": bool}: gradient fill.
    - {"type": "image", "image": path-or-base64, "tile": bool}: picture
      background through the shared media pipeline (dedup by content hash,
      extension Default registered); stretch to fill by default, tiled
      when tile=True.

    The p:bg element lands as the FIRST child of p:cSld (schema position;
    anywhere else is a repair dialog). "none" (a transparent background) is
    refused: PowerPoint has no transparent slide background, use "inherit"
    to fall back to the layout/master. Text contrast against the new
    background is the caller's to verify (check_layout's contrast check
    resolves slide backgrounds)."""
    from . import geometry as _g
    from .read import resolve_slide as _resolve

    rec = _resolve(pkg, slide)
    part = rec["part"]
    csld = pkg.root(part).find(qn("p:cSld"))
    if csld is None:
        raise UnsupportedStructure(f"{part} has no p:cSld")
    existing = csld.find(qn("p:bg"))
    had_override = existing is not None

    if fill is None or fill == "inherit":
        if existing is not None:
            csld.remove(existing)
            pkg.mark_dirty(part)
        return {
            "slide_index": rec["index"],
            "slide_id": rec["slide_id"],
            "background": "inherited",
            "changed": had_override,
            "previous_override_removed": had_override,
        }
    if fill == "none" or (isinstance(fill, dict) and fill.get("type") == "none"):
        raise PptMcpError(
            'a slide background cannot be "none" (PowerPoint has no '
            'transparent slide background); use "inherit" to fall back to '
            "the layout/master background"
        )

    media_part = None
    kind: str
    if isinstance(fill, dict) and fill.get("type") == "image":
        image = fill.get("image")
        if not image:
            raise PptMcpError(
                'image background spec needs an "image" key (file path or '
                "base64)"
            )
        from .media import _add_media, _image_rel, _load_image

        data, fmt = _load_image(image)
        media_part, _reused = _add_media(pkg, data, fmt)
        rid = _image_rel(pkg, part, media_part)
        fill_el = etree.Element(qn("a:blipFill"))
        blip = etree.SubElement(fill_el, qn("a:blip"))
        blip.set(qn("r:embed"), rid)
        if fill.get("tile"):
            tile = etree.SubElement(fill_el, qn("a:tile"))
            tile.set("tx", "0")
            tile.set("ty", "0")
            tile.set("sx", "100000")
            tile.set("sy", "100000")
            tile.set("flip", "none")
            tile.set("algn", "tl")
        else:
            stretch = etree.SubElement(fill_el, qn("a:stretch"))
            etree.SubElement(stretch, qn("a:fillRect"))
        kind = "image"
    else:
        fill_el = _g.fill_element(fill)
        if fill_el is None:
            raise PptMcpError(f"invalid background fill spec {fill!r}")
        if fill_el.tag == qn("a:noFill"):
            raise PptMcpError(
                'a slide background cannot be "none"; use "inherit" to '
                "fall back to the layout/master background"
            )
        kind = (
            "gradient" if fill_el.tag == qn("a:gradFill") else "solid"
        )

    bg = etree.Element(qn("p:bg"))
    bgpr = etree.SubElement(bg, qn("p:bgPr"))
    bgpr.append(fill_el)
    etree.SubElement(bgpr, qn("a:effectLst"))
    if existing is not None:
        csld.remove(existing)
    csld.insert(0, bg)  # p:bg is the first child of p:cSld by schema
    pkg.mark_dirty(part)
    result = {
        "slide_index": rec["index"],
        "slide_id": rec["slide_id"],
        "background": kind,
        "changed": True,
        "replaced_previous_override": had_override,
    }
    if media_part is not None:
        result["media_part"] = media_part
    return result
